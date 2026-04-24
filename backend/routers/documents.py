"""Document management endpoints: upload, download, delete, URL import."""

import os
import io
import json
import time
import logging
import mimetypes
import traceback
from uuid import uuid4

from fastapi import APIRouter, UploadFile, File, Depends, HTTPException, Request, BackgroundTasks
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session

from auth import verify_token
from database import get_db, Document, DocumentChunk, AgentShare, SessionLocal
from helpers.agent_helpers import _user_can_access_agent, _user_can_edit_agent
from helpers.tenant import _get_caller_company_id
from helpers.rate_limiting import _check_api_rate_limit, _API_UPLOAD_LIMIT
from rag_engine import process_document_for_user
from redis_client import get_redis
from utils import logger as _app_logger, event_tracker
from validation import (
    UrlUploadValidated,
    validate_file_extension,
    validate_file_content,
    validate_file_size,
    sanitize_filename,
    MAX_FILE_SIZE,
    ALLOWED_FILE_EXTENSIONS,
)
from google.cloud import storage

logger = logging.getLogger(__name__)
router = APIRouter()


def _fetch_and_parse_url(url: str) -> tuple[str, str]:
    """Fetch a URL, parse HTML, and return (cleaned_text_content, filename).

    Raises HTTPException on failure.
    """
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8",
        "Accept-Language": "fr-FR,fr;q=0.9,en-US;q=0.8,en;q=0.7",
        "Accept-Encoding": "gzip, deflate, br",
        "Connection": "keep-alive",
        "Upgrade-Insecure-Requests": "1",
        "Sec-Fetch-Dest": "document",
        "Sec-Fetch-Mode": "navigate",
        "Sec-Fetch-Site": "none",
        "Cache-Control": "max-age=0",
    }

    import requests as http_requests

    max_retries = 3
    retry_delay = 2
    html = None
    last_error = None

    def _is_safe_redirect(redirect_url: str) -> bool:
        blocked_patterns = [
            "localhost", "127.0.0.1", "0.0.0.0", "192.168.", "10.",
            "172.16.", "172.17.", "172.18.", "172.19.", "172.20.",
            "172.21.", "172.22.", "172.23.", "172.24.", "172.25.",
            "172.26.", "172.27.", "172.28.", "172.29.", "172.30.",
            "172.31.", "169.254.", "[::1]", "[fc", "[fd",
            "metadata.google.internal",
        ]
        return not any(pattern in redirect_url.lower() for pattern in blocked_patterns)

    for attempt in range(max_retries):
        try:
            response = http_requests.get(url, headers=headers, timeout=20, allow_redirects=False, verify=True)
            redirect_count = 0
            while response.is_redirect and redirect_count < 5:
                redirect_url = response.headers.get("Location", "")
                if not redirect_url or not _is_safe_redirect(redirect_url):
                    raise http_requests.exceptions.ConnectionError("Redirect to blocked destination")
                response = http_requests.get(redirect_url, headers=headers, timeout=20, allow_redirects=False, verify=True)
                redirect_count += 1
            response.raise_for_status()
            if response.encoding:
                html = response.text
            else:
                response.encoding = response.apparent_encoding
                html = response.text
            break
        except http_requests.exceptions.SSLError as e:
            logger.warning(f"SSL error on attempt {attempt + 1} for {url}: {e}")
            last_error = e
            if attempt < max_retries - 1:
                time.sleep(retry_delay)
        except (http_requests.exceptions.Timeout, http_requests.exceptions.ConnectionError) as e:
            logger.warning(f"Connection error on attempt {attempt + 1} for {url}: {e}")
            last_error = e
            if attempt < max_retries - 1:
                time.sleep(retry_delay)
        except http_requests.exceptions.HTTPError as e:
            logger.error(f"HTTP error for {url}: {e}")
            last_error = e
            break
        except Exception as e:
            logger.error(f"Unexpected error on attempt {attempt + 1} for {url}: {e}")
            last_error = e
            if attempt < max_retries - 1:
                time.sleep(retry_delay)

    if html is None:
        error_msg = f"Failed to fetch URL after {max_retries} attempts"
        if last_error:
            error_msg += f": {str(last_error)}"
        logger.error(error_msg)
        raise HTTPException(status_code=400, detail="Unable to fetch the provided URL. Please check the URL and try again.")

    from bs4 import BeautifulSoup

    try:
        from readability import Document as ReadabilityDocument
        use_readability = True
    except Exception:
        use_readability = False

    title = ""
    meta_desc = ""
    main_text = ""

    try:
        soup = BeautifulSoup(html, "lxml")
        if soup.title and soup.title.string:
            title = soup.title.string.strip()
        md = soup.find("meta", attrs={"name": "description"})
        if md and md.get("content"):
            meta_desc = md.get("content").strip()

        if use_readability:
            try:
                doc = ReadabilityDocument(html)
                main_html = doc.summary()
                main_soup = BeautifulSoup(main_html, "lxml")
                main_text = "\n".join(
                    [p.get_text(separator=" ", strip=True) for p in main_soup.find_all(["p", "h1", "h2", "h3"])]
                )
            except Exception:
                use_readability = False

        if not main_text:
            body = soup.body
            if body:
                for tag in body.find_all(["script", "style", "nav", "footer", "aside", "header", "form", "noscript"]):
                    tag.decompose()
                paragraphs = [
                    p.get_text(separator=" ", strip=True)
                    for p in body.find_all(["p", "h1", "h2", "h3"])
                    if p.get_text(strip=True)
                ]
                main_text = "\n".join(paragraphs)

        cleaned = []
        if title:
            cleaned.append(f"Title: {title}")
        if meta_desc:
            cleaned.append(f"Description: {meta_desc}")
        if main_text:
            cleaned.append("Content:\n" + main_text)

        content = "\n\n".join(cleaned)
        if not content.strip():
            content = soup.get_text(separator="\n", strip=True)

    except Exception as e:
        logger.warning(f"Failed to parse HTML for useful content, falling back to raw. Error: {e}")
        content = html

    filename = url.split("//")[-1][:100].replace("/", "_") + ".txt"

    max_chars = 200000
    if len(content) > max_chars:
        content = content[:max_chars]

    return content, filename


@router.post("/upload-url")
async def upload_url(request: UrlUploadValidated, user_id: str = Depends(verify_token), db: Session = Depends(get_db)):
    """Ajoute une URL comme document/source pour le RAG"""
    try:
        content, filename = _fetch_and_parse_url(request.url)

        doc_id = process_document_for_user(
            filename,
            content.encode("utf-8", errors="ignore"),
            int(user_id),
            db,
            agent_id=request.agent_id,
            company_id=_get_caller_company_id(user_id, db),
        )

        # Store the source URL on the document
        doc = db.query(Document).filter(Document.id == doc_id).first()
        if doc:
            doc.source_url = request.url
            db.commit()

        logger.info(f"URL ajoutée pour user {user_id}, agent {request.agent_id}: {request.url}")
        event_tracker.track_document_upload(int(user_id), request.url, len(content))

        return {"url": request.url, "document_id": doc_id, "agent_id": request.agent_id, "status": "uploaded"}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Erreur lors de l'ajout d'URL: {e}")
        raise HTTPException(status_code=500, detail="Erreur lors de l'ajout de l'URL")


@router.post("/documents/{document_id}/refresh-url")
async def refresh_document_url(document_id: int, user_id: str = Depends(verify_token), db: Session = Depends(get_db)):
    """Re-fetch a URL-sourced document and re-embed its content."""
    try:
        uid = int(user_id)

        # Find document and verify access
        document = db.query(Document).filter(Document.id == document_id).first()
        if not document:
            raise HTTPException(status_code=404, detail="Document not found")

        # Check ownership or edit access via AgentShare
        is_owner = document.user_id == uid
        if not is_owner and document.agent_id:
            share = (
                db.query(AgentShare)
                .filter(AgentShare.agent_id == document.agent_id, AgentShare.user_id == uid, AgentShare.can_edit == True)
                .first()
            )
            if not share:
                raise HTTPException(status_code=403, detail="Access denied")
        elif not is_owner:
            raise HTTPException(status_code=403, detail="Access denied")

        # Verify this document has a source URL
        if not document.source_url:
            raise HTTPException(status_code=400, detail="This document has no source URL to refresh")

        # Re-fetch and parse the URL
        content, filename = _fetch_and_parse_url(document.source_url)

        if not content.strip():
            raise HTTPException(status_code=400, detail="No content could be extracted from the URL")

        # Delete old chunks
        db.query(DocumentChunk).filter(DocumentChunk.document_id == document_id).delete()

        # Update document content
        document.content = content
        document.filename = filename
        db.commit()

        # Re-chunk and re-embed
        from file_loader import chunk_text
        from mistral_embeddings import get_embedding_fast
        import numpy as np

        chunks = chunk_text(content)

        def split_for_embedding(chunk, max_tokens=8192):
            chunk = chunk.replace("\x00", "")
            max_chars = max_tokens * 4
            return [chunk[i : i + max_chars] for i in range(0, len(chunk), max_chars)]

        for i, chunk in enumerate(chunks):
            sub_chunks = split_for_embedding(chunk, 8192)
            embeddings = []
            for sub in sub_chunks:
                embedding = get_embedding_fast(sub)
                embeddings.append(embedding)
            if embeddings:
                avg_embedding = list(np.mean(np.array(embeddings), axis=0))
            else:
                raise ValueError("No sub-chunks produced for embedding")
            doc_chunk = DocumentChunk(
                document_id=document_id,
                company_id=document.company_id,
                chunk_text=chunk,
                embedding_vec=avg_embedding,
                chunk_index=i,
            )
            db.add(doc_chunk)

        db.commit()

        logger.info(f"Document {document_id} refreshed from URL: {document.source_url} ({len(chunks)} chunks)")
        event_tracker.track_user_action(uid, f"url_refresh:{document.source_url}")

        return {
            "document_id": document_id,
            "status": "refreshed",
            "source_url": document.source_url,
            "chunks": len(chunks),
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error refreshing document {document_id}: {e}")
        db.rollback()
        raise HTTPException(status_code=500, detail="Error refreshing URL content")


@router.post("/upload")
async def upload_file(
    file: UploadFile = File(...), user_id: str = Depends(verify_token), db: Session = Depends(get_db)
):
    """Upload and process document for a specific agent, extracting full text from any supported file type

    Security: Validates file size, extension, and sanitizes filename.
    """
    if not _check_api_rate_limit(user_id, "upload", _API_UPLOAD_LIMIT):
        raise HTTPException(status_code=429, detail="Too many uploads. Please try again later.")
    logger.info(f"Appel reçu sur /upload : filename={file.filename if file else 'None'}")
    try:
        # Validate filename exists
        if not file.filename:
            raise HTTPException(status_code=400, detail="Filename is required")

        # Sanitize filename to prevent path traversal attacks
        original_filename = file.filename
        safe_filename = sanitize_filename(original_filename)

        # Validate file extension against whitelist
        if not validate_file_extension(safe_filename):
            allowed_exts = ", ".join(sorted(ALLOWED_FILE_EXTENSIONS))
            raise HTTPException(status_code=400, detail=f"File type not allowed. Allowed types: {allowed_exts}")

        logger.info(
            f"Début import PJ : filename={safe_filename}, content_type={file.content_type if hasattr(file, 'content_type') else 'unknown'}"
        )

        # Read content and validate size
        content = await file.read()
        content_size = len(content)

        # Validate file size
        if not validate_file_size(content_size):
            max_size_mb = MAX_FILE_SIZE / (1024 * 1024)
            raise HTTPException(status_code=413, detail=f"File too large (max {max_size_mb:.0f}MB)")

        # Validate file content matches extension (magic bytes)
        if not validate_file_content(content, safe_filename):
            raise HTTPException(status_code=400, detail="File content does not match its extension")

        logger.info(f"PJ reçue : filename={safe_filename}, taille={content_size} octets")
        filename = safe_filename.lower()
        text = None

        if filename.endswith(".pdf"):
            MAX_PDF_PAGES = 500
            logger.info("Tentative extraction PDF (pdfplumber)")
            import pdfplumber
            import tempfile

            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp.write(content)
                tmp.flush()
                logger.info(f"Fichier PDF temporaire créé : {tmp.name}")
                try:
                    with pdfplumber.open(tmp.name) as pdf:
                        if len(pdf.pages) > MAX_PDF_PAGES:
                            raise HTTPException(
                                status_code=400, detail=f"PDF too large ({len(pdf.pages)} pages, max {MAX_PDF_PAGES})"
                            )
                        text = "\n".join([page.extract_text() or "" for page in pdf.pages])
                    logger.info(
                        f"Texte PDF extrait (pdfplumber) : longueur={len(text) if text else 0}, aperçu='{text[:200] if text else ''}'"
                    )
                except Exception as e:
                    logger.error(f"PDF extraction error: {e}")
                # Si le texte est vide, tente l'OCR sur chaque page
                if not text or not text.strip():
                    logger.info("PDF vide ou non textuel, tentative OCR (pytesseract)")
                    try:
                        from PIL import Image
                        import pytesseract

                        with pdfplumber.open(tmp.name) as pdf:
                            ocr_text = ""
                            for i, page in enumerate(pdf.pages):
                                img = page.to_image(resolution=300)
                                pil_img = img.original
                                page_ocr = pytesseract.image_to_string(pil_img, lang="fra")
                                logger.info(f"OCR page {i + 1}: longueur={len(page_ocr)}, aperçu='{page_ocr[:100]}'")
                                ocr_text += page_ocr + "\n"
                        text = ocr_text
                        logger.info(f"OCR PDF extrait: longueur={len(text)}, aperçu='{text[:200]}'")
                    except Exception as e:
                        logger.error(f"PDF OCR extraction error: {e}")
            os.unlink(tmp.name)
        elif filename.endswith(".docx"):
            logger.info("Tentative extraction DOCX")
            from docx import Document as DocxDocument
            import tempfile
            import tempfile

            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                tmp.write(content)
                tmp.flush()
                logger.info(f"Fichier DOCX temporaire créé : {tmp.name}")
                try:
                    doc = DocxDocument(tmp.name)
                    text = "\n".join([p.text for p in doc.paragraphs])
                    logger.info(
                        f"Texte DOCX extrait : longueur={len(text) if text else 0}, aperçu='{text[:200] if text else ''}'"
                    )
                except Exception as e:
                    logger.error(f"DOCX extraction error: {e}")
            os.unlink(tmp.name)
        elif filename.endswith(".pptx"):
            logger.info("Tentative extraction PPTX")
            from pptx import Presentation
            import tempfile
            import tempfile

            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                tmp.write(content)
                tmp.flush()
                logger.info(f"Fichier PPTX temporaire créé : {tmp.name}")
                try:
                    pres = Presentation(tmp.name)
                    text = "\n".join(
                        [shape.text for slide in pres.slides for shape in slide.shapes if hasattr(shape, "text")]
                    )
                    logger.info(
                        f"Texte PPTX extrait : longueur={len(text) if text else 0}, aperçu='{text[:200] if text else ''}'"
                    )
                except Exception as e:
                    logger.error(f"PPTX extraction error: {e}")
            os.unlink(tmp.name)
        elif filename.endswith(".xlsx"):
            logger.info("Tentative extraction XLSX")
            import openpyxl
            import tempfile
            import tempfile

            with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
                tmp.write(content)
                tmp.flush()
                logger.info(f"Fichier XLSX temporaire créé : {tmp.name}")
                try:
                    wb = openpyxl.load_workbook(tmp.name, data_only=True)
                    text = ""
                    for sheet in wb.worksheets:
                        for row in sheet.iter_rows(values_only=True):
                            text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                    logger.info(
                        f"Texte XLSX extrait : longueur={len(text) if text else 0}, aperçu='{text[:200] if text else ''}'"
                    )
                except Exception as e:
                    logger.error(f"XLSX extraction error: {e}")
            os.unlink(tmp.name)
        elif filename.endswith(".txt") or filename.endswith(".csv") or filename.endswith(".ics"):
            logger.info("Tentative extraction fichier texte/csv/ics")
            try:
                text = content.decode("utf-8", errors="ignore")
                logger.info(
                    f"Texte fichier texte/csv/ics extrait : longueur={len(text) if text else 0}, aperçu='{text[:200] if text else ''}'"
                )
            except Exception as e:
                logger.error(f"Text file decode error: {e}")
        else:
            raise HTTPException(status_code=400, detail="File type not supported")

        logger.info(
            f"Texte extrait de la PJ ({file.filename}): longueur={len(text) if text else 0}, aperçu='{text[:200] if text else ''}'"
        )
        if not text or not text.strip():
            raise HTTPException(
                status_code=400,
                detail="Aucun texte détecté dans la pièce jointe. Vérifiez que le document contient du texte sélectionnable (pas une image ou un scan).",
            )

        # Process document with extracted text
        doc_id = process_document_for_user(
            file.filename,
            text.encode("utf-8", errors="ignore"),
            int(user_id),
            db,
            agent_id=None,
            company_id=_get_caller_company_id(user_id, db),
        )

        logger.info(f"Document uploaded for user {user_id}: {file.filename}")
        event_tracker.track_document_upload(int(user_id), file.filename, len(text))

        return {"filename": file.filename, "document_id": doc_id, "status": "uploaded"}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error uploading document: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")


def _process_document_background(
    task_id: str, filename: str, content: bytes, user_id: int, agent_id: int, company_id: int = None
):
    """Background worker for async document processing. Uses its own DB session."""
    db = SessionLocal()
    try:
        # Update status to processing
        r = get_redis()
        if r:
            r.setex(
                f"doc_task:{task_id}",
                3600,
                json.dumps(
                    {
                        "task_id": task_id,
                        "status": "processing",
                        "filename": filename,
                        "document_id": None,
                        "error": None,
                    }
                ),
            )

        doc_id = process_document_for_user(filename, content, user_id, db, agent_id, company_id=company_id)

        logger.info(f"Background processing completed: {filename} -> doc_id={doc_id}")
        event_tracker.track_document_upload(user_id, filename, len(content))

        if r:
            r.setex(
                f"doc_task:{task_id}",
                3600,
                json.dumps(
                    {
                        "task_id": task_id,
                        "status": "completed",
                        "filename": filename,
                        "document_id": doc_id,
                        "error": None,
                    }
                ),
            )
    except Exception as e:
        logger.error(f"Background document processing failed for {filename}: {e}")
        r = get_redis()
        if r:
            r.setex(
                f"doc_task:{task_id}",
                3600,
                json.dumps(
                    {
                        "task_id": task_id,
                        "status": "failed",
                        "filename": filename,
                        "document_id": None,
                        "error": str(e),
                    }
                ),
            )
    finally:
        db.close()


@router.post("/upload-agent")
async def upload_file_for_agent(
    request: Request,
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db),
):
    """Upload and process document for a specific agent (async when Redis available)"""
    if not _check_api_rate_limit(user_id, "upload", _API_UPLOAD_LIMIT):
        raise HTTPException(status_code=429, detail="Too many uploads. Please try again later.")
    try:
        # Get agent_id from form data
        form = await request.form()
        logger.info(f"Form data received in /upload-agent: {dict(form)}")
        agent_id = form.get("agent_id")
        # Fallback: extract agent_id from 'data' if present (Zapier edge case)
        if not agent_id and "data" in form:
            # Try to parse agent_id from string like 'agent_id=23'
            data_value = form.get("data")
            if isinstance(data_value, str) and data_value.startswith("agent_id="):
                agent_id = data_value.split("=", 1)[1]

        if not agent_id:
            logger.error(f"agent_id missing in form: {dict(form)}")
            raise HTTPException(status_code=400, detail="agent_id is required")

        agent_id = int(agent_id)
        # Check file size (10MB limit)
        if file.size > 10 * 1024 * 1024:
            raise HTTPException(status_code=413, detail="File too large (max 10MB)")

        # Check file type
        allowed_types = [".pdf", ".txt", ".docx", ".ics"]
        if not any(file.filename.lower().endswith(ext) for ext in allowed_types):
            raise HTTPException(status_code=400, detail="File type not supported")

        # Verify agent belongs to the user or user has edit permission
        agent = _user_can_edit_agent(int(user_id), agent_id, db)

        content = await file.read()

        # If Redis is available, process in background
        r = get_redis()
        if r is not None:
            task_id = str(uuid4())
            r.setex(
                f"doc_task:{task_id}",
                3600,
                json.dumps(
                    {
                        "task_id": task_id,
                        "status": "processing",
                        "filename": file.filename,
                        "document_id": None,
                        "error": None,
                    }
                ),
            )
            caller_cid = _get_caller_company_id(user_id, db)
            background_tasks.add_task(
                _process_document_background, task_id, file.filename, content, int(user_id), agent_id, caller_cid
            )
            logger.info(f"Document queued for async processing: {file.filename} (task_id={task_id})")
            return {"filename": file.filename, "task_id": task_id, "agent_id": agent_id, "status": "processing"}

        # Fallback: synchronous processing when Redis is unavailable
        doc_id = process_document_for_user(
            file.filename, content, int(user_id), db, agent_id, company_id=_get_caller_company_id(user_id, db)
        )
        logger.info(f"Document uploaded (sync) for user {user_id}, agent {agent_id}: {file.filename}")
        event_tracker.track_document_upload(int(user_id), file.filename, len(content))
        return {"filename": file.filename, "document_id": doc_id, "agent_id": agent_id, "status": "uploaded"}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error uploading document: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")


@router.get("/upload-status/{task_id}")
async def get_upload_status(
    task_id: str,
    user_id: str = Depends(verify_token),
):
    """Poll the status of an async document upload task."""
    r = get_redis()
    if r is None:
        raise HTTPException(status_code=503, detail="Status tracking unavailable")
    data = r.get(f"doc_task:{task_id}")
    if data is None:
        raise HTTPException(status_code=404, detail="Task not found or expired")
    return json.loads(data)


@router.get("/user/documents")
async def get_user_documents(user_id: str = Depends(verify_token), db: Session = Depends(get_db), agent_id: int = None):
    """Get user's documents, optionally filtered by agent"""
    try:
        logger.info(f"Fetching documents for user {user_id}, agent {agent_id}")

        # Build query
        uid = int(user_id)
        if agent_id is not None:
            # Check if user has edit access to this agent (owner or can_edit share)
            share = (
                db.query(AgentShare)
                .filter(AgentShare.agent_id == agent_id, AgentShare.user_id == uid, AgentShare.can_edit == True)
                .first()
            )
            if share:
                # Shared user with edit access: show docs for the agent
                query = db.query(Document).filter(Document.agent_id == agent_id)
            else:
                query = db.query(Document).filter(Document.user_id == uid, Document.agent_id == agent_id)
        else:
            query = db.query(Document).filter(Document.user_id == uid)

        documents = query.order_by(Document.created_at.desc()).all()
        logger.info(f"Found {len(documents)} documents for user {user_id}, agent {agent_id}")

        result = []

        for doc in documents:
            try:
                doc_data = {
                    "id": doc.id,
                    "filename": doc.filename,
                    "created_at": doc.created_at.isoformat(),
                    "gcs_url": doc.gcs_url,
                    "notion_link_id": doc.notion_link_id,
                    "drive_link_id": getattr(doc, "drive_link_id", None),
                    "source_url": getattr(doc, "source_url", None),
                }
                # Safely try to add agent_id if it exists
                if hasattr(doc, "agent_id"):
                    doc_data["agent_id"] = doc.agent_id
                result.append(doc_data)
            except Exception as doc_error:
                logger.error(f"Error processing document {doc.id}: {doc_error}")
                continue

        return {"documents": result}

    except Exception as e:
        logger.error(f"Error fetching documents: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Internal server error")


@router.delete("/documents/{document_id}")
async def delete_document(document_id: int, user_id: str = Depends(verify_token), db: Session = Depends(get_db)):
    """Delete a user's document"""
    try:
        # Check if document exists and belongs to user
        document = db.query(Document).filter(Document.id == document_id, Document.user_id == int(user_id)).first()

        if not document:
            raise HTTPException(status_code=404, detail="Document not found")

        # Delete document
        db.delete(document)
        db.commit()

        logger.info(f"Document {document_id} deleted by user {user_id}")
        event_tracker.track_user_action(int(user_id), f"document_deleted:{document.filename}")

        return {"message": "Document deleted successfully"}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error deleting document: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")


# test
@router.get("/documents/{document_id}/download-url")
async def get_signed_download_url(
    document_id: int, user_id: str = Depends(verify_token), db: Session = Depends(get_db)
):
    """Retourne une URL signée pour télécharger le document depuis GCS"""
    from urllib.parse import urlparse

    _logger = logging.getLogger("main.download_url")

    try:
        document = db.query(Document).filter(Document.id == document_id).first()
        if not document or not document.gcs_url:
            raise HTTPException(status_code=404, detail="Document non trouvé ou pas de fichier GCS")
        # Owner or has access to the agent
        if document.user_id != int(user_id):
            if document.agent_id:
                _user_can_access_agent(int(user_id), document.agent_id, db)
            else:
                raise HTTPException(status_code=403, detail="Access denied")

        gcs_url = document.gcs_url
        _logger.info(f"Generating signed URL for document {document_id}, gcs_url={gcs_url}")

        # Parse bucket and blob name (supports storage.googleapis.com and gs:// formats)
        from urllib.parse import unquote

        parsed = urlparse(gcs_url)
        if gcs_url.startswith("gs://"):
            parts = gcs_url[5:].split("/", 1)
            bucket_name = parts[0]
            blob_name = parts[1] if len(parts) > 1 else ""
        else:
            path = parsed.path.lstrip("/")
            path_parts = path.split("/")
            bucket_name = path_parts[0]
            blob_name_encoded = "/".join(path_parts[1:])
            # URL-decode the blob name (handles %C3%A9, %2B, etc.)
            blob_name = unquote(blob_name_encoded)

        _logger.info(f"Blob name (encoded)={locals().get('blob_name_encoded', None)}, decoded={blob_name}")

        _logger.info(f"Parsed bucket={bucket_name}, blob={blob_name}")

        client = storage.Client()
        bucket = client.bucket(bucket_name)
        blob = bucket.blob(blob_name)

        # Existence check
        try:
            exists = blob.exists()
        except Exception as e:
            _logger.exception("Error checking blob existence (possible permission issue)")
            raise HTTPException(
                status_code=500,
                detail="Erreur lors de la vérification de l'existence du fichier GCS (vérifiez les permissions du service account)",
            )

        if not exists:
            _logger.error(f"Blob not found: {bucket_name}/{blob_name}")
            raise HTTPException(status_code=404, detail="Fichier introuvable dans le bucket GCS")

        try:
            url = blob.generate_signed_url(version="v4", expiration=600, method="GET")
        except Exception as e:
            _logger.exception("Error generating signed URL (permission or signing issue)")
            # Provide a helpful hint without exposing sensitive info
            detail_msg = (
                "Impossible de générer le lien signé. Vérifiez que le service account a les droits GCS "
                "et la capacité de signer des URL (roles/storage.objectViewer et permissions de signature)."
            )
            # Fallback: offer a proxied download endpoint (secure, authenticated)
            proxy_url = f"/documents/{document_id}/download"
            _logger.info(f"Falling back to proxy download for document {document_id}")
            return {"proxy_url": proxy_url, "note": "Signed URL generation failed; using authenticated proxy download."}

        _logger.info(f"Signed URL generated for document {document_id}")
        return {"signed_url": url}

    except HTTPException:
        raise
    except Exception as e:
        tb = traceback.format_exc()
        _logger = logging.getLogger("main.download_url")
        _logger.error(f"Unexpected error generating signed URL for document {document_id}: {e}\n{tb}")
        raise HTTPException(
            status_code=500,
            detail="Erreur interne lors de la génération du lien de téléchargement. Vérifiez les logs du backend.",
        )


@router.get("/documents/{document_id}/download")
async def proxy_download_document(
    document_id: int, user_id: str = Depends(verify_token), db: Session = Depends(get_db)
):
    """Stream the object from GCS through the backend as an authenticated proxy.
    This is a secure fallback when signed URL generation is not possible from the environment.
    """
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document or not document.gcs_url:
        raise HTTPException(status_code=404, detail="Document non trouvé ou pas de fichier GCS")
    # Owner or has access to the agent
    if document.user_id != int(user_id):
        if document.agent_id:
            _user_can_access_agent(int(user_id), document.agent_id, db)
        else:
            raise HTTPException(status_code=403, detail="Access denied")

    from urllib.parse import urlparse, unquote

    gcs_url = document.gcs_url
    parsed = urlparse(gcs_url)
    path = parsed.path.lstrip("/")
    path_parts = path.split("/")
    bucket_name = path_parts[0]
    blob_name_encoded = "/".join(path_parts[1:])
    blob_name = unquote(blob_name_encoded)
    _logger = logging.getLogger("main.download_url")
    _logger.info(f"Proxy download: bucket={bucket_name}, blob_encoded={blob_name_encoded}, blob_decoded={blob_name}")

    client = storage.Client()
    bucket = client.bucket(bucket_name)

    # Attempt to get blob and check existence
    def get_blob(name: str):
        return bucket.blob(name)

    blob = get_blob(blob_name)

    # Existence check (may raise if permission issues)
    exists = None
    try:
        exists = blob.exists()
    except Exception:
        _logger.exception(f"Error checking existence for blob {bucket_name}/{blob_name} (possible permission issue)")
        # keep exists as None and try download below

    if exists is False:
        _logger.error(f"Blob not found for proxy: {bucket_name}/{blob_name}")
        raise HTTPException(status_code=404, detail="Fichier introuvable dans le bucket GCS")

    data = None
    # Try direct download
    try:
        data = blob.download_as_bytes()
    except Exception as exc:
        _logger.exception(f"Initial download attempt failed for {bucket_name}/{blob_name}: {exc}")

        # If download failed, try unicode normalization variants (NFC/NFD)
        try:
            import unicodedata

            tried = []
            for norm in ("NFC", "NFD"):
                alt_name = unicodedata.normalize(norm, blob_name)
                if alt_name in tried or alt_name == blob_name:
                    continue
                tried.append(alt_name)
                _logger.info(f"Retrying download with normalized blob name ({norm}): {alt_name}")
                alt_blob = get_blob(alt_name)
                try:
                    data = alt_blob.download_as_bytes()
                    # successful: use this blob
                    blob = alt_blob
                    blob_name = alt_name
                    _logger.info(f"Download succeeded with normalized name ({norm})")
                    break
                except Exception as exc2:
                    _logger.exception(f"Download failed with normalized name {alt_name}: {exc2}")
        except Exception as norm_exc:
            _logger.exception(f"Error during unicode normalization retries: {norm_exc}")

    if data is None:
        # Determine if likely permission issue vs not found
        # If exists is None, we couldn't determine existence due to permission; respond with 403 hint
        if exists is None:
            _logger.error(
                f"Download failed and existence unknown for {bucket_name}/{blob_name}. Likely permission issue."
            )
            raise HTTPException(
                status_code=403,
                detail="Le service n'a pas les permissions nécessaires pour lire l'objet GCS. Vérifiez roles/storage.objectViewer.",
            )
        else:
            _logger.error(f"All attempts to download blob failed for {bucket_name}/{blob_name}")
            raise HTTPException(status_code=500, detail="Impossible de récupérer le fichier depuis GCS")

    # Guess mimetype
    mime, _ = mimetypes.guess_type(document.filename)
    if not mime:
        mime = "application/octet-stream"

    from io import BytesIO

    # Ensure filename is safe; use the stored document filename
    safe_filename = document.filename or os.path.basename(blob_name)
    headers = {"Content-Disposition": f'attachment; filename="{safe_filename}"'}
    return StreamingResponse(BytesIO(data), media_type=mime, headers=headers)
