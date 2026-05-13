"""Email ingestion endpoints: extractText, email ingest, attachment upload."""

import json
import hmac
import os
import logging
import traceback

from fastapi import APIRouter, Depends, File, HTTPException, Request, UploadFile
from sqlalchemy.orm import Session
from typing import List

from auth import verify_token
from database import get_db, Agent, Document, DocumentChunk, set_current_company_id, engine
from helpers.rate_limiting import _check_api_rate_limit, _API_EXTRACT_LIMIT
from schemas.email_ingest import EmailIngestRequest

logger = logging.getLogger(__name__)
router = APIRouter()


# ============================================================================
# HELPER FUNCTIONS
# ============================================================================


def extract_email_tags_from_title(title: str) -> List[str]:
    """Extrait les @tags du titre de l'email (case-insensitive)"""
    import re

    pattern = r"@([a-zA-Z0-9_-]+)"
    matches = re.findall(pattern, title)
    # Normaliser en lowercase avec @
    return [f"@{tag.lower()}" for tag in matches]


class _MatchedAgent:
    """Lightweight stand-in for Agent ORM objects returned by find_agents_by_email_tags.

    We avoid returning ORM-managed Agent instances because the session's RLS
    context changes (service_bypass → company_id) and SQLAlchemy may try to
    lazy-reload attributes, causing 'row not present' errors.
    """

    __slots__ = ("id", "name", "user_id", "company_id")

    def __init__(self, id, name, user_id, company_id):
        self.id = id
        self.name = name
        self.user_id = user_id
        self.company_id = company_id


def find_agents_by_email_tags(db: Session, tags: List[str]):
    """Trouve tous les agents dont email_tags contient au moins un des tags.

    Uses a direct engine connection to bypass RLS entirely — email ingestion
    is authenticated by API key (not JWT) so no tenant context exists.
    A direct connection from engine.connect() is not subject to the
    SessionLocal after_begin listener and we explicitly SET LOCAL
    app.service_bypass = 'true' for safety.

    Returns _MatchedAgent objects (not ORM instances) and sets the correct
    tenant context on the *session* for subsequent operations.
    """
    if not tags:
        return []

    lower_tags = [t.lower() for t in tags]

    from sqlalchemy import text

    # Use a direct engine connection to query across ALL companies.
    # This avoids any RLS filtering that the session might impose.
    # IMPORTANT: wrap in conn.begin() so SET LOCAL takes effect within a transaction.
    try:
        with engine.connect() as conn:
            with conn.begin():
                conn.execute(text("SET LOCAL app.service_bypass = 'true'"))
                rows = conn.execute(
                    text("SELECT id, name, user_id, email_tags, company_id FROM agents "
                         "WHERE email_tags IS NOT NULL AND LENGTH(email_tags) > 2")
                ).fetchall()
    except Exception as e:
        print(f"[EMAIL_INGEST] Failed to query agents for email_tags: {e}", flush=True)
        logger.error(f"Failed to query agents for email_tags: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        return []

    matched_agents = []
    matched_company_id = None
    print(f"[EMAIL_INGEST] [TAG_MATCH] Searching for tags {lower_tags} across {len(rows)} agents with email_tags", flush=True)
    logger.info(f"[TAG_MATCH] Searching for tags {lower_tags} across {len(rows)} agents with email_tags")
    for row in rows:
        agent_id, agent_name, user_id, email_tags_raw, company_id = row
        try:
            agent_tags = json.loads(email_tags_raw) if isinstance(email_tags_raw, str) else []
            agent_tags_lower = [t.lower() for t in agent_tags if isinstance(t, str)]
            logger.info(f"[TAG_MATCH] Agent {agent_id} ({agent_name}): tags={agent_tags_lower}, looking_for={lower_tags}, match={any(tag in agent_tags_lower for tag in lower_tags)}")
            if any(tag in agent_tags_lower for tag in lower_tags):
                matched_agents.append(_MatchedAgent(agent_id, agent_name, user_id, company_id))
                if company_id is not None:
                    matched_company_id = company_id
        except (json.JSONDecodeError, TypeError) as e:
            logger.warning(f"[TAG_MATCH] Agent {agent_id}: failed to parse email_tags '{email_tags_raw}': {e}")
            continue

    if not matched_agents:
        print(f"[EMAIL_INGEST] No agents matched tags {lower_tags} (checked {len(rows)} agents with email_tags)", flush=True)
        logger.info(f"No agents matched tags {lower_tags} (checked {len(rows)} agents with email_tags)")
        return []

    # Set the tenant context on the request session so subsequent ORM
    # operations (document insert, dedup check, etc.) work under RLS.
    if matched_company_id is not None:
        set_current_company_id(matched_company_id)
        db.execute(text("SET LOCAL app.company_id = :cid"), {"cid": str(int(matched_company_id))})
        logger.info(f"Set tenant context to company_id={matched_company_id} for email ingestion")

    return matched_agents



def verify_email_api_key(request: Request) -> bool:
    """Vérifie l'API Key pour l'ingestion d'emails

    Security: Uses constant-time comparison (hmac.compare_digest) to prevent
    timing attacks that could be used to guess the API key character by character.
    """
    api_key = request.headers.get("X-API-Key", "")
    expected_key = os.getenv("EMAIL_INGEST_API_KEY", "")

    if not expected_key:
        logger.error("EMAIL_INGEST_API_KEY not configured in environment")
        raise HTTPException(status_code=500, detail="API Key not configured on server")

    if not api_key:
        logger.warning("Missing API Key for email ingestion")
        raise HTTPException(status_code=401, detail="Invalid API Key")

    # Constant-time comparison to prevent timing attacks
    if not hmac.compare_digest(api_key, expected_key):
        logger.warning("Invalid API Key attempt for email ingestion")
        raise HTTPException(status_code=401, detail="Invalid API Key")

    return True


# ============================================================================
# ENDPOINTS
# ============================================================================


# Endpoint pour extraire le texte des fichiers uploadés (PDF, TXT, DOCX, XLSX, PPTX, etc.)


# Nouvelle version : accepte un seul fichier UploadFile
@router.post("/api/agent/extractText")
async def extract_text_from_file(file: UploadFile = File(...), user_id: str = Depends(verify_token)):
    """Extrait le texte d'un fichier uploadé et renvoie le texte extrait.

    Security: Requires authentication to prevent unauthorized file processing.
    """
    if not _check_api_rate_limit(user_id, "extract", _API_EXTRACT_LIMIT):
        raise HTTPException(status_code=429, detail="Too many extraction requests. Please try again later.")
    logger = logging.getLogger("extractText")
    ext = file.filename.lower().split(".")[-1]
    logger.info(
        f"Appel reçu sur /api/agent/extractText : filename={file.filename}, ext={ext}, content_type={getattr(file, 'content_type', 'unknown')}"
    )
    text = ""
    try:
        if ext == "pdf":
            MAX_PDF_PAGES = 500
            logger.info(f"Tentative extraction PDF: {file.filename}")
            try:
                import pdfplumber

                with pdfplumber.open(file.file) as pdf:
                    if len(pdf.pages) > MAX_PDF_PAGES:
                        raise HTTPException(
                            status_code=400, detail=f"PDF too large ({len(pdf.pages)} pages, max {MAX_PDF_PAGES})"
                        )
                    for i, page in enumerate(pdf.pages):
                        page_text = page.extract_text()
                        logger.info(
                            f"Page {i + 1} PDF: longueur={len(page_text) if page_text else 0}, aperçu='{page_text[:100] if page_text else ''}'"
                        )
                        if page_text:
                            text += page_text + "\n"
            except Exception as e:
                logger.error(f"Erreur extraction PDF {file.filename}: {e}")
        elif ext in ["txt", "md", "json", "xml", "csv"]:
            logger.info(f"Tentative extraction texte: {file.filename}")
            try:
                raw = await file.read()
                text = raw.decode(errors="ignore")
                logger.info(f"Texte extrait: longueur={len(text)}, aperçu='{text[:200]}'")
            except Exception as e:
                logger.error(f"Erreur extraction texte {file.filename}: {e}")
        elif ext in ["doc", "docx"]:
            logger.info(f"Tentative extraction DOCX: {file.filename}")
            try:
                from docx import Document as DocxDocument

                doc = DocxDocument(file.file)
                text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                logger.info(f"Texte DOCX extrait: longueur={len(text)}, aperçu='{text[:200]}'")
            except Exception as e:
                logger.error(f"Erreur extraction DOCX {file.filename}: {e}")
        elif ext in ["xls", "xlsx"]:
            logger.info(f"Tentative extraction XLSX: {file.filename}")
            try:
                import openpyxl

                wb = openpyxl.load_workbook(file.file, read_only=True)
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                        text += row_text + "\n"
                logger.info(f"Texte XLSX extrait: longueur={len(text)}, aperçu='{text[:200]}'")
            except Exception as e:
                logger.error(f"Erreur extraction XLSX {file.filename}: {e}")
        elif ext in ["ppt", "pptx"]:
            logger.info(f"Tentative extraction PPTX: {file.filename}")
            try:
                from pptx import Presentation

                pres = Presentation(file.file)
                for slide in pres.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                logger.info(f"Texte PPTX extrait: longueur={len(text)}, aperçu='{text[:200]}'")
            except Exception as e:
                logger.error(f"Erreur extraction PPTX {file.filename}: {e}")
        else:
            logger.warning(f"Type de fichier non supporté: {file.filename}")
            text = f"[Type de fichier non supporté: {file.filename}]"
    except Exception as e:
        logger.error(f"Erreur extraction {file.filename}: {e}")
        text = f"[Erreur extraction {file.filename}: {e}]"
    logger.info(f"Résultat extraction {file.filename}: longueur={len(text.strip())}, aperçu='{text.strip()[:200]}'")
    return {"text": text.strip()}


# ============================================================================
# EMAIL INGESTION API - Cloud Function Integration
# ============================================================================


@router.post("/api/emails/ingest")
async def ingest_email(payload: EmailIngestRequest, request: Request, db: Session = Depends(get_db)):
    """
    Ingère un email depuis la Cloud Function Gmail.

    Authentification via X-API-Key header.
    Route l'email vers les agents basé sur les @tags dans le titre.
    Si aucun tag trouvé ou aucun agent matché, ne rien faire.
    """
    # Vérifier l'API Key
    verify_email_api_key(request)

    try:
        logger.info(f"Ingesting email: {payload.title} (source_id: {payload.source_id})")

        # 1. Extraire les @tags du titre
        extracted_tags = extract_email_tags_from_title(payload.title)
        logger.info(f"Extracted tags from title: {extracted_tags}")

        # 2. Trouver les agents correspondants
        if extracted_tags:
            target_agents = find_agents_by_email_tags(db, extracted_tags)
            logger.info(f"Found {len(target_agents)} agents matching tags: {[a.name for a in target_agents]}")
        else:
            target_agents = []
            logger.info("No tags found in email title")

        # 3. Si aucun tag ou aucun agent matché, ne rien faire
        if not target_agents:
            logger.info(f"No matching agents for email: {payload.title} - skipping ingestion")
            return {
                "success": True,
                "document_ids": [],
                "agents_matched": 0,
                "message": "Aucun tag trouvé ou aucun companion correspondant - email ignoré",
            }

        # 4. Construire le contenu enrichi avec les métadonnées
        enriched_content = f"Sujet: {payload.title}\n"
        if payload.metadata:
            if payload.metadata.from_email:
                enriched_content += f"De: {payload.metadata.from_email}\n"
            if payload.metadata.date:
                enriched_content += f"Date: {payload.metadata.date}\n"
        enriched_content += f"\n{payload.content}"

        # Importer les fonctions nécessaires
        from file_loader import chunk_text
        from mistral_embeddings import get_embedding_fast
        import numpy as np

        # Fonction pour découper les gros chunks
        def split_for_embedding(chunk, max_tokens=8192):
            chunk = chunk.replace("\x00", "")
            max_chars = max_tokens * 4
            return [chunk[i : i + max_chars] for i in range(0, len(chunk), max_chars)]

        # Préparer les chunks une seule fois (réutilisés pour chaque agent)
        chunks = chunk_text(enriched_content)
        logger.info(f"Created {len(chunks)} chunks for email")

        # Calculer les embeddings une seule fois avec Mistral
        chunk_embeddings = []
        max_immediate_chunks = 20
        for i, chunk in enumerate(chunks):
            if i < max_immediate_chunks:
                try:
                    sub_chunks = split_for_embedding(chunk, 8192)
                    embeddings = []
                    for sub in sub_chunks:
                        embedding = get_embedding_fast(sub)
                        embeddings.append(embedding)
                    if embeddings:
                        avg_embedding = list(np.mean(np.array(embeddings), axis=0))
                    else:
                        raise ValueError("No sub-chunks produced for embedding")
                except Exception as e:
                    logger.error(f"Failed to get Mistral embedding for chunk {i}: {e}")
                    raise
            else:
                avg_embedding = None
            chunk_embeddings.append(avg_embedding)

        # 5. Créer un document pour CHAQUE agent trouvé
        from sqlalchemy import text as _text
        document_ids = []
        for agent in target_agents:
            # Dedup via direct engine connection (bypasses RLS entirely)
            unique_id = f"email_{payload.source_id}_agent_{agent.id}"
            trace_unique_id = f"email_trace_{payload.source_id}_agent_{agent.id}"

            rag_exists = False
            trace_exists = False
            existing_rag_id = None

            with engine.connect() as conn:
                with conn.begin():
                    conn.execute(_text("SET LOCAL app.service_bypass = 'true'"))
                    row = conn.execute(
                        _text("SELECT id FROM documents WHERE source_url = :uid LIMIT 1"),
                        {"uid": unique_id}
                    ).first()
                    if row:
                        rag_exists = True
                        existing_rag_id = row[0]
                    trace_row = conn.execute(
                        _text("SELECT id FROM documents WHERE source_url = :uid LIMIT 1"),
                        {"uid": trace_unique_id}
                    ).first()
                    if trace_row:
                        trace_exists = True

            if rag_exists and trace_exists:
                logger.info(f"[DEDUP] Email already fully ingested for agent {agent.name}: {payload.source_id}")
                document_ids.append(existing_rag_id)
                continue

            if rag_exists:
                logger.info(f"[DEDUP] RAG doc exists for agent {agent.name}, checking traceability")
                document_ids.append(existing_rag_id)
            else:
                # Créer le document RAG
                document = Document(
                    filename=payload.title,
                    content=enriched_content,
                    user_id=agent.user_id,
                    agent_id=agent.id,
                    company_id=agent.company_id,
                    source_url=unique_id,
                )
                db.add(document)
                db.commit()
                db.refresh(document)

                logger.info(f"RAG document created for agent {agent.name} with ID: {document.id}")

                # Créer les chunks avec les embeddings pré-calculés
                for i, chunk in enumerate(chunks):
                    doc_chunk = DocumentChunk(
                        document_id=document.id,
                        company_id=agent.company_id,
                        chunk_text=chunk,
                        embedding_vec=chunk_embeddings[i] if chunk_embeddings[i] else None,
                        chunk_index=i,
                    )
                    db.add(doc_chunk)

                db.commit()
                document_ids.append(document.id)

            # Document de traçabilité
            if not trace_exists:
                trace_doc = Document(
                    filename=f"[Email] {payload.title}",
                    content=enriched_content,
                    user_id=agent.user_id,
                    agent_id=agent.id,
                    company_id=agent.company_id,
                    source_url=trace_unique_id,
                    document_type="traceability",
                )
                db.add(trace_doc)
                db.commit()

        logger.info(f"Email ingested successfully to {len(document_ids)} agents: {payload.title}")

        return {
            "success": True,
            "document_ids": document_ids,
            "agents_matched": len(target_agents),
            "agents": [{"id": a.id, "name": a.name} for a in target_agents],
            "tags_extracted": extracted_tags,
            "message": f"Email ingéré avec succès vers {len(target_agents)} companion(s)",
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error ingesting email: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        db.rollback()
        raise HTTPException(status_code=500, detail="Erreur lors de l'ingestion de l'email")


@router.post("/api/emails/upload-attachment")
async def upload_email_attachment(request: Request, file: UploadFile = File(...), db: Session = Depends(get_db)):
    """
    Upload une pièce jointe d'email comme document séparé.
    Le fichier est stocké sur GCS et devient téléchargeable.
    Route vers les agents basé sur les @tags passés dans le form data.

    Authentification via X-API-Key header.

    Creates exactly:
      - 1 RAG document  (chunked + embedded, filename prefixed with [Email PJ])
      - 1 traceability document (full text, same prefix, document_type='traceability')
    per matched agent.  Dedup uses source_id (Gmail message ID) + filename + agent_id
    via a direct engine connection to bypass RLS entirely.
    """
    verify_email_api_key(request)

    try:
        form = await request.form()
        email_subject = form.get("email_subject", "")
        source_id = form.get("source_id", "")

        logger.info(f"Uploading attachment: {file.filename} for email: {email_subject} (source_id: {source_id})")

        # --- 1. Find target agents by @tags (cross-tenant, uses engine.connect) ---
        extracted_tags = extract_email_tags_from_title(email_subject) if email_subject else []
        logger.info(f"Extracted tags from email subject: {extracted_tags}")

        target_agents = find_agents_by_email_tags(db, extracted_tags) if extracted_tags else []

        if not target_agents:
            logger.info(f"No matching agents for attachment: {file.filename} - skipping")
            return {
                "success": True,
                "document_ids": [],
                "agents_matched": 0,
                "message": "Aucun companion correspondant - pièce jointe ignorée",
            }

        # --- 2. Validate file ---
        allowed_types = [".pdf", ".txt", ".docx", ".json"]
        if not any(file.filename.lower().endswith(ext) for ext in allowed_types):
            logger.warning(f"Unsupported file type: {file.filename}")
            return {"success": False, "message": f"Type de fichier non supporté: {file.filename}"}

        content = await file.read()
        if len(content) > 10 * 1024 * 1024:
            raise HTTPException(status_code=413, detail="Fichier trop volumineux (max 10MB)")

        # --- 3. Process each agent ---
        from rag_engine import process_document_for_user
        from sqlalchemy import text as _text

        prefixed_filename = f"[Email PJ] {file.filename}"
        document_ids = []

        for agent in target_agents:
            try:
                # Stable dedup keys stored in source_url (gcs_url keeps the real GCS path)
                rag_key = f"email_pj_{source_id}_{file.filename}_agent_{agent.id}" if source_id else ""
                trace_key = f"email_pj_trace_{source_id}_{file.filename}_agent_{agent.id}" if source_id else ""

                # Single connection to check both RAG and traceability existence
                rag_exists = False
                trace_exists = False
                existing_rag_id = None

                if rag_key:
                    with engine.connect() as conn:
                        with conn.begin():
                            conn.execute(_text("SET LOCAL app.service_bypass = 'true'"))
                            row = conn.execute(
                                _text("SELECT id FROM documents WHERE source_url = :key LIMIT 1"),
                                {"key": rag_key}
                            ).first()
                            if row:
                                rag_exists = True
                                existing_rag_id = row[0]
                            trace_row = conn.execute(
                                _text("SELECT id FROM documents WHERE source_url = :key LIMIT 1"),
                                {"key": trace_key}
                            ).first()
                            if trace_row:
                                trace_exists = True

                if rag_exists and trace_exists:
                    logger.info(f"[DEDUP] Both docs exist for agent {agent.name}: {prefixed_filename}")
                    document_ids.append(existing_rag_id)
                    continue

                # Create RAG document if needed
                if not rag_exists:
                    doc_id = process_document_for_user(
                        filename=prefixed_filename, content=content, user_id=agent.user_id,
                        db=db, agent_id=agent.id, company_id=agent.company_id,
                    )
                    document_ids.append(doc_id)
                    logger.info(f"RAG doc created for agent {agent.name}: {prefixed_filename} (doc_id: {doc_id})")

                    # Store dedup key in source_url (gcs_url already has the real GCS path)
                    if rag_key:
                        db.execute(
                            _text("UPDATE documents SET source_url = :key WHERE id = :did"),
                            {"key": rag_key, "did": doc_id}
                        )
                        db.commit()
                else:
                    document_ids.append(existing_rag_id)
                    logger.info(f"[DEDUP] RAG doc exists for agent {agent.name} (doc_id: {existing_rag_id})")

                # Create traceability document if needed
                if not trace_exists:
                    trace_text = ""
                    try:
                        if file.filename.lower().endswith(".pdf"):
                            import tempfile as _tempfile
                            with _tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                                tmp.write(content)
                                tmp_path = tmp.name
                            try:
                                from file_loader import load_text_from_pdf
                                trace_text = load_text_from_pdf(tmp_path) or ""
                            finally:
                                import os as _os
                                _os.unlink(tmp_path)
                        else:
                            trace_text = content.decode("utf-8", errors="replace")
                    except Exception:
                        trace_text = ""

                    trace_doc = Document(
                        filename=prefixed_filename,
                        content=trace_text,
                        user_id=agent.user_id,
                        agent_id=agent.id,
                        company_id=agent.company_id,
                        source_url=trace_key,
                        document_type="traceability",
                    )
                    db.add(trace_doc)
                    db.commit()
                    logger.info(f"Traceability doc created for agent {agent.name}: {prefixed_filename}")

            except Exception as e:
                logger.error(f"Failed to process attachment for agent {agent.name}: {e}")
                logger.error(f"Traceback: {traceback.format_exc()}")
                continue

        return {
            "success": True,
            "document_ids": document_ids,
            "agents_matched": len(target_agents),
            "agents": [{"id": a.id, "name": a.name} for a in target_agents],
            "filename": file.filename,
            "message": f"Pièce jointe uploadée vers {len(document_ids)} companion(s)",
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error uploading attachment: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Erreur lors de l'upload de la pièce jointe: {e}")
