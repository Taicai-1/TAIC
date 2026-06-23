"""Missions automation endpoints: CRUD, planning parse/events, documents, recaps, chat."""

import json
import logging
from datetime import datetime
from uuid import uuid4

from fastapi import APIRouter, BackgroundTasks, Depends, File, HTTPException, Query, UploadFile
from sqlalchemy.orm import Session

from auth import verify_token
from database import Agent, Document, Mission, MissionEvent, MissionRecap, MissionRecapSchedule, get_db
from permissions import require_role
from schemas.missions import (
    EventCreate,
    EventsBulk,
    EventUpdate,
    MissionChatRequest,
    MissionCreate,
    MissionUpdate,
    RecapScheduleCreate,
    RecapScheduleUpdate,
)

logger = logging.getLogger(__name__)
router = APIRouter()

MAX_FILE_SIZE = 25 * 1024 * 1024
ALLOWED_EXT = (".pdf", ".txt", ".csv", ".docx", ".xlsx", ".pptx", ".json")


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #


def _get_mission_or_404(mission_id: int, user_id: int, company_id: int, db: Session) -> Mission:
    """Fetch a mission scoped to its creator + company (private to creator)."""
    mission = (
        db.query(Mission)
        .filter(Mission.id == mission_id, Mission.company_id == company_id, Mission.user_id == user_id)
        .first()
    )
    if not mission:
        raise HTTPException(status_code=404, detail="Mission not found")
    return mission


def _validate_agent(agent_id, company_id: int, db: Session):
    """Verify the companion belongs to the company; return it or None."""
    if agent_id is None:
        return None
    agent = db.query(Agent).filter(Agent.id == agent_id, Agent.company_id == company_id).first()
    if not agent:
        raise HTTPException(status_code=400, detail="Companion introuvable dans votre organisation")
    return agent


def _mission_detail(mission: Mission, db: Session) -> dict:
    last_recap = (
        db.query(MissionRecap)
        .filter(MissionRecap.mission_id == mission.id)
        .order_by(MissionRecap.created_at.desc())
        .first()
    )
    return {
        "id": mission.id,
        "name": mission.name,
        "objective": mission.objective,
        "agent_id": mission.agent_id,
        "status": mission.status,
        "recap_enabled": mission.recap_enabled,
        "recap_weekday": mission.recap_weekday,
        "recap_hour": mission.recap_hour,
        "recap_prompt": mission.recap_prompt,
        "created_at": mission.created_at.isoformat() if mission.created_at else None,
        "last_recap_at": last_recap.created_at.isoformat() if last_recap and last_recap.created_at else None,
    }


# --------------------------------------------------------------------------- #
# Mission CRUD
# --------------------------------------------------------------------------- #


@router.get("/api/automations/missions")
async def list_missions(user_id: int = Depends(verify_token), db: Session = Depends(get_db)):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    missions = (
        db.query(Mission)
        .filter(Mission.company_id == membership.company_id, Mission.user_id == user_id)
        .order_by(Mission.created_at.desc())
        .all()
    )
    return {"missions": [_mission_detail(m, db) for m in missions]}


@router.post("/api/automations/missions")
async def create_mission(body: MissionCreate, user_id: int = Depends(verify_token), db: Session = Depends(get_db)):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    _validate_agent(body.agent_id, membership.company_id, db)

    mission = Mission(
        company_id=membership.company_id,
        user_id=user_id,
        agent_id=body.agent_id,
        name=body.name.strip(),
        objective=body.objective.strip(),
        recap_enabled=body.recap_enabled,
        recap_weekday=body.recap_weekday,
        recap_hour=body.recap_hour,
    )
    db.add(mission)
    db.commit()
    db.refresh(mission)
    return {"mission": _mission_detail(mission, db)}


@router.get("/api/automations/missions/{mission_id}")
async def get_mission(mission_id: int, user_id: int = Depends(verify_token), db: Session = Depends(get_db)):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    return {"mission": _mission_detail(mission, db)}


@router.put("/api/automations/missions/{mission_id}")
async def update_mission(
    mission_id: int, body: MissionUpdate, user_id: int = Depends(verify_token), db: Session = Depends(get_db)
):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    _validate_agent(body.agent_id, membership.company_id, db)

    mission.name = body.name.strip()
    mission.objective = body.objective.strip()
    mission.agent_id = body.agent_id
    mission.status = body.status
    mission.recap_enabled = body.recap_enabled
    mission.recap_weekday = body.recap_weekday
    mission.recap_hour = body.recap_hour
    mission.recap_prompt = body.recap_prompt
    mission.updated_at = datetime.utcnow()
    db.commit()
    db.refresh(mission)
    return {"mission": _mission_detail(mission, db)}


@router.delete("/api/automations/missions/{mission_id}")
async def delete_mission(mission_id: int, user_id: int = Depends(verify_token), db: Session = Depends(get_db)):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    db.delete(mission)
    db.commit()
    return {"success": True}


# --------------------------------------------------------------------------- #
# Planning parse + events
# --------------------------------------------------------------------------- #


def _extract_text_from_upload(filename: str, content: bytes, db: Session, user_id: int) -> str:
    """Extract raw text from an uploaded planning file using the existing loaders."""
    import os
    import tempfile

    name = filename.lower()
    if name.endswith(".pdf"):
        from file_loader import load_text_from_pdf

        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp_path = tmp.name
                tmp.write(content)
            return load_text_from_pdf(tmp_path) or ""
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.unlink(tmp_path)
    if name.endswith(".docx"):
        from docx import Document as DocxDocument

        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            tmp.write(content)
            tmp.flush()
            path = tmp.name
        try:
            doc = DocxDocument(path)
            return "\n".join(p.text for p in doc.paragraphs)
        finally:
            os.unlink(path)
    if name.endswith(".xlsx"):
        import openpyxl

        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
            tmp.write(content)
            tmp.flush()
            path = tmp.name
        try:
            wb = openpyxl.load_workbook(path, data_only=True)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += "\t".join(str(c) if c is not None else "" for c in row) + "\n"
            return text
        finally:
            os.unlink(path)
    if name.endswith(".pptx"):
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(content)
            tmp.flush()
            path = tmp.name
        try:
            pres = Presentation(path)
            return "\n".join(shape.text for slide in pres.slides for shape in slide.shapes if hasattr(shape, "text"))
        finally:
            os.unlink(path)
    # txt / csv / json
    return content.decode("utf-8", errors="ignore")


def _parse_planning_text(raw_text: str, agent, db: Session) -> tuple[list, int]:
    """LLM-parse raw planning text into events. Returns (valid_events, skipped_count)."""
    from openai_client import get_chat_response_json
    from schemas.missions import ParsedEvent
    from weekly_recap import get_model_id_for_agent

    current_year = datetime.utcnow().year
    system = (
        "Tu extrais des évènements datés depuis un planning brut. "
        "Réponds UNIQUEMENT avec un tableau JSON d'objets "
        '{"date": "YYYY-MM-DD", "title": "...", "description": "..."}. '
        f"Utilise l'année {current_year} par défaut si l'année est absente. "
        "Les dates DOIVENT être au format ISO YYYY-MM-DD. "
        "Ignore tout ce qui n'est pas un évènement daté. "
        "Si aucun évènement, réponds []."
    )
    messages = [
        {"role": "system", "content": system},
        {"role": "user", "content": raw_text[:20000]},
    ]
    model_id = get_model_id_for_agent(agent) if agent else "mistral:mistral-large-latest"

    parsed = get_chat_response_json(messages, model_id=model_id)
    if isinstance(parsed, dict):
        for key in ("events", "data", "items", "result"):
            if isinstance(parsed.get(key), list):
                parsed = parsed[key]
                break
        else:
            parsed = []
    if not isinstance(parsed, list):
        parsed = []

    valid, skipped = [], 0
    for item in parsed:
        try:
            ev = ParsedEvent(**item)
            valid.append({"date": ev.date.isoformat(), "title": ev.title, "description": ev.description})
        except Exception:
            skipped += 1
    return valid, skipped


@router.post("/api/automations/missions/{mission_id}/planning/parse")
async def parse_planning(
    mission_id: int,
    file: UploadFile = File(...),
    user_id: int = Depends(verify_token),
    db: Session = Depends(get_db),
):
    """Parse an uploaded planning into events WITHOUT writing anything."""
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    if mission.status != "active":
        raise HTTPException(status_code=400, detail="Mission archivée : modification impossible")

    if not any(file.filename.lower().endswith(ext) for ext in ALLOWED_EXT):
        raise HTTPException(status_code=400, detail="Type de fichier non supporté")
    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="Fichier trop volumineux")

    raw_text = _extract_text_from_upload(file.filename, content, db, user_id)
    if not raw_text or not raw_text.strip():
        raise HTTPException(status_code=400, detail="Aucun texte détecté dans le fichier")

    agent = db.query(Agent).filter(Agent.id == mission.agent_id).first() if mission.agent_id else None
    events, skipped = _parse_planning_text(raw_text, agent, db)
    if not events:
        raise HTTPException(
            status_code=422,
            detail="Aucun évènement daté n'a pu être extrait. Vous pouvez les saisir manuellement.",
        )
    return {"events": events, "skipped": skipped}


@router.get("/api/automations/missions/{mission_id}/events")
async def list_events(mission_id: int, user_id: int = Depends(verify_token), db: Session = Depends(get_db)):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    rows = (
        db.query(MissionEvent)
        .filter(MissionEvent.mission_id == mission.id)
        .order_by(MissionEvent.event_date.asc())
        .all()
    )
    return {
        "events": [
            {
                "id": e.id,
                "date": e.event_date.isoformat(),
                "title": e.title,
                "description": e.description,
                "source": e.source,
            }
            for e in rows
        ]
    }


@router.post("/api/automations/missions/{mission_id}/events/bulk")
async def bulk_events(
    mission_id: int, body: EventsBulk, user_id: int = Depends(verify_token), db: Session = Depends(get_db)
):
    """Insert validated events (source='upload'). replace_upload purges prior uploads first."""
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    if mission.status != "active":
        raise HTTPException(status_code=400, detail="Mission archivée : modification impossible")

    if body.replace_upload:
        db.query(MissionEvent).filter(MissionEvent.mission_id == mission.id, MissionEvent.source == "upload").delete(
            synchronize_session=False
        )

    for ev in body.events:
        db.add(
            MissionEvent(
                mission_id=mission.id,
                company_id=mission.company_id,
                event_date=ev.date,
                title=ev.title.strip(),
                description=ev.description,
                source="upload",
            )
        )
    db.commit()
    return {"inserted": len(body.events)}


@router.post("/api/automations/missions/{mission_id}/events")
async def create_event(
    mission_id: int, body: EventCreate, user_id: int = Depends(verify_token), db: Session = Depends(get_db)
):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    if mission.status != "active":
        raise HTTPException(status_code=400, detail="Mission archivée : modification impossible")
    event = MissionEvent(
        mission_id=mission.id,
        company_id=mission.company_id,
        event_date=body.date,
        title=body.title.strip(),
        description=body.description,
        source="manual",
    )
    db.add(event)
    db.commit()
    db.refresh(event)
    return {"id": event.id}


@router.put("/api/automations/missions/{mission_id}/events/{event_id}")
async def update_event(
    mission_id: int,
    event_id: int,
    body: EventUpdate,
    user_id: int = Depends(verify_token),
    db: Session = Depends(get_db),
):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    if mission.status != "active":
        raise HTTPException(status_code=400, detail="Mission archivée : modification impossible")
    event = db.query(MissionEvent).filter(MissionEvent.id == event_id, MissionEvent.mission_id == mission.id).first()
    if not event:
        raise HTTPException(status_code=404, detail="Event not found")
    event.event_date = body.date
    event.title = body.title.strip()
    event.description = body.description
    db.commit()
    return {"success": True}


@router.delete("/api/automations/missions/{mission_id}/events/{event_id}")
async def delete_event(
    mission_id: int, event_id: int, user_id: int = Depends(verify_token), db: Session = Depends(get_db)
):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    if mission.status != "active":
        raise HTTPException(status_code=400, detail="Mission archivée : modification impossible")
    event = db.query(MissionEvent).filter(MissionEvent.id == event_id, MissionEvent.mission_id == mission.id).first()
    if not event:
        raise HTTPException(status_code=404, detail="Event not found")
    db.delete(event)
    db.commit()
    return {"success": True}


# --------------------------------------------------------------------------- #
# Recap schedules (recurring + one-shot)
# --------------------------------------------------------------------------- #


def _schedule_detail(s: MissionRecapSchedule) -> dict:
    return {
        "id": s.id,
        "kind": s.kind,
        "weekday": s.weekday,
        "run_date": s.run_date.isoformat() if s.run_date else None,
        "hour": s.hour,
        "enabled": s.enabled,
        "last_run_at": s.last_run_at.isoformat() if s.last_run_at else None,
    }


@router.get("/api/automations/missions/{mission_id}/recap-schedules")
async def list_recap_schedules(mission_id: int, user_id: int = Depends(verify_token), db: Session = Depends(get_db)):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    rows = (
        db.query(MissionRecapSchedule)
        .filter(MissionRecapSchedule.mission_id == mission.id)
        .order_by(MissionRecapSchedule.created_at.asc())
        .all()
    )
    return {"schedules": [_schedule_detail(s) for s in rows]}


@router.post("/api/automations/missions/{mission_id}/recap-schedules")
async def create_recap_schedule(
    mission_id: int,
    body: RecapScheduleCreate,
    user_id: int = Depends(verify_token),
    db: Session = Depends(get_db),
):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    if mission.status != "active":
        raise HTTPException(status_code=400, detail="Mission archivée : modification impossible")
    schedule = MissionRecapSchedule(
        mission_id=mission.id,
        company_id=mission.company_id,
        kind=body.kind,
        weekday=body.weekday if body.kind == "recurring" else None,
        run_date=body.run_date if body.kind == "once" else None,
        hour=body.hour,
        enabled=body.enabled,
    )
    db.add(schedule)
    db.commit()
    db.refresh(schedule)
    return {"id": schedule.id}


@router.put("/api/automations/missions/{mission_id}/recap-schedules/{schedule_id}")
async def update_recap_schedule(
    mission_id: int,
    schedule_id: int,
    body: RecapScheduleUpdate,
    user_id: int = Depends(verify_token),
    db: Session = Depends(get_db),
):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    if mission.status != "active":
        raise HTTPException(status_code=400, detail="Mission archivée : modification impossible")
    schedule = (
        db.query(MissionRecapSchedule)
        .filter(MissionRecapSchedule.id == schedule_id, MissionRecapSchedule.mission_id == mission.id)
        .first()
    )
    if not schedule:
        raise HTTPException(status_code=404, detail="Schedule not found")
    schedule.kind = body.kind
    schedule.weekday = body.weekday if body.kind == "recurring" else None
    schedule.run_date = body.run_date if body.kind == "once" else None
    schedule.hour = body.hour
    schedule.enabled = body.enabled
    db.commit()
    return {"success": True}


@router.delete("/api/automations/missions/{mission_id}/recap-schedules/{schedule_id}")
async def delete_recap_schedule(
    mission_id: int,
    schedule_id: int,
    user_id: int = Depends(verify_token),
    db: Session = Depends(get_db),
):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    if mission.status != "active":
        raise HTTPException(status_code=400, detail="Mission archivée : modification impossible")
    schedule = (
        db.query(MissionRecapSchedule)
        .filter(MissionRecapSchedule.id == schedule_id, MissionRecapSchedule.mission_id == mission.id)
        .first()
    )
    if not schedule:
        raise HTTPException(status_code=404, detail="Schedule not found")
    db.delete(schedule)
    db.commit()
    return {"success": True}


# --------------------------------------------------------------------------- #
# Mission documents (siloed RAG sources)
# --------------------------------------------------------------------------- #


@router.post("/api/automations/missions/{mission_id}/documents")
async def upload_mission_document(
    mission_id: int,
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    user_id: int = Depends(verify_token),
    db: Session = Depends(get_db),
):
    """Upload a source document siloed to the mission (async via Redis when available)."""
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)

    if not any(file.filename.lower().endswith(ext) for ext in ALLOWED_EXT):
        raise HTTPException(status_code=400, detail="Type de fichier non supporté")
    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="Fichier trop volumineux")

    from redis_client import get_redis
    from routers.documents import _process_document_background

    r = get_redis()
    if r is not None:
        task_id = str(uuid4())
        r.setex(
            f"doc_task:{task_id}",
            3600,
            json.dumps(
                {
                    "task_id": task_id,
                    "status": "processing",
                    "filename": file.filename,
                    "document_id": None,
                    "error": None,
                }
            ),
        )
        background_tasks.add_task(
            _process_document_background,
            task_id,
            file.filename,
            content,
            user_id,
            None,  # agent_id
            mission.company_id,
            mission.id,  # mission_id
        )
        return {"filename": file.filename, "task_id": task_id, "status": "processing"}

    from rag_engine import process_document_for_user

    doc_id = process_document_for_user(
        file.filename, content, user_id, db, agent_id=None, company_id=mission.company_id, mission_id=mission.id
    )
    return {"filename": file.filename, "document_id": doc_id, "status": "uploaded"}


@router.get("/api/automations/missions/{mission_id}/documents")
async def list_mission_documents(mission_id: int, user_id: int = Depends(verify_token), db: Session = Depends(get_db)):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    docs = (
        db.query(Document)
        .filter(Document.mission_id == mission.id, Document.is_mission_recap_source.is_(False))
        .order_by(Document.created_at.desc())
        .all()
    )
    return {
        "documents": [
            {"id": d.id, "filename": d.filename, "created_at": d.created_at.isoformat() if d.created_at else None}
            for d in docs
        ]
    }


@router.delete("/api/automations/missions/{mission_id}/documents/{document_id}")
async def delete_mission_document(
    mission_id: int, document_id: int, user_id: int = Depends(verify_token), db: Session = Depends(get_db)
):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    doc = db.query(Document).filter(Document.id == document_id, Document.mission_id == mission.id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    db.delete(doc)  # chunks cascade via the relationship
    db.commit()
    return {"success": True}


# --------------------------------------------------------------------------- #
# Mission recap-source documents (siloed, excluded from normal RAG)
# --------------------------------------------------------------------------- #


@router.post("/api/automations/missions/{mission_id}/recap-documents")
async def upload_mission_recap_document(
    mission_id: int,
    file: UploadFile = File(...),
    user_id: int = Depends(verify_token),
    db: Session = Depends(get_db),
):
    """Upload a document used ONLY as a source for this mission's recaps."""
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)

    if not any(file.filename.lower().endswith(ext) for ext in ALLOWED_EXT):
        raise HTTPException(status_code=400, detail="Type de fichier non supporté")
    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="Fichier trop volumineux")

    from rag_engine import process_document_for_user

    doc_id = process_document_for_user(
        file.filename, content, user_id, db, agent_id=None, company_id=mission.company_id, mission_id=mission.id
    )
    db.query(Document).filter(Document.id == doc_id).update({Document.is_mission_recap_source: True})
    db.commit()
    return {"filename": file.filename, "document_id": doc_id, "status": "uploaded"}


@router.get("/api/automations/missions/{mission_id}/recap-documents")
async def list_mission_recap_documents(
    mission_id: int, user_id: int = Depends(verify_token), db: Session = Depends(get_db)
):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    docs = (
        db.query(Document)
        .filter(Document.mission_id == mission.id, Document.is_mission_recap_source.is_(True))
        .order_by(Document.created_at.desc())
        .all()
    )
    return {
        "documents": [
            {"id": d.id, "filename": d.filename, "created_at": d.created_at.isoformat() if d.created_at else None}
            for d in docs
        ]
    }


@router.delete("/api/automations/missions/{mission_id}/recap-documents/{document_id}")
async def delete_mission_recap_document(
    mission_id: int, document_id: int, user_id: int = Depends(verify_token), db: Session = Depends(get_db)
):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    doc = (
        db.query(Document)
        .filter(
            Document.id == document_id,
            Document.mission_id == mission.id,
            Document.is_mission_recap_source.is_(True),
        )
        .first()
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    db.delete(doc)
    db.commit()
    return {"success": True}


# --------------------------------------------------------------------------- #
# Recaps
# --------------------------------------------------------------------------- #


@router.get("/api/automations/missions/{mission_id}/recaps")
async def list_recaps(
    mission_id: int,
    limit: int = Query(20, ge=1, le=100),
    offset: int = Query(0, ge=0),
    user_id: int = Depends(verify_token),
    db: Session = Depends(get_db),
):
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    base = db.query(MissionRecap).filter(MissionRecap.mission_id == mission.id)
    total = base.count()
    rows = base.order_by(MissionRecap.created_at.desc()).offset(offset).limit(limit).all()
    return {
        "recaps": [
            {
                "id": rcp.id,
                "period_start": rcp.period_start.isoformat(),
                "period_end": rcp.period_end.isoformat(),
                "content": rcp.content,
                "status": rcp.status,
                "trigger": rcp.trigger,
                "email_sent": rcp.email_sent,
                "created_at": rcp.created_at.isoformat() if rcp.created_at else None,
            }
            for rcp in rows
        ],
        "total": total,
    }


@router.post("/api/automations/missions/{mission_id}/recaps/generate")
async def generate_recap_now(mission_id: int, user_id: int = Depends(verify_token), db: Session = Depends(get_db)):
    """Generate a recap on demand (synchronous, no email, no scheduler impact)."""
    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    if not mission.agent_id:
        raise HTTPException(status_code=400, detail="Connectez un companion à la mission d'abord")

    from mission_recap import process_mission_recap

    result = process_mission_recap(mission, db, trigger="manual")
    if result["status"] == "no_data":
        raise HTTPException(status_code=400, detail="Aucun évènement à venir cette semaine")
    if result["status"] == "error":
        raise HTTPException(status_code=502, detail="La génération du récap a échoué")
    return {"recap_id": result["recap_id"], "content": result["content"]}


# --------------------------------------------------------------------------- #
# Mission chat
# --------------------------------------------------------------------------- #


@router.post("/api/automations/missions/{mission_id}/chat")
async def mission_chat(
    mission_id: int, body: MissionChatRequest, user_id: int = Depends(verify_token), db: Session = Depends(get_db)
):
    """Chat with the connected companion in the mission's context."""
    from datetime import date, timedelta

    from database import Conversation, Message
    from mistral_embeddings import get_embedding_fast
    from openai_client import get_chat_response
    from rag_engine import search_similar_texts_for_user
    from weekly_recap import get_model_id_for_agent

    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    if not mission.agent_id:
        raise HTTPException(status_code=400, detail="Connectez un companion à la mission d'abord")
    agent = db.query(Agent).filter(Agent.id == mission.agent_id).first()

    # Resolve or create the conversation (scoped to this mission).
    conversation = None
    if body.conversation_id:
        conversation = (
            db.query(Conversation)
            .filter(Conversation.id == body.conversation_id, Conversation.mission_id == mission.id)
            .first()
        )
    if conversation is None:
        conversation = Conversation(
            agent_id=mission.agent_id,
            user_id=user_id,
            company_id=mission.company_id,
            mission_id=mission.id,
            title=body.message[:60],
        )
        db.add(conversation)
        db.commit()
        db.refresh(conversation)

    # Build mission context: objective + events [-7, +7] + RAG snippets.
    today = date.today()
    events = (
        db.query(MissionEvent)
        .filter(
            MissionEvent.mission_id == mission.id,
            MissionEvent.event_date >= today - timedelta(days=7),
            MissionEvent.event_date <= today + timedelta(days=7),
        )
        .order_by(MissionEvent.event_date.asc())
        .all()
    )
    events_text = (
        "\n".join(
            f"- {e.event_date.isoformat()} : {e.title}" + (f" — {e.description}" if e.description else "")
            for e in events
        )
        or "(aucun évènement proche)"
    )

    rag_text = ""
    try:
        emb = get_embedding_fast(body.message)
        results = search_similar_texts_for_user(
            emb, user_id=user_id, db=db, top_k=4, company_id=mission.company_id, mission_id=mission.id
        )
        rag_text = "\n".join(f"> {r['text'][:800]}" for r in results)
    except Exception as e:
        logger.warning(f"Mission {mission.id} chat RAG failed: {e}")
    rag_text = rag_text or "(aucun extrait documentaire pertinent)"

    agent_name = (agent.name if agent else None) or "Assistant"
    agent_context = (getattr(agent, "contexte", "") if agent else "") or ""
    system_prompt = (
        f"Tu es {agent_name}, un assistant IA d'entreprise. {agent_context}\n\n"
        f'Tu assistes sur une mission dont l\'objectif est :\n"""{mission.objective.strip()}"""\n\n'
        f"Évènements proches (±7 jours) :\n{events_text}\n\n"
        f"Extraits de documents de la mission :\n{rag_text}\n\n"
        "Réponds en t'appuyant sur ce contexte. N'invente pas d'évènements absents."
    )

    db.add(Message(conversation_id=conversation.id, company_id=mission.company_id, role="user", content=body.message))
    db.commit()

    model_id = get_model_id_for_agent(agent) if agent else "mistral:mistral-large-latest"
    answer = get_chat_response(
        [{"role": "system", "content": system_prompt}, {"role": "user", "content": body.message}],
        model_id=model_id,
    )

    db.add(Message(conversation_id=conversation.id, company_id=mission.company_id, role="agent", content=answer))
    db.commit()

    return {"conversation_id": conversation.id, "answer": answer}


@router.get("/api/automations/missions/{mission_id}/conversations")
async def list_mission_conversations(
    mission_id: int, user_id: int = Depends(verify_token), db: Session = Depends(get_db)
):
    from database import Conversation, Message

    user_id = int(user_id)
    membership = require_role(user_id, db, "member")
    mission = _get_mission_or_404(mission_id, user_id, membership.company_id, db)
    convs = (
        db.query(Conversation)
        .filter(Conversation.mission_id == mission.id, Conversation.user_id == user_id)
        .order_by(Conversation.created_at.desc())
        .all()
    )
    result = []
    for c in convs:
        messages = db.query(Message).filter(Message.conversation_id == c.id).order_by(Message.timestamp.asc()).all()
        result.append(
            {
                "id": c.id,
                "title": c.title,
                "messages": [{"role": m.role, "content": m.content} for m in messages],
            }
        )
    return {"conversations": result}
