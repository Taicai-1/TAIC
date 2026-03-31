# Standard library
import os
import io
import re
import time
import json
import logging
import mimetypes
import traceback
import threading
import hmac
import hashlib
from uuid import uuid4
from datetime import datetime, timedelta
from collections import deque

# Third-party
import requests
import httpx
import openai
import urllib3
import redis
import bcrypt
from fastapi import FastAPI, UploadFile, File, Depends, HTTPException, Request, Body, Form, Response
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from fastapi.staticfiles import StaticFiles
from sqlalchemy.orm import Session, joinedload
from sqlalchemy import text
from pydantic import BaseModel
from typing import List, Optional
from google.cloud import storage
import google.auth
from google.auth.transport.requests import AuthorizedSession

# Local modules
from auth import create_access_token, verify_token, hash_password, verify_password, hash_reset_token, verify_pre_2fa_token, verify_setup_token
from email_service import send_password_reset_email, send_invitation_email, send_verification_email, send_feedback_email
from database import get_db, init_db, ensure_columns, migrate_existing_company_memberships, User, Document, Agent, Team, Base, engine, Conversation, Message, PasswordResetToken, Company, CompanyMembership, CompanyInvitation, WeeklyRecapLog, NotionLink, AgentShare
from rag_engine import get_answer, get_answer_with_files, process_document_for_user
from openai_client import get_embedding
from file_generator import FileGenerator
from utils import logger, event_tracker
from utils_ai import normalize_model_output, extract_json_object_from_text
from validation import (
    UserCreateValidated,
    AgentCreateValidated,
    MessageCreateValidated,
    QuestionRequestValidated,
    ConversationTitleValidated,
    UrlUploadValidated,
    TeamCreateValidated,
    ChangePasswordRequest,
    validate_id_parameter,
    validate_file_extension,
    validate_file_content,
    validate_file_size,
    sanitize_filename,
    sanitize_text,
    MAX_FILE_SIZE,
    ALLOWED_FILE_EXTENSIONS
)

# Configuration du logger

# Setup Google Cloud Logging si en production, sinon basicConfig
if os.getenv("GOOGLE_CLOUD_PROJECT"):
    try:
        from google.cloud import logging as cloud_logging
        client = cloud_logging.Client()
        client.setup_logging()
        logger = logging.getLogger("app")
    except ImportError:
        # Fallback to basic logging if Google Cloud Logging not available
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s %(levelname)s %(name)s %(message)s',
        )
        logger = logging.getLogger("app")
else:
    # Development: use basic logging
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s %(levelname)s %(name)s %(message)s',
    )
    logger = logging.getLogger("app")

app = FastAPI(title="TAIC Companion API", version="1.0.0")

# --- Agent type helpers ---

_AGENT_TYPE_MODEL_MAP = {
    "recherche_live": ("PERPLEXITY_MODEL", "perplexity:sonar"),
    "visuel": (None, "imagen:imagen-3.0-generate-002"),
}
_DEFAULT_MODEL = ("MISTRAL_MODEL", "mistral:mistral-small-latest")

_AGENT_TYPE_PROVIDER_MAP = {
    "recherche_live": "perplexity",
    "visuel": "imagen",
}


def resolve_model_id(agent) -> str:
    """Return the model_id for an agent based on its type."""
    atype = getattr(agent, "type", "conversationnel")
    env_var, default = _AGENT_TYPE_MODEL_MAP.get(atype, _DEFAULT_MODEL)
    return os.getenv(env_var, default) if env_var else default


def resolve_llm_provider(agent_type: str) -> str:
    """Return the llm_provider string for an agent type."""
    return _AGENT_TYPE_PROVIDER_MAP.get(agent_type, "mistral")


# Security headers middleware
@app.middleware("http")
async def add_security_headers(request: Request, call_next):
    """Add security headers to all responses"""
    response = await call_next(request)

    # HSTS: Force HTTPS for 1 year (only in production)
    if not request.url.hostname in ["localhost", "127.0.0.1"]:
        response.headers["Strict-Transport-Security"] = "max-age=31536000; includeSubDomains"

    # CSP: Prevent XSS attacks
    # Strict policy: no unsafe-eval, no unsafe-inline for scripts
    response.headers["Content-Security-Policy"] = (
        "default-src 'none'; "
        "script-src 'self' https://cdn.jsdelivr.net; "
        "style-src 'self' https://fonts.googleapis.com; "
        "font-src 'self' https://fonts.gstatic.com; "
        "img-src 'self' data: https:; "
        "connect-src 'self' https://api.openai.com https://api.mistral.ai https://generativelanguage.googleapis.com; "
        "frame-ancestors 'none';"
    )

    # X-Frame-Options: Prevent clickjacking
    response.headers["X-Frame-Options"] = "DENY"

    # X-Content-Type-Options: Prevent MIME sniffing
    response.headers["X-Content-Type-Options"] = "nosniff"

    # X-XSS-Protection: Enable browser XSS protection (legacy but still useful)
    response.headers["X-XSS-Protection"] = "1; mode=block"

    # Referrer-Policy: Control referrer information
    response.headers["Referrer-Policy"] = "strict-origin-when-cross-origin"

    # Permissions-Policy: Disable unnecessary browser features
    response.headers["Permissions-Policy"] = "geolocation=(), microphone=(), camera=()"

    return response


# ============================================================================
# REDIS CONNECTION FOR DISTRIBUTED RATE LIMITING
# ============================================================================

def get_redis_client():
    """Get Redis client for distributed rate limiting (Phase 1 security upgrade)

    Security: Distributed rate limiting prevents bypass via instance restart
    or multi-instance Cloud Run scaling.
    """
    redis_host = os.getenv("REDIS_HOST", "redis")  # "redis" for docker-compose, Cloud Memorystore for production
    redis_port = int(os.getenv("REDIS_PORT", "6379"))
    redis_password = os.getenv("REDIS_PASSWORD", None)  # Optional for Cloud Memorystore

    try:
        client = redis.Redis(
            host=redis_host,
            port=redis_port,
            password=redis_password,
            decode_responses=True,
            socket_connect_timeout=5,
            socket_timeout=5
        )
        # Test connection
        client.ping()
        return client
    except Exception as e:
        logger.error(f"Redis connection failed: {e}. Falling back to in-memory rate limiting.")
        return None

# Initialize Redis client
redis_client = get_redis_client()

# Fallback in-memory storage (only used if Redis unavailable)
_auth_rate_limit_fallback = {}
_public_chat_rate_fallback = {}
_api_rate_limit_fallback = {}

# Rate limiting configuration
_AUTH_LIMIT = 5  # max failed attempts per window
_AUTH_WINDOW = 3600  # 1 hour in seconds

# Per-user rate limiting for API endpoints (upload, ask, extractText)
_API_UPLOAD_LIMIT = 30  # max uploads per window per user
_API_ASK_LIMIT = 60  # max /ask calls per window per user
_API_EXTRACT_LIMIT = 30  # max extractText calls per window per user
_API_WINDOW = 3600  # 1 hour in seconds


def _check_api_rate_limit(user_id: str, action: str, limit: int) -> bool:
    """
    Check if user has exceeded rate limit for a specific API action using Redis.
    Returns True if allowed, False if rate limited.
    """
    key = f"rate_limit:{action}:{user_id}"

    if redis_client:
        try:
            current = redis_client.get(key)
            current = int(current) if current else 0
            if current >= limit:
                return False
            pipe = redis_client.pipeline()
            pipe.incr(key)
            pipe.expire(key, _API_WINDOW)
            pipe.execute()
            return True
        except Exception as e:
            logger.error(f"Redis rate limit check failed for {action}: {e}. Using fallback.")

    # Fallback to in-memory
    now = time.time()
    fallback_key = f"{action}:{user_id}"
    attempts = _api_rate_limit_fallback.get(fallback_key, [])
    attempts = [t for t in attempts if now - t < _API_WINDOW]
    if len(attempts) >= limit:
        return False
    attempts.append(now)
    _api_rate_limit_fallback[fallback_key] = attempts
    return True


def _check_auth_rate_limit(ip: str) -> bool:
    """
    Check if IP has exceeded rate limit for auth endpoints.
    Returns True if allowed, False if rate limited.

    Only failed attempts are counted (via _record_auth_failure).
    Max 5 failures per hour per IP.
    """
    key = f"rate_limit:auth:{ip}"

    # Try Redis first
    if redis_client:
        try:
            current = redis_client.get(key)
            current = int(current) if current else 0
            return current < _AUTH_LIMIT
        except Exception as e:
            logger.error(f"Redis rate limit check failed: {e}. Using fallback.")

    # Fallback to in-memory
    now = time.time()
    attempts = _auth_rate_limit_fallback.get(ip, [])
    attempts = [t for t in attempts if now - t < _AUTH_WINDOW]
    _auth_rate_limit_fallback[ip] = attempts
    return len(attempts) < _AUTH_LIMIT


def _record_auth_failure(ip: str):
    """Record a failed auth attempt for rate limiting."""
    key = f"rate_limit:auth:{ip}"

    if redis_client:
        try:
            pipe = redis_client.pipeline()
            pipe.incr(key)
            pipe.expire(key, _AUTH_WINDOW)
            pipe.execute()
            return
        except Exception as e:
            logger.error(f"Redis rate limit record failed: {e}. Using fallback.")

    # Fallback to in-memory
    now = time.time()
    attempts = _auth_rate_limit_fallback.get(ip, [])
    attempts = [t for t in attempts if now - t < _AUTH_WINDOW]
    attempts.append(now)
    _auth_rate_limit_fallback[ip] = attempts


# Ajout d'un endpoint pour ajouter une URL comme source
# Use validated model from validation.py
# class UrlUploadRequest is now UrlUploadValidated

# Expose le dossier profile_photos en statique après la création de l'app
from fastapi.staticfiles import StaticFiles
import os
if not os.path.exists("profile_photos"):
    os.makedirs("profile_photos")
app.mount("/profile_photos", StaticFiles(directory="profile_photos"), name="profile_photos")

@app.post("/upload-url")
async def upload_url(
    request: UrlUploadValidated,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Ajoute une URL comme document/source pour le RAG"""
    try:
        # Headers to mimic a real browser and avoid being blocked
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8',
            'Accept-Language': 'fr-FR,fr;q=0.9,en-US;q=0.8,en;q=0.7',
            'Accept-Encoding': 'gzip, deflate, br',
            'Connection': 'keep-alive',
            'Upgrade-Insecure-Requests': '1',
            'Sec-Fetch-Dest': 'document',
            'Sec-Fetch-Mode': 'navigate',
            'Sec-Fetch-Site': 'none',
            'Cache-Control': 'max-age=0'
        }

        # Télécharger le contenu de l'URL avec retry logic
        max_retries = 3
        retry_delay = 2
        html = None
        last_error = None

        def _is_safe_redirect(url: str) -> bool:
            """Check that a redirect URL doesn't point to internal/private networks."""
            from urllib.parse import urlparse
            blocked_patterns = [
                'localhost', '127.0.0.1', '0.0.0.0',
                '192.168.', '10.', '172.16.', '172.17.', '172.18.', '172.19.',
                '172.20.', '172.21.', '172.22.', '172.23.', '172.24.', '172.25.',
                '172.26.', '172.27.', '172.28.', '172.29.', '172.30.', '172.31.',
                '169.254.', '[::1]', '[fc', '[fd', 'metadata.google.internal',
            ]
            return not any(pattern in url.lower() for pattern in blocked_patterns)

        for attempt in range(max_retries):
            try:
                response = requests.get(
                    request.url,
                    headers=headers,
                    timeout=20,
                    allow_redirects=False,
                    verify=True
                )
                # Manually follow redirects with SSRF validation (max 5 hops)
                redirect_count = 0
                while response.is_redirect and redirect_count < 5:
                    redirect_url = response.headers.get("Location", "")
                    if not redirect_url or not _is_safe_redirect(redirect_url):
                        raise requests.exceptions.ConnectionError("Redirect to blocked destination")
                    response = requests.get(
                        redirect_url,
                        headers=headers,
                        timeout=20,
                        allow_redirects=False,
                        verify=True
                    )
                    redirect_count += 1
                response.raise_for_status()
                # Try to get encoding from response headers or detect it
                if response.encoding:
                    html = response.text
                else:
                    response.encoding = response.apparent_encoding
                    html = response.text
                break
            except requests.exceptions.SSLError as e:
                logger.warning(f"SSL error on attempt {attempt + 1} for {request.url}: {e}")
                last_error = e
                if attempt < max_retries - 1:
                    time.sleep(retry_delay)
                    # Try without SSL verification on last attempt
                    # Security: Removed insecure verify=False fallback to prevent MITM attacks
                    # If SSL verification fails, the request should fail rather than be vulnerable
            except (requests.exceptions.Timeout, requests.exceptions.ConnectionError) as e:
                logger.warning(f"Connection error on attempt {attempt + 1} for {request.url}: {e}")
                last_error = e
                if attempt < max_retries - 1:
                    time.sleep(retry_delay)
            except requests.exceptions.HTTPError as e:
                logger.error(f"HTTP error for {request.url}: {e}")
                last_error = e
                break
            except Exception as e:
                logger.error(f"Unexpected error on attempt {attempt + 1} for {request.url}: {e}")
                last_error = e
                if attempt < max_retries - 1:
                    time.sleep(retry_delay)

        if html is None:
            error_msg = f"Failed to fetch URL after {max_retries} attempts"
            if last_error:
                error_msg += f": {str(last_error)}"
            logger.error(error_msg)
            raise HTTPException(status_code=400, detail="Unable to fetch the provided URL. Please check the URL and try again.")

        # Extraire uniquement les informations utiles : titre, meta description et contenu principal
        from bs4 import BeautifulSoup
        try:
            from readability import Document as ReadabilityDocument
            use_readability = True
        except Exception:
            use_readability = False

        title = ""
        meta_desc = ""
        main_text = ""

        try:
            soup = BeautifulSoup(html, "lxml")
            # Title
            if soup.title and soup.title.string:
                title = soup.title.string.strip()
            # Meta description
            md = soup.find("meta", attrs={"name": "description"})
            if md and md.get("content"):
                meta_desc = md.get("content").strip()

            # Try Readability first (better extraction of main article)
            if use_readability:
                try:
                    doc = ReadabilityDocument(html)
                    main_html = doc.summary()
                    main_soup = BeautifulSoup(main_html, "lxml")
                    # Get visible text
                    main_text = "\n".join([p.get_text(separator=" ", strip=True) for p in main_soup.find_all(["p", "h1", "h2", "h3"])])
                except Exception:
                    use_readability = False

            # Fallback: extract visible text from body, but filter out navigation/footer links
            if not main_text:
                body = soup.body
                if body:
                    # Remove scripts, styles, nav, footer, aside
                    for tag in body.find_all(["script", "style", "nav", "footer", "aside", "header", "form", "noscript"]):
                        tag.decompose()
                    # Collect paragraphs and headings
                    paragraphs = [p.get_text(separator=" ", strip=True) for p in body.find_all(["p", "h1", "h2", "h3"]) if p.get_text(strip=True)]
                    main_text = "\n".join(paragraphs)

            # Build a cleaned text that contains only useful metadata + main content (limit length)
            cleaned = []
            if title:
                cleaned.append(f"Title: {title}")
            if meta_desc:
                cleaned.append(f"Description: {meta_desc}")
            if main_text:
                cleaned.append("Content:\n" + main_text)

            content = "\n\n".join(cleaned)
            if not content.strip():
                # If nothing meaningful found, fallback to raw text (but cleaned)
                content = soup.get_text(separator="\n", strip=True)

        except Exception as e:
            logger.warning(f"Failed to parse HTML for useful content, falling back to raw. Error: {e}")
            content = html

        # Shorten the filename
        filename = request.url.split("//")[-1][:100].replace("/", "_") + ".txt"

        # Truncate content to a reasonable length to avoid huge token usage (e.g., 200k chars)
        max_chars = 200000
        if len(content) > max_chars:
            content = content[:max_chars]

        # Indexer le document comme pour un upload classique (send cleaned text)
        doc_id = process_document_for_user(filename, content.encode("utf-8", errors="ignore"), int(user_id), db, agent_id=request.agent_id)

        logger.info(f"URL ajoutée pour user {user_id}, agent {request.agent_id}: {request.url}")
        event_tracker.track_document_upload(int(user_id), request.url, len(content))

        return {"url": request.url, "document_id": doc_id, "agent_id": request.agent_id, "status": "uploaded"}
    except Exception as e:
        logger.error(f"Erreur lors de l'ajout d'URL: {e}")
        raise HTTPException(status_code=500, detail="Erreur lors de l'ajout de l'URL")
    


# CORS configuration - SÉCURISÉ
# Security: Restrict origins to only trusted domains to prevent unauthorized API access
# Separate CORS configuration for development vs production environments
ENVIRONMENT = os.getenv("ENVIRONMENT", "development")

# Base production origins
allowed_origins = [
    "https://taic.ai",
    "https://www.taic.ai",
    "https://dev-taic-frontend-817946451913.europe-west1.run.app",
    "https://applydi-frontend-817946451913.europe-west1.run.app",
]

# Add localhost only in development (CRITICAL: prevents CSRF in production)
if ENVIRONMENT == "development":
    allowed_origins.extend([
        "http://localhost:3000",
        "http://localhost:8080"
    ])
    logger.info("CORS: Development mode - localhost origins enabled")
else:
    logger.info("CORS: Production mode - localhost origins disabled")

app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_origin_regex=r"^https://([\w-]+\.taic\.ai|[\w-]+-817946451913\.europe-west1\.run\.app)$" if ENVIRONMENT == "production" else None,
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS", "PATCH"],
    allow_headers=["Content-Type", "Authorization", "Accept", "Origin", "X-Requested-With"],
    expose_headers=["Content-Length", "Content-Type"],
    max_age=600,  # Cache preflight requests for 10 minutes
)

# Initialize database on startup
@app.on_event("startup")
async def startup_event():
    """Initialize database and create tables on startup

    Note: Database migrations should be run separately via scripts/migrations/
    before deploying new versions.
    """
    try:
        logger.info("Initializing database...")
        init_db()
        logger.info("Creating database tables...")
        Base.metadata.create_all(bind=engine)
        ensure_columns()
        migrate_existing_company_memberships()
        logger.info("Database initialization completed successfully")
    except Exception as e:
        logger.error(f"Database initialization failed: {e}")
        # Don't raise exception to allow the app to start, but log the error

# Health check endpoints
@app.get("/")
async def root():
    """Root endpoint"""
    return {"message": "TAIC Companion API is running", "status": "ok"}

@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {"status": "healthy", "service": "TAIC Companion API"}

# Pydantic models
# Use validated model from validation.py
# class UserCreate is now UserCreateValidated

class UserLogin(BaseModel):
    username: str
    password: str

# Use validated model from validation.py
# class QuestionRequest is now QuestionRequestValidated

class AgentCreate(BaseModel):
    name: str
    contexte: str = None
    biographie: str = None
    profile_photo: str = None  # URL or filename
    email: str
    password: str

class AgentResponse(BaseModel):
    id: int
    name: str
    contexte: str = None
    biographie: str = None
    profile_photo: str = None
    email: str
    user_id: int
    created_at: datetime
    class Config:
        from_attributes = True

# Routes
@app.post("/register")
async def register(user: UserCreateValidated, request: Request, db: Session = Depends(get_db)):
    """Register new user"""
    # Rate limiting: prevent brute force account creation
    ip = request.client.host if hasattr(request, 'client') and request.client else 'unknown'
    if not _check_auth_rate_limit(ip):
        logger.warning(f"Rate limit exceeded for registration from IP: {ip}")
        raise HTTPException(status_code=429, detail="Too many attempts. Please try again in 1 hour.")

    try:
        # Check if user exists
        if db.query(User).filter(User.username == user.username).first():
            raise HTTPException(status_code=400, detail="Username already registered")

        if db.query(User).filter(User.email == user.email).first():
            raise HTTPException(status_code=400, detail="Email already registered")

        # Create new user
        hashed_password = hash_password(user.password)
        db_user = User(
            username=user.username,
            email=user.email,
            hashed_password=hashed_password
        )
        db.add(db_user)
        db.commit()
        db.refresh(db_user)

        # Handle invite_code: join org on registration
        if user.invite_code:
            company = db.query(Company).filter(
                Company.invite_code == user.invite_code,
                Company.invite_code_enabled == True
            ).first()
            if company:
                db_user.company_id = company.id
                membership = CompanyMembership(
                    user_id=db_user.id,
                    company_id=company.id,
                    role="member"
                )
                db.add(membership)
                db.commit()
                logger.info(f"User {user.username} joined company {company.name} via invite_code at registration")

        logger.info(f"User registered: {user.username}")
        event_tracker.track_user_action(db_user.id, "user_registered")

        # Send verification email
        try:
            verify_token = create_access_token(
                data={"sub": str(db_user.id), "type": "email_verify"},
                expires_delta=timedelta(hours=24)
            )
            frontend_url = os.getenv("FRONTEND_URL", "https://taic.ai")
            verify_link = f"{frontend_url}/verify-email?token={verify_token}"
            send_verification_email(db_user.email, verify_link)
            logger.info(f"Verification email sent to {db_user.email}")
        except Exception as e:
            logger.error(f"Failed to send verification email to {db_user.email}: {e}")

        return {"message": "User created successfully. Please verify your email."}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Registration error: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")

@app.post("/login")
async def login(user: UserLogin, request: Request, response: Response, db: Session = Depends(get_db)):
    """Login user with HttpOnly secure cookie"""
    # Rate limiting: prevent brute force password attacks
    ip = request.client.host if hasattr(request, 'client') and request.client else 'unknown'
    if not _check_auth_rate_limit(ip):
        logger.warning(f"Rate limit exceeded for login from IP: {ip}")
        raise HTTPException(status_code=429, detail="Too many failed attempts. Please try again in 1 hour.")

    try:
        # Permet la connexion avec username OU email
        db_user = db.query(User).filter(
            (User.username == user.username) | (User.email == user.username)
        ).first()
        if not db_user:
            _record_auth_failure(ip)
            raise HTTPException(status_code=401, detail="Invalid credentials")

        # OAuth-only user (no password) trying to login with password
        if not db_user.hashed_password:
            raise HTTPException(status_code=400, detail="This account uses Google sign-in. Please use the Google button.")

        if not verify_password(user.password, db_user.hashed_password):
            _record_auth_failure(ip)
            raise HTTPException(status_code=401, detail="Invalid credentials")

        # Email verification check (before 2FA)
        if not db_user.email_verified:
            # Auto-send verification email
            try:
                verify_token = create_access_token(
                    data={"sub": str(db_user.id), "type": "email_verify"},
                    expires_delta=timedelta(hours=24)
                )
                frontend_url = os.getenv("FRONTEND_URL", "https://taic.ai")
                verify_link = f"{frontend_url}/verify-email?token={verify_token}"
                send_verification_email(db_user.email, verify_link)
                logger.info(f"Verification email auto-sent to {db_user.email} on login")
            except Exception as e:
                logger.error(f"Failed to auto-send verification email on login: {e}")

            # Mask email: j***@gmail.com
            email = db_user.email
            at_idx = email.index("@")
            masked = email[0] + "***" + email[at_idx:]
            return {
                "requires_email_verification": True,
                "email": masked
            }

        # 2FA flow: check user's TOTP state
        if db_user.totp_enabled:
            # User has 2FA enabled → issue restricted pre_2fa token (5 min)
            pre_2fa_token = create_access_token(
                data={"sub": str(db_user.id), "type": "pre_2fa"},
                expires_delta=timedelta(minutes=5)
            )
            logger.info(f"User {user.username} requires 2FA verification")
            return {
                "requires_2fa": True,
                "pre_2fa_token": pre_2fa_token,
                "token_type": "bearer"
            }

        if not getattr(db_user, 'totp_setup_completed_at', None):
            # User has NOT set up 2FA yet → issue restricted setup token (30 min)
            setup_token = create_access_token(
                data={"sub": str(db_user.id), "type": "needs_2fa_setup"},
                expires_delta=timedelta(minutes=30)
            )
            logger.info(f"User {user.username} needs 2FA setup")
            return {
                "requires_2fa_setup": True,
                "setup_token": setup_token,
                "token_type": "bearer"
            }

        # 2FA completed → issue full access token
        access_token = create_access_token(data={"sub": str(db_user.id)})

        # Security: Set HttpOnly secure cookie to prevent XSS token theft
        response.set_cookie(
            key="token",
            value=access_token,
            httponly=True,
            secure=True,
            samesite="none",
            max_age=28800,
            path="/"
        )

        logger.info(f"User logged in: {user.username}")
        event_tracker.track_user_action(db_user.id, "user_login")

        return {"access_token": access_token, "token_type": "bearer"}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Login error: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")


############################
# Feedback Endpoint
############################

class FeedbackRequest(BaseModel):
    type: str  # "bug", "feature", "feedback", "other"
    message: str

@app.post("/feedback")
async def submit_feedback(req: FeedbackRequest, request: Request, db: Session = Depends(get_db)):
    """Submit user feedback via email to contact@taic.co."""
    from auth import verify_token_from_cookie

    user_id = verify_token_from_cookie(request)
    db_user = db.query(User).filter(User.id == int(user_id)).first()
    if not db_user:
        raise HTTPException(status_code=401, detail="User not found")

    if not req.message.strip():
        raise HTTPException(status_code=400, detail="Message cannot be empty")

    if len(req.message) > 5000:
        raise HTTPException(status_code=400, detail="Message too long (max 5000 characters)")

    try:
        send_feedback_email(db_user.email, db_user.username, req.type, req.message.strip())
        return {"message": "Feedback sent successfully"}
    except Exception as e:
        logger.error(f"Failed to send feedback email: {e}")
        raise HTTPException(status_code=500, detail="Failed to send feedback")


############################
# Email Verification Endpoints
############################

class VerifyEmailRequest(BaseModel):
    token: str

class ResendVerificationRequest(BaseModel):
    email: str

@app.post("/auth/verify-email")
async def verify_email(req: VerifyEmailRequest, db: Session = Depends(get_db)):
    """Verify user email with JWT token from verification link."""
    from auth import SECRET_KEY, ALGORITHM
    import jwt as pyjwt

    try:
        payload = pyjwt.decode(req.token, SECRET_KEY, algorithms=[ALGORITHM])
        if payload.get("type") != "email_verify":
            raise HTTPException(status_code=400, detail="Invalid verification token")

        user_id = payload.get("sub")
        if not user_id:
            raise HTTPException(status_code=400, detail="Invalid verification token")

        db_user = db.query(User).filter(User.id == int(user_id)).first()
        if not db_user:
            raise HTTPException(status_code=404, detail="User not found")

        if db_user.email_verified:
            return {"message": "Email already verified"}

        db_user.email_verified = True
        db.commit()
        logger.info(f"Email verified for user {db_user.username}")
        return {"message": "Email verified successfully"}

    except pyjwt.ExpiredSignatureError:
        raise HTTPException(status_code=400, detail="Verification link has expired")
    except pyjwt.PyJWTError:
        raise HTTPException(status_code=400, detail="Invalid verification token")

@app.post("/auth/resend-verification")
async def resend_verification(req: ResendVerificationRequest, request: Request, db: Session = Depends(get_db)):
    """Resend email verification link. Rate limited."""
    ip = request.client.host if hasattr(request, 'client') and request.client else 'unknown'
    if not _check_auth_rate_limit(ip):
        raise HTTPException(status_code=429, detail="Too many attempts. Please try again later.")

    # Always return success to prevent email enumeration
    db_user = db.query(User).filter(User.email == req.email).first()
    if not db_user or db_user.email_verified:
        return {"message": "If this email exists and is unverified, a verification link has been sent."}

    try:
        verify_token = create_access_token(
            data={"sub": str(db_user.id), "type": "email_verify"},
            expires_delta=timedelta(hours=24)
        )
        frontend_url = os.getenv("FRONTEND_URL", "https://taic.ai")
        verify_link = f"{frontend_url}/verify-email?token={verify_token}"
        send_verification_email(db_user.email, verify_link)
        logger.info(f"Verification email resent to {db_user.email}")
    except Exception as e:
        logger.error(f"Failed to resend verification email: {e}")

    return {"message": "If this email exists and is unverified, a verification link has been sent."}


############################
# Google OAuth2 Endpoint
############################

class GoogleOAuthRequest(BaseModel):
    credential: str
    invite_code: Optional[str] = None

@app.post("/auth/google")
async def google_oauth(req: GoogleOAuthRequest, request: Request, response: Response, db: Session = Depends(get_db)):
    """Authenticate with Google OAuth2 ID token."""
    from google.oauth2 import id_token as google_id_token
    from google.auth.transport import requests as google_requests

    google_client_id = os.getenv("GOOGLE_OAUTH_CLIENT_ID")
    if not google_client_id:
        raise HTTPException(status_code=500, detail="Google OAuth not configured")

    try:
        # Verify the Google ID token
        idinfo = google_id_token.verify_oauth2_token(
            req.credential,
            google_requests.Request(),
            google_client_id
        )

        email = idinfo.get("email")
        name = idinfo.get("name", "")
        if not email:
            raise HTTPException(status_code=400, detail="Google account has no email")

        # Check for existing user by email
        db_user = db.query(User).filter(User.email == email).first()

        if db_user:
            # Existing user — link Google if not already done
            if not db_user.email_verified:
                db_user.email_verified = True
            if not db_user.oauth_provider:
                db_user.oauth_provider = "google"
            db.commit()
        else:
            # New user — create account
            username = email.split("@")[0]
            # Ensure username uniqueness
            base_username = username
            counter = 1
            while db.query(User).filter(User.username == username).first():
                username = f"{base_username}{counter}"
                counter += 1

            db_user = User(
                username=username,
                email=email,
                hashed_password=None,
                email_verified=True,
                oauth_provider="google"
            )
            db.add(db_user)
            db.commit()
            db.refresh(db_user)

            # Handle invite_code
            if req.invite_code:
                company = db.query(Company).filter(
                    Company.invite_code == req.invite_code,
                    Company.invite_code_enabled == True
                ).first()
                if company:
                    db_user.company_id = company.id
                    membership = CompanyMembership(
                        user_id=db_user.id,
                        company_id=company.id,
                        role="member"
                    )
                    db.add(membership)
                    db.commit()

            logger.info(f"Google OAuth user created: {username}")
            event_tracker.track_user_action(db_user.id, "user_registered_google")

        # Issue full access token (skip 2FA for OAuth)
        access_token = create_access_token(data={"sub": str(db_user.id)})

        response.set_cookie(
            key="token",
            value=access_token,
            httponly=True,
            secure=True,
            samesite="none",
            max_age=28800,
            path="/"
        )

        logger.info(f"Google OAuth login: {email}")
        event_tracker.track_user_action(db_user.id, "user_login_google")

        return {"access_token": access_token, "token_type": "bearer"}

    except ValueError as e:
        logger.error(f"Google OAuth token verification failed: {e}")
        raise HTTPException(status_code=401, detail="Invalid Google token")
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Google OAuth error: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")


@app.get("/auth/verify")
async def verify_auth(request: Request, db: Session = Depends(get_db)):
    """Verify authentication via HttpOnly cookie

    This endpoint checks if the user is authenticated by verifying the HttpOnly cookie.
    Used by frontend pages to check auth without localStorage.

    Returns:
        - 200: User authenticated, returns user info
        - 401: Not authenticated
    """
    from auth import verify_token_from_cookie

    try:
        user_id = verify_token_from_cookie(request)

        # Get user info
        db_user = db.query(User).filter(User.id == int(user_id)).first()
        if not db_user:
            raise HTTPException(status_code=401, detail="User not found")

        return {
            "authenticated": True,
            "user": {
                "id": db_user.id,
                "username": db_user.username,
                "email": db_user.email,
                "totp_enabled": db_user.totp_enabled
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Auth verification error: {e}")
        raise HTTPException(status_code=401, detail="Not authenticated")


############################
# 2FA (TOTP) Endpoints
############################

# 2FA rate limiting configuration
_2FA_LIMIT = 5  # max attempts per window
_2FA_WINDOW = 300  # 5 minutes in seconds
_2fa_rate_limit_fallback = {}


def _check_2fa_rate_limit(user_id: str) -> bool:
    """Check if user has exceeded 2FA verification rate limit.
    Returns True if allowed, False if rate limited.
    """
    key = f"rate_limit:2fa:{user_id}"

    if redis_client:
        try:
            current = redis_client.get(key)
            current = int(current) if current else 0
            if current >= _2FA_LIMIT:
                return False
            pipe = redis_client.pipeline()
            pipe.incr(key)
            pipe.expire(key, _2FA_WINDOW)
            pipe.execute()
            return True
        except Exception as e:
            logger.error(f"Redis 2FA rate limit check failed: {e}. Using fallback.")

    # Fallback to in-memory
    now = time.time()
    attempts = _2fa_rate_limit_fallback.get(user_id, [])
    attempts = [t for t in attempts if now - t < _2FA_WINDOW]
    if len(attempts) >= _2FA_LIMIT:
        return False
    attempts.append(now)
    _2fa_rate_limit_fallback[user_id] = attempts
    return True


class TwoFactorVerifyRequest(BaseModel):
    code: str


class TwoFactorConfirmSetupRequest(BaseModel):
    code: str


@app.post("/auth/2fa/setup")
async def setup_2fa(request: Request, db: Session = Depends(get_db)):
    """Generate TOTP secret and QR code for 2FA setup.
    Requires a setup token (needs_2fa_setup).
    """
    import pyotp
    import qrcode
    import io
    import base64

    user_id = verify_setup_token(request)
    db_user = db.query(User).filter(User.id == int(user_id)).first()
    if not db_user:
        raise HTTPException(status_code=404, detail="User not found")

    # Generate TOTP secret
    secret = pyotp.random_base32()

    # Store encrypted secret (not yet enabled)
    db_user.totp_secret = secret
    db.commit()

    # Generate provisioning URI for QR code
    totp = pyotp.TOTP(secret)
    provisioning_uri = totp.provisioning_uri(
        name=db_user.email,
        issuer_name="TAIC Companion"
    )

    # Generate QR code as base64 image
    qr = qrcode.QRCode(version=1, box_size=10, border=5)
    qr.add_data(provisioning_uri)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")

    buffer = io.BytesIO()
    img.save(buffer, format="PNG")
    buffer.seek(0)
    qr_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')

    return {
        "secret": secret,
        "qr_code": f"data:image/png;base64,{qr_base64}"
    }


@app.post("/auth/2fa/confirm-setup")
async def confirm_2fa_setup(
    body: TwoFactorConfirmSetupRequest,
    request: Request,
    response: Response,
    db: Session = Depends(get_db)
):
    """Confirm 2FA setup by verifying the first TOTP code.
    Activates 2FA, generates backup codes, returns full access token.
    """
    import pyotp

    user_id = verify_setup_token(request)
    db_user = db.query(User).filter(User.id == int(user_id)).first()
    if not db_user:
        raise HTTPException(status_code=404, detail="User not found")

    secret = db_user.totp_secret
    if not secret:
        raise HTTPException(status_code=400, detail="2FA setup not initiated")

    # Verify the code
    totp = pyotp.TOTP(secret)
    if not totp.verify(body.code, valid_window=1):
        raise HTTPException(status_code=400, detail="Invalid verification code")

    # Activate 2FA
    db_user.totp_enabled = True
    db_user.totp_setup_completed_at = datetime.utcnow()
    db.commit()

    # Issue full access token
    access_token = create_access_token(data={"sub": str(db_user.id)})

    response.set_cookie(
        key="token",
        value=access_token,
        httponly=True,
        secure=True,
        samesite="none",
        max_age=28800,
        path="/"
    )

    logger.info(f"2FA setup completed for user {db_user.username}")
    event_tracker.track_user_action(db_user.id, "2fa_setup_completed")

    return {
        "access_token": access_token,
        "token_type": "bearer"
    }


@app.post("/auth/2fa/verify")
async def verify_2fa(
    body: TwoFactorVerifyRequest,
    request: Request,
    response: Response,
    db: Session = Depends(get_db)
):
    """Verify TOTP code or backup code during login.
    Requires a pre_2fa token. Returns full access token on success.
    """
    import pyotp

    user_id = verify_pre_2fa_token(request)

    # Rate limiting per user
    if not _check_2fa_rate_limit(user_id):
        raise HTTPException(status_code=429, detail="Too many 2FA attempts. Please try again in 5 minutes.")

    db_user = db.query(User).filter(User.id == int(user_id)).first()
    if not db_user:
        raise HTTPException(status_code=404, detail="User not found")

    if not db_user.totp_enabled or not db_user.totp_secret:
        raise HTTPException(status_code=400, detail="2FA is not enabled")

    # Verify TOTP code
    secret = db_user.totp_secret
    totp = pyotp.TOTP(secret)
    if not totp.verify(body.code.strip(), valid_window=1):
        raise HTTPException(status_code=400, detail="Invalid verification code")

    # Issue full access token
    access_token = create_access_token(data={"sub": str(db_user.id)})

    response.set_cookie(
        key="token",
        value=access_token,
        httponly=True,
        secure=True,
        samesite="none",
        max_age=28800,
        path="/"
    )

    logger.info(f"2FA verified for user {db_user.username}")
    event_tracker.track_user_action(db_user.id, "2fa_verified")

    return {
        "access_token": access_token,
        "token_type": "bearer"
    }


@app.get("/auth/2fa/status")
async def get_2fa_status(request: Request, db: Session = Depends(get_db)):
    """Get 2FA status for the current user. Requires full access token."""
    from auth import verify_token_from_cookie

    user_id = verify_token_from_cookie(request)
    db_user = db.query(User).filter(User.id == int(user_id)).first()
    if not db_user:
        raise HTTPException(status_code=404, detail="User not found")

    return {
        "totp_enabled": db_user.totp_enabled,
        "setup_completed_at": db_user.totp_setup_completed_at.isoformat() if db_user.totp_setup_completed_at else None
    }


@app.post("/logout")
async def logout(response: Response):
    """Logout user by clearing HttpOnly cookie

    Clears the authentication cookie to log out the user.
    """
    response.delete_cookie(key="token", path="/")
    return {"message": "Logged out successfully"}


# Nouvelle version de l'endpoint /ask : utilise toujours la mémoire (historique) et le modèle fine-tuné si dispo

@app.post("/ask")
async def ask_question(
    request: QuestionRequestValidated,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Ask question to RAG system (toujours avec mémoire et bon modèle)"""
    if not _check_api_rate_limit(user_id, "ask", _API_ASK_LIMIT):
        raise HTTPException(status_code=429, detail="Too many requests. Please try again later.")
    start_time = time.time()
    try:
        logger.info(f"Processing question from user {user_id}: {request.question}")
        logger.info(f"Selected documents: {request.selected_documents}")

        # Récupérer l'historique complet de la conversation si conversation_id fourni
        history = []
        if hasattr(request, 'conversation_id') and request.conversation_id:
            msgs = db.query(Message).filter(Message.conversation_id == request.conversation_id).order_by(Message.created_at.asc()).all()
            history = [{"role": m.role, "content": m.content} for m in msgs]
        elif hasattr(request, 'history') and request.history:
            # fallback: si le frontend envoie déjà l'historique
            history = request.history

        answer = None
        agent = None
        model_id = None
        # Si agent_id fourni, comportement agent classique
        if request.agent_id:
            from database import Agent
            agent = _user_can_access_agent(int(user_id), request.agent_id, db)
            model_id = agent.finetuned_model_id or resolve_model_id(agent)
            logger.info(f"[LLM ROUTING] Agent '{agent.name}' type={getattr(agent, 'type', 'unknown')} -> model_id={model_id}")
            question_finale = request.question
            prompt = f"Sachant le contexte et la discussion en cours, réponds à cette question : {question_finale}"
            answer = get_answer(
                prompt,
                int(user_id),
                db,
                selected_doc_ids=request.selected_documents,
                agent_id=request.agent_id,
                history=history,
                model_id=model_id
            )
        # Si team_id fourni, on va chercher le chef d'équipe et les sous-agents
        elif request.team_id:
            from database import Team, Agent
            import numpy as np
            team = db.query(Team).filter(Team.id == request.team_id).first()
            if not team:
                raise HTTPException(status_code=404, detail="Team not found")
            if team.user_id != int(user_id):
                raise HTTPException(status_code=403, detail="Access denied to this team")
            leader = db.query(Agent).filter(Agent.id == team.leader_agent_id).first()
            if not leader:
                raise HTTPException(status_code=404, detail="Leader agent not found")

            # Récupérer TOUS les sous-agents (actionnables ET conversationnels)
            all_sub_agent_ids = []
            if team.action_agent_ids:
                try:
                    action_ids = json.loads(team.action_agent_ids) if isinstance(team.action_agent_ids, str) else team.action_agent_ids
                    all_sub_agent_ids.extend(action_ids)
                except:
                    pass

            # Récupérer les sous-agents depuis la base
            sub_agents = []
            if all_sub_agent_ids:
                sub_agents = db.query(Agent).filter(Agent.id.in_(all_sub_agent_ids)).all()

            best_agent = None
            best_score = -1

            # Si on a des sous-agents avec embeddings, faire le matching sémantique
            if sub_agents:
                try:
                    prompt_embedding = get_embedding(request.question)
                    for a in sub_agents:
                        if not a.embedding:
                            continue
                        try:
                            emb = np.array(json.loads(a.embedding))
                            score = float(np.dot(prompt_embedding, emb) / (np.linalg.norm(prompt_embedding) * np.linalg.norm(emb)))
                            if score > best_score:
                                best_score = score
                                best_agent = a
                        except Exception:
                            continue
                except Exception as e:
                    logger.warning(f"Error during semantic matching: {e}")

            # Si aucun sous-agent qualifié, utiliser le leader directement
            if not best_agent:
                logger.info(f"No sub-agent matched, using leader agent {leader.name}")
                best_agent = leader

            # Appel get_answer avec l'agent sélectionné
            if best_agent.finetuned_model_id:
                model_id = best_agent.finetuned_model_id
            else:
                atype = getattr(best_agent, 'type', 'conversationnel')
                if atype == 'recherche_live':
                    model_id = os.getenv('PERPLEXITY_MODEL', 'perplexity:sonar')
                else:
                    model_id = os.getenv('MISTRAL_MODEL', 'mistral:mistral-small-latest')
            prompt = f"Sachant le contexte et la discussion en cours, réponds à cette question : {request.question}"
            agent_answer = get_answer(
                prompt,
                int(user_id),
                db,
                selected_doc_ids=request.selected_documents,
                agent_id=best_agent.id,
                history=history,
                model_id=model_id
            )

            # Réponse formatée
            if best_agent.id == leader.id:
                answer = agent_answer
            else:
                answer = f"Pour répondre à votre question, j'ai fait appel à l'agent {best_agent.name}. Voici sa réponse :\n{agent_answer}"
            agent = best_agent

        if answer is None:
            raise HTTPException(status_code=400, detail="Aucun agent ou équipe valide fourni.")

        response_time = time.time() - start_time
        logger.info(f"Question answered for user {user_id} in {response_time:.2f}s")
        event_tracker.track_question_asked(int(user_id), request.question, response_time)
        return {"answer": answer}
    except Exception as e:
        logger.error(f"Error answering question for user {user_id}: {e}")
        return {"answer": "Désolé, une erreur s'est produite lors du traitement de votre question. Veuillez réessayer."}


_ACTIONNABLE_REMOVED = """
                    {
                        "name": "create_google_doc",
                        "description": "Create a Google Doc and return its URL",
                        "parameters": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "content": {"type": "string"},
                                "folder_id": {"type": "string"}
                            },
                            "required": ["title"]
                        }
                    },
                    {
                        "name": "create_google_sheet",
                        "description": "Create a Google Sheet and optionally populate structured sheets",
                        "parameters": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "sheets": {
                                    "type": "array",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "title": {"type": "string"},
                                            "headers": {"type": "array", "items": {"type": "string"}},
                                            "rows": {"type": "array", "items": {"type": "array", "items": {"type": ["string","number"]}}},
                                            "formulas": {"type": "array"},
                                            "conditional_formats": {"type": "array"}
                                        },
                                        "required": ["title","headers"]
                                    }
                                },
                                "folder_id": {"type": "string"}
                            },
                            "required": ["title","sheets"]
                        }
                    },
                    {
                        "name": "echo",
                        "description": "Echo back a message",
                        "parameters": {"type": "object", "properties": {"text": {"type": "string"}}}
                    },
                    {
                        "name": "write_local_file",
                        "description": "Write a local debug file on the server",
                        "parameters": {"type": "object", "properties": {"filename": {"type": "string"}, "content": {"type": "string"}}}
                    }
                ]

                # Build messages for the structured call: include a couple of short examples (few-shot)
                struct_messages = []
                # Strong instruction for models (Gemini) used by actionnable agents:
                # If an action is appropriate, respond with ONLY a JSON object in the shape
                # {"name": "<action_name>", "arguments": {...}} and nothing else. If no action
                # is required, reply with normal assistant text.
                struct_messages.append({
                    "role": "system",
                    "content": (
                        "Quand une action doit être exécutée, réponds STRICTEMENT et UNIQUEMENT avec un objet JSON de la forme :\n"
                        "{\n  \"function_call\": {\n    \"name\": \"<nom_de_l_action>\",\n    \"arguments\": { ... }\n  }\n}\n"
                        "N’ajoute aucune explication, texte ou commentaire. Si aucune action n’est requise, réponds avec un texte normal d’assistant."
                    )
                })
                example_user_sheet = (
                    "Exemple: Crée un Google Sheet intitulé \"Tableau RH exemple\" avec une feuille 'Employés'"
                    " contenant les colonnes Nom, Département, Poste, Salaire mensuel et une ligne d'exemple: Alice, IT, Dev, 4000."
                )
                example_assistant_sheet = (
                    '{'
                    '"function_call": {'
                    '"name": "create_google_sheet",'
                    '"arguments": {'
                    '"title":"Tableau RH exemple",'
                    '"sheets":[{"title":"Employés","headers":["Nom","Département","Poste","Salaire mensuel"],'
                    '"rows":[["Alice","IT","Dev",4000]]}]'
                    '}'
                    '}'
                    '}'
                )
                example_user_doc = (
                    "Exemple: Crée un Google Doc intitulé \"Note projet\" qui contient un court résumé et des actions à mener."
                )
                example_assistant_doc = (
                    '{'
                    '"function_call": {'
                    '"name": "create_google_doc",'
                    '"arguments": {'
                    '"title":"Note projet",'
                    '"content":"Résumé: ...\\nActions:\\n- Action 1\\n- Action 2"'
                    '}'
                    '}'
                    '}'
                )

                if agent and getattr(agent, 'contexte', None):
                    struct_messages.append({"role": "system", "content": agent.contexte})
                struct_messages.append({"role": "user", "content": example_user_sheet})
                struct_messages.append({"role": "assistant", "content": example_assistant_sheet})
                struct_messages.append({"role": "user", "content": example_user_doc})
                struct_messages.append({"role": "assistant", "content": example_assistant_doc})

                # For actionnable agents (Gemini), normalize the assistant answer to plain text for the structured call
                try:
                    assistant_context = normalize_model_output(answer)
                except Exception:
                    assistant_context = str(answer)
                struct_messages.append({"role": "assistant", "content": assistant_context})
                struct_messages.append({"role": "user", "content": request.question})

                # Enforce Gemini-only for actionnable agents' structured calls
                struct_gemini_only = True

                # Heuristic: if the user's question strongly indicates a spreadsheet/table request
                # prefer calling create_google_sheet; if it indicates a document, prefer create_google_doc.
                forced_call = None
                try:
                    q = (request.question or "").lower()
                    sheet_keywords = ["table", "tableau", "spreadsheet", "sheet", "tableur", "excel", "csv"]
                    doc_keywords = ["document", "doc", "note", "rapport", "résumé", "resume"]
                    if any(k in q for k in sheet_keywords):
                        forced_call = {"name": "create_google_sheet"}
                    elif any(k in q for k in doc_keywords):
                        forced_call = {"name": "create_google_doc"}
                except Exception:
                    forced_call = None

                message = get_chat_response_structured(struct_messages, functions=functions, function_call=forced_call, model_id=model_id, gemini_only=struct_gemini_only)

                action_results = []
                # If model requested a function call, execute it
                if hasattr(message, 'function_call') and message.function_call:
                    fc = message.function_call
                else:
                    # Fallback: some providers (or Gemini text outputs) return JSON text rather
                    # than a structured function_call attribute. Try to parse message.content
                    # as JSON and extract a function_call-shaped object.
                    fc = None
                    try:
                        raw_text = getattr(message, 'content', None) if hasattr(message, 'content') else str(message)
                        logger.info(f"Structured call: no function_call attribute; raw model output: {raw_text}")
                        # Try a more robust JSON extraction/parsing from the raw text
                        try:
                            parsed = extract_json_object_from_text(raw_text)
                            if parsed is None:
                                logger.debug("No JSON object could be extracted from raw model output")
                            else:
                                logger.debug(f"Parsed JSON object from model output (type={type(parsed).__name__})")
                                if isinstance(parsed, dict):
                                    if 'function_call' in parsed and isinstance(parsed['function_call'], dict):
                                        fc = parsed['function_call']
                                    elif 'name' in parsed and ('arguments' in parsed or 'params' in parsed or 'parameters' in parsed):
                                        fc = {'name': parsed.get('name'), 'arguments': parsed.get('arguments') or parsed.get('params') or parsed.get('parameters')}
                                    elif 'action' in parsed and 'params' in parsed:
                                        fc = {'name': parsed.get('action'), 'arguments': parsed.get('params')}
                                    else:
                                        # If parsed dict looks like arguments only, attempt to infer action name
                                        # by searching for common action names in the raw text.
                                        for candidate in ('create_google_sheet', 'create_google_doc', 'echo', 'write_local_file'):
                                            if candidate in raw_text:
                                                fc = {'name': candidate, 'arguments': parsed}
                                                break
                        except Exception as e:
                            logger.debug(f"Could not parse raw model output as JSON function_call fallback: {e}")
                    except Exception as e:
                        logger.debug(f"Could not parse raw model output as JSON function_call fallback: {e}")
                    try:
                        if isinstance(fc, dict):
                            name = fc.get('name')
                            arguments = fc.get('arguments')
                        else:
                            name = getattr(fc, 'name', None) or getattr(fc, 'function_name', None)
                            arguments = getattr(fc, 'arguments', None) or getattr(fc, 'params', None)
                        payload = {"name": name, "arguments": arguments}
                        logger.info(f"Prepared action payload: {payload}")

                        try:
                            args_obj = payload.get("arguments")
                            if isinstance(args_obj, str):
                                try:
                                    args_parsed = json.loads(args_obj)
                                except Exception:
                                    args_parsed = {"_raw": args_obj}
                            elif isinstance(args_obj, dict):
                                args_parsed = args_obj
                            else:
                                args_parsed = {"_raw": str(args_obj)}

                            if isinstance(args_parsed, dict) and ("_raw" not in args_parsed or not args_parsed.get("_raw")):
                                args_parsed["_raw"] = request.question

                            payload["arguments"] = args_parsed
                        except Exception:
                            pass

                        # If create_google_sheet, attempt validation and one repair pass
                        try:
                            if payload.get("name") == "create_google_sheet":
                                try:
                                    import jsonschema
                                    sheet_schema = {
                                        "type": "object",
                                        "properties": {
                                            "title": {"type": "string"},
                                            "sheets": {"type": "array"},
                                            "folder_id": {"type": "string"}
                                        },
                                        "required": ["title","sheets"]
                                    }
                                    if isinstance(payload.get("arguments"), dict):
                                        jsonschema.validate(payload.get("arguments"), sheet_schema)
                                    else:
                                        raise jsonschema.ValidationError("Arguments not an object")
                                except Exception:
                                    try:
                                        from openai_client import get_chat_response_json
                                        repair_msgs = [
                                            {"role": "system", "content": "You must return ONLY a JSON object that matches the requested schema for the function arguments. No explanation."},
                                            {"role": "user", "content": f"The user asked: {request.question}\nPlease return the function arguments JSON that matches this schema: {json.dumps(sheet_schema, ensure_ascii=False)}"}
                                        ]
                                        corrected = get_chat_response_json(repair_msgs, schema=sheet_schema, model_id=model_id, gemini_only=True)
                                        if isinstance(corrected, dict):
                                            payload["arguments"] = corrected
                                    except Exception:
                                        pass
                        except Exception:
                            pass
                    except Exception:
                        payload = {"name": None, "arguments": None}

                    if payload.get("name"):
                        result = parse_and_execute_actions(payload, db=db, agent_id=request.agent_id, user_id=int(user_id))
                        action_results.append({"action": payload.get("name"), "result": result})
                    else:
                        action_results.append({"status": "error", "error": "Could not parse function_call from model response"})

                # If any action succeeded, ask the model to generate a short assistant confirmation
                try:
                    confirmations = []
                    for ar in action_results:
                        res = ar.get("result") if isinstance(ar, dict) else None
                        if isinstance(res, dict) and res.get("status") == "ok":
                            payload_result = res.get("result") if isinstance(res.get("result"), dict) else None
                            if payload_result:
                                url = payload_result.get("url") or payload_result.get("webViewLink")
                                if url:
                                    confirmations.append(url)
                                    continue
                                doc_id = payload_result.get("document_id") or payload_result.get("spreadsheetId")
                                if doc_id:
                                    confirmations.append(str(doc_id))
                                    continue
                                path = payload_result.get("path")
                                if path:
                                    confirmations.append(path)

                    if confirmations:
                        messages_for_model = []
                        if agent and getattr(agent, 'contexte', None):
                            messages_for_model.append({"role": "system", "content": agent.contexte})
                        messages_for_model.append({"role": "assistant", "content": answer})
                        links_text = "\n".join([f"- {u}" for u in confirmations])
                        user_instruction = (
                            f'Suite à ce prompt "{request.question}" les actions suivantes ont été exécutées :\n'
                            + links_text + "\n\n"
                            + "Génère une réponse d'assistant courte et affirmative en français confirmant que l'action a été réalisée."
                        )
                        messages_for_model.append({"role": "user", "content": user_instruction})
                        try:
                            crafted = get_chat_response(messages_for_model, model_id=model_id, gemini_only=True)
                            # Normalize and strip markdown-style links so the frontend shows full URLs
                            try:
                                raw_crafted = normalize_model_output(crafted)
                            except Exception:
                                raw_crafted = str(crafted)
                            try:
                                                        # Replace markdown links [text](url) with the raw url
                                raw_crafted = re.sub(r"\[([^\]]+)\]\((https?://[^)]+)\)", r"\2", raw_crafted)
                                # Also replace simple HTML anchors <a href="url">text</a>
                                raw_crafted = re.sub(r"<a\s+[^>]*href=[\"'](https?://[^\"']+)[\"'][^>]*>.*?<\/a>", r"\1", raw_crafted, flags=re.IGNORECASE)
                            except Exception:
                                pass
                            answer = raw_crafted
                        except Exception:
                            answer = answer + "\n\n" + "\n".join([f"Action exécutée : {u}" for u in confirmations])
                except Exception:
                    pass

                try:
                    answer = normalize_model_output(answer)
                except Exception:
                    answer = str(answer)

                return {"answer": answer, "action_results": action_results}
            except Exception as e:
                logger.error(f"Error while checking/executing actions: {e}")
                try:
                    answer = normalize_model_output(answer)
                except Exception:
                    answer = str(answer)
                return {"answer": answer, "action_results": [{"status": "error", "error": "Action execution failed"}]}
        else:
            # Non-actionnable agents: do not attempt function-calling or action execution; return the original answer
"""


@app.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Upload and process document for a specific agent, extracting full text from any supported file type

    Security: Validates file size, extension, and sanitizes filename.
    """
    if not _check_api_rate_limit(user_id, "upload", _API_UPLOAD_LIMIT):
        raise HTTPException(status_code=429, detail="Too many uploads. Please try again later.")
    logger.info(f"Appel reçu sur /upload : filename={file.filename if file else 'None'}")
    try:
        # Validate filename exists
        if not file.filename:
            raise HTTPException(status_code=400, detail="Filename is required")

        # Sanitize filename to prevent path traversal attacks
        original_filename = file.filename
        safe_filename = sanitize_filename(original_filename)

        # Validate file extension against whitelist
        if not validate_file_extension(safe_filename):
            allowed_exts = ', '.join(sorted(ALLOWED_FILE_EXTENSIONS))
            raise HTTPException(
                status_code=400,
                detail=f"File type not allowed. Allowed types: {allowed_exts}"
            )

        logger.info(f"Début import PJ : filename={safe_filename}, content_type={file.content_type if hasattr(file, 'content_type') else 'unknown'}")

        # Read content and validate size
        content = await file.read()
        content_size = len(content)

        # Validate file size
        if not validate_file_size(content_size):
            max_size_mb = MAX_FILE_SIZE / (1024 * 1024)
            raise HTTPException(
                status_code=413,
                detail=f"File too large (max {max_size_mb:.0f}MB)"
            )

        # Validate file content matches extension (magic bytes)
        if not validate_file_content(content, safe_filename):
            raise HTTPException(
                status_code=400,
                detail="File content does not match its extension"
            )

        logger.info(f"PJ reçue : filename={safe_filename}, taille={content_size} octets")
        filename = safe_filename.lower()
        text = None

        if filename.endswith('.pdf'):
            MAX_PDF_PAGES = 500
            logger.info("Tentative extraction PDF (pdfplumber)")
            import pdfplumber
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
                tmp.write(content)
                tmp.flush()
                logger.info(f"Fichier PDF temporaire créé : {tmp.name}")
                try:
                    with pdfplumber.open(tmp.name) as pdf:
                        if len(pdf.pages) > MAX_PDF_PAGES:
                            raise HTTPException(status_code=400, detail=f"PDF too large ({len(pdf.pages)} pages, max {MAX_PDF_PAGES})")
                        text = '\n'.join([page.extract_text() or '' for page in pdf.pages])
                    logger.info(f"Texte PDF extrait (pdfplumber) : longueur={len(text) if text else 0}, aperçu='{text[:200] if text else ''}'")
                except Exception as e:
                    logger.error(f"PDF extraction error: {e}")
                # Si le texte est vide, tente l'OCR sur chaque page
                if not text or not text.strip():
                    logger.info("PDF vide ou non textuel, tentative OCR (pytesseract)")
                    try:
                        from PIL import Image
                        import pytesseract
                        with pdfplumber.open(tmp.name) as pdf:
                            ocr_text = ''
                            for i, page in enumerate(pdf.pages):
                                img = page.to_image(resolution=300)
                                pil_img = img.original
                                page_ocr = pytesseract.image_to_string(pil_img, lang='fra')
                                logger.info(f"OCR page {i+1}: longueur={len(page_ocr)}, aperçu='{page_ocr[:100]}'")
                                ocr_text += page_ocr + '\n'
                        text = ocr_text
                        logger.info(f"OCR PDF extrait: longueur={len(text)}, aperçu='{text[:200]}'")
                    except Exception as e:
                        logger.error(f"PDF OCR extraction error: {e}")
            os.unlink(tmp.name)
        elif filename.endswith('.docx'):
            logger.info("Tentative extraction DOCX")
            from docx import Document as DocxDocument
            import tempfile
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
                tmp.write(content)
                tmp.flush()
                logger.info(f"Fichier DOCX temporaire créé : {tmp.name}")
                try:
                    doc = DocxDocument(tmp.name)
                    text = '\n'.join([p.text for p in doc.paragraphs])
                    logger.info(f"Texte DOCX extrait : longueur={len(text) if text else 0}, aperçu='{text[:200] if text else ''}'")
                except Exception as e:
                    logger.error(f"DOCX extraction error: {e}")
            os.unlink(tmp.name)
        elif filename.endswith('.pptx'):
            logger.info("Tentative extraction PPTX")
            from pptx import Presentation
            import tempfile
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
                tmp.write(content)
                tmp.flush()
                logger.info(f"Fichier PPTX temporaire créé : {tmp.name}")
                try:
                    pres = Presentation(tmp.name)
                    text = '\n'.join([shape.text for slide in pres.slides for shape in slide.shapes if hasattr(shape, "text")])
                    logger.info(f"Texte PPTX extrait : longueur={len(text) if text else 0}, aperçu='{text[:200] if text else ''}'")
                except Exception as e:
                    logger.error(f"PPTX extraction error: {e}")
            os.unlink(tmp.name)
        elif filename.endswith('.xlsx'):
            logger.info("Tentative extraction XLSX")
            import openpyxl
            import tempfile
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
                tmp.write(content)
                tmp.flush()
                logger.info(f"Fichier XLSX temporaire créé : {tmp.name}")
                try:
                    wb = openpyxl.load_workbook(tmp.name, data_only=True)
                    text = ''
                    for sheet in wb.worksheets:
                        for row in sheet.iter_rows(values_only=True):
                            text += '\t'.join([str(cell) if cell is not None else '' for cell in row]) + '\n'
                    logger.info(f"Texte XLSX extrait : longueur={len(text) if text else 0}, aperçu='{text[:200] if text else ''}'")
                except Exception as e:
                    logger.error(f"XLSX extraction error: {e}")
            os.unlink(tmp.name)
        elif filename.endswith('.txt') or filename.endswith('.csv') or filename.endswith('.ics'):
            logger.info("Tentative extraction fichier texte/csv/ics")
            try:
                text = content.decode('utf-8', errors='ignore')
                logger.info(f"Texte fichier texte/csv/ics extrait : longueur={len(text) if text else 0}, aperçu='{text[:200] if text else ''}'")
            except Exception as e:
                logger.error(f"Text file decode error: {e}")
        else:
            raise HTTPException(status_code=400, detail="File type not supported")


        logger.info(f"Texte extrait de la PJ ({file.filename}): longueur={len(text) if text else 0}, aperçu='{text[:200] if text else ''}'")
        if not text or not text.strip():
            raise HTTPException(status_code=400, detail="Aucun texte détecté dans la pièce jointe. Vérifiez que le document contient du texte sélectionnable (pas une image ou un scan).")

        # Process document with extracted text
        doc_id = process_document_for_user(file.filename, text.encode('utf-8', errors='ignore'), int(user_id), db, agent_id=None)

        logger.info(f"Document uploaded for user {user_id}: {file.filename}")
        event_tracker.track_document_upload(int(user_id), file.filename, len(text))

        return {"filename": file.filename, "document_id": doc_id, "status": "uploaded"}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error uploading document: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")

@app.post("/upload-agent")
async def upload_file_for_agent(
    request: Request,
    file: UploadFile = File(...),
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Upload and process document for a specific agent"""
    if not _check_api_rate_limit(user_id, "upload", _API_UPLOAD_LIMIT):
        raise HTTPException(status_code=429, detail="Too many uploads. Please try again later.")
    try:
        # Get agent_id from form data
        form = await request.form()
        logger.info(f"Form data received in /upload-agent: {dict(form)}")
        agent_id = form.get("agent_id")
        # Fallback: extract agent_id from 'data' if present (Zapier edge case)
        if not agent_id and "data" in form:
            # Try to parse agent_id from string like 'agent_id=23'
            data_value = form.get("data")
            if isinstance(data_value, str) and data_value.startswith("agent_id="):
                agent_id = data_value.split("=", 1)[1]
        
        if not agent_id:
            logger.error(f"agent_id missing in form: {dict(form)}")
            raise HTTPException(status_code=400, detail="agent_id is required")
        
        agent_id = int(agent_id)
        # Check file size (10MB limit)
        if file.size > 10 * 1024 * 1024:
            raise HTTPException(status_code=413, detail="File too large (max 10MB)")
        
        # Check file type
        allowed_types = ['.pdf', '.txt', '.docx', '.ics']
        if not any(file.filename.lower().endswith(ext) for ext in allowed_types):
            raise HTTPException(status_code=400, detail="File type not supported")
        
        # Verify agent belongs to the user or user has edit permission
        agent = _user_can_edit_agent(int(user_id), agent_id, db)
        
        content = await file.read()
        doc_id = process_document_for_user(file.filename, content, int(user_id), db, agent_id)
        
        logger.info(f"Document uploaded for user {user_id}, agent {agent_id}: {file.filename}")
        event_tracker.track_document_upload(int(user_id), file.filename, len(content))
        
        return {"filename": file.filename, "document_id": doc_id, "agent_id": agent_id, "status": "uploaded"}
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error uploading document: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")

# Security: Debug endpoints /test-jwt, /test-auth, /test-openai, /debug/whoami have been removed
# These endpoints exposed sensitive information (JWT secret length, ADC credentials, API connectivity)
# and should not exist in production code

@app.get("/user/documents")
async def get_user_documents(
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db),
    agent_id: int = None
):
    """Get user's documents, optionally filtered by agent"""
    try:
        logger.info(f"Fetching documents for user {user_id}, agent {agent_id}")

        # Build query
        uid = int(user_id)
        if agent_id is not None:
            # Check if user has edit access to this agent (owner or can_edit share)
            share = db.query(AgentShare).filter(
                AgentShare.agent_id == agent_id,
                AgentShare.user_id == uid,
                AgentShare.can_edit == True
            ).first()
            if share:
                # Shared user with edit access: show docs for the agent
                query = db.query(Document).filter(Document.agent_id == agent_id)
            else:
                query = db.query(Document).filter(Document.user_id == uid, Document.agent_id == agent_id)
        else:
            query = db.query(Document).filter(Document.user_id == uid)
        
        documents = query.order_by(Document.created_at.desc()).all()
        logger.info(f"Found {len(documents)} documents for user {user_id}, agent {agent_id}")
        
        result = []

        for doc in documents:
            try:
                doc_data = {
                    "id": doc.id,
                    "filename": doc.filename,
                    "created_at": doc.created_at.isoformat(),
                    "gcs_url": doc.gcs_url,
                    "notion_link_id": doc.notion_link_id
                }
                # Safely try to add agent_id if it exists
                if hasattr(doc, 'agent_id'):
                    doc_data["agent_id"] = doc.agent_id
                result.append(doc_data)
            except Exception as doc_error:
                logger.error(f"Error processing document {doc.id}: {doc_error}")
                continue

        return {"documents": result}
        
    except Exception as e:
        logger.error(f"Error fetching documents: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Internal server error")

@app.delete("/documents/{document_id}")
async def delete_document(
    document_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Delete a user's document"""
    try:
        # Check if document exists and belongs to user
        document = db.query(Document).filter(
            Document.id == document_id,
            Document.user_id == int(user_id)
        ).first()
        
        if not document:
            raise HTTPException(status_code=404, detail="Document not found")
        
        # Delete document
        db.delete(document)
        db.commit()
        
        logger.info(f"Document {document_id} deleted by user {user_id}")
        event_tracker.track_user_action(int(user_id), f"document_deleted:{document.filename}")
        
        return {"message": "Document deleted successfully"}
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error deleting document: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")

# Endpoints pour les agents
@app.get("/agents")
async def get_agents(
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get user's own agents + agents shared with them"""
    try:
        uid = int(user_id)
        own_agents = db.query(Agent).filter(Agent.user_id == uid).order_by(Agent.created_at.desc()).all()
        result = [{
            "id": a.id, "name": a.name, "type": a.type, "statut": a.statut,
            "profile_photo": a.profile_photo, "llm_provider": a.llm_provider,
            "neo4j_enabled": a.neo4j_enabled, "email_tags": a.email_tags,
            "weekly_recap_enabled": a.weekly_recap_enabled,
            "created_at": a.created_at.isoformat() if a.created_at else None,
            "shared": False
        } for a in own_agents]

        # Add agents shared with this user
        shared = db.query(Agent, User, AgentShare).join(AgentShare, AgentShare.agent_id == Agent.id).join(
            User, Agent.user_id == User.id
        ).filter(AgentShare.user_id == uid).order_by(Agent.created_at.desc()).all()
        for a, owner, share in shared:
            result.append({
                "id": a.id, "name": a.name, "type": a.type, "statut": a.statut,
                "profile_photo": a.profile_photo, "llm_provider": a.llm_provider,
                "neo4j_enabled": a.neo4j_enabled, "email_tags": a.email_tags,
                "weekly_recap_enabled": a.weekly_recap_enabled,
                "created_at": a.created_at.isoformat() if a.created_at else None,
                "shared": True,
                "can_edit": share.can_edit,
                "owner_username": owner.username
            })

        return {"agents": result}
    except Exception as e:
        logger.error(f"Error getting agents: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")


@app.post("/agents")
async def create_agent(
    name: str = Form(...),
    contexte: str = Form(None),
    biographie: str = Form(None),
    type: str = Form("conversationnel"),
    email_tags: str = Form(None),  # JSON array ou liste séparée par virgules
    neo4j_enabled: str = Form("false"),
    neo4j_person_name: str = Form(None),
    neo4j_depth: str = Form("1"),
    weekly_recap_enabled: str = Form("false"),
    profile_photo: UploadFile = File(None),
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Create a new agent with optional profile photo upload"""
    try:
        logger.info(f"[CREATE_AGENT] Champs reçus: name={name}, contexte={contexte}, biographie={biographie}, type={type}, profile_photo={profile_photo.filename if profile_photo else None}, user_id={user_id}")
        # --- GCS UPLOAD UTILS ---
        GCS_BUCKET_NAME = os.getenv("GCS_BUCKET_NAME", "applydi-agent-photos")
        def upload_profile_photo_to_gcs(file: UploadFile) -> str:
            """Upload a file to Google Cloud Storage and return its public URL."""
            client = storage.Client()
            bucket = client.bucket(GCS_BUCKET_NAME)
            safe_name = sanitize_filename(file.filename)
            filename = f"{int(time.time())}_{safe_name}"
            blob = bucket.blob(filename)
            blob.upload_from_file(file.file, content_type=file.content_type)
            try:
                # Make the object publicly readable so the browser can load it directly
                blob.make_public()
            except Exception:
                logger.exception("Failed to make uploaded profile photo public; object may remain private")
            public_url = blob.public_url
            logger.info(f"Uploaded profile photo to GCS and set public URL: {public_url}")
            return public_url

        photo_url = None
        if profile_photo is not None:
            try:
                photo_url = upload_profile_photo_to_gcs(profile_photo)
                logger.info(f"[CREATE_AGENT] Photo de profil uploadée sur GCS: {photo_url}")
            except Exception as file_err:
                logger.error(f"[CREATE_AGENT] Erreur lors de l'upload GCS: {file_err}")
                raise HTTPException(status_code=500, detail="Erreur lors de l'upload de la photo")

        # Parser les email_tags
        parsed_email_tags = None
        if email_tags:
            try:
                # Essayer de parser comme JSON
                parsed_email_tags = json.loads(email_tags)
            except json.JSONDecodeError:
                # Sinon, traiter comme liste séparée par virgules
                tags_list = [t.strip() for t in email_tags.split(',') if t.strip()]
                # Normaliser avec @ prefix
                parsed_email_tags = [f"@{t.lstrip('@').lower()}" for t in tags_list]
            parsed_email_tags = json.dumps(parsed_email_tags) if parsed_email_tags else None

        # Auto-calculate llm_provider from type
        effective_llm_provider = resolve_llm_provider(type)
        db_agent = Agent(
            name=name,
            contexte=contexte,
            biographie=biographie,
            profile_photo=photo_url,
            statut="privé",
            type=type,
            llm_provider=effective_llm_provider,
            email_tags=parsed_email_tags,
            neo4j_enabled=neo4j_enabled.lower() in ("true", "1", "yes"),
            neo4j_person_name=neo4j_person_name if neo4j_person_name and neo4j_person_name.strip() else None,
            neo4j_depth=int(neo4j_depth) if neo4j_depth else 1,
            weekly_recap_enabled=weekly_recap_enabled.lower() in ("true", "1", "yes"),
            user_id=int(user_id)
        )
        db.add(db_agent)
        db.commit()
        db.refresh(db_agent)
        # Génère et stocke l'embedding si le contexte n'est pas vide
        if db_agent.contexte and db_agent.contexte.strip():
            update_agent_embedding(db_agent, db)
        logger.info(f"[CREATE_AGENT] Agent créé avec succès: id={db_agent.id}, statut={db_agent.statut}")
        return {"agent": db_agent}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[CREATE_AGENT] Erreur inattendue: {e}")
        raise HTTPException(status_code=500, detail="Erreur lors de la création de l'agent")

def _user_can_access_agent(user_id: int, agent_id: int, db: Session):
    """Return the agent if the user is owner OR has an AgentShare. Otherwise raise 403."""
    agent = db.query(Agent).filter(Agent.id == agent_id).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")
    if agent.user_id == user_id:
        return agent
    share = db.query(AgentShare).filter(
        AgentShare.agent_id == agent_id,
        AgentShare.user_id == user_id
    ).first()
    if share:
        return agent
    raise HTTPException(status_code=403, detail="Access denied to this agent")


def _user_can_edit_agent(user_id: int, agent_id: int, db: Session):
    """Return the agent if the user is owner OR has an AgentShare with can_edit=True. Otherwise raise 403."""
    agent = db.query(Agent).filter(Agent.id == agent_id).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")
    if agent.user_id == user_id:
        return agent
    share = db.query(AgentShare).filter(
        AgentShare.agent_id == agent_id,
        AgentShare.user_id == user_id,
        AgentShare.can_edit == True
    ).first()
    if share:
        return agent
    raise HTTPException(status_code=403, detail="You do not have edit permission on this agent")


def _delete_agent_and_related_data(agent: Agent, owner_user_id: int, db: Session):
    """Delete an agent and all its related data (conversations, actions, teams, shares)."""
    agent_id = agent.id

    # Delete all agent shares
    db.query(AgentShare).filter(AgentShare.agent_id == agent_id).delete()
    db.flush()

    # Delete all conversations related to this agent
    conversations = db.query(Conversation).filter(Conversation.agent_id == agent_id).all()
    for conv in conversations:
        db.delete(conv)
    db.flush()

    # Delete all agent actions related to this agent
    from database import AgentAction
    actions = db.query(AgentAction).filter(AgentAction.agent_id == agent_id).all()
    for action in actions:
        db.delete(action)
    db.flush()

    # Check if this agent is used in any teams
    teams = db.query(Team).filter(Team.user_id == owner_user_id).all()
    teams_to_delete = []
    for team in teams:
        if team.leader_agent_id == agent_id:
            logger.warning(f"Agent {agent_id} is leader of team {team.id}, deleting team")
            teams_to_delete.append(team)
        else:
            try:
                action_ids = json.loads(team.action_agent_ids) if team.action_agent_ids else []
                if agent_id in action_ids:
                    action_ids = [aid for aid in action_ids if aid != agent_id]
                    team.action_agent_ids = json.dumps(action_ids)
                    db.add(team)
            except Exception as e:
                logger.error(f"Error updating team {team.id}: {e}")

    if teams_to_delete:
        team_ids = [team.id for team in teams_to_delete]

        ct_result = db.execute(
            text("SELECT id FROM conversations_teams WHERE team_id = ANY(:team_ids)"),
            {"team_ids": team_ids}
        )
        ct_ids = [row[0] for row in ct_result]

        if ct_ids:
            db.execute(
                text("DELETE FROM messages_teams WHERE conversation_id = ANY(:ct_ids)"),
                {"ct_ids": ct_ids}
            )
            db.flush()
            db.execute(
                text("DELETE FROM conversations_teams WHERE id = ANY(:ct_ids)"),
                {"ct_ids": ct_ids}
            )
            db.flush()

        conv_result = db.execute(
            text("SELECT id FROM conversations WHERE team_id = ANY(:team_ids)"),
            {"team_ids": team_ids}
        )
        conv_ids = [row[0] for row in conv_result]

        if conv_ids:
            db.execute(
                text("DELETE FROM messages WHERE conversation_id = ANY(:conv_ids)"),
                {"conv_ids": conv_ids}
            )
            db.flush()

        db.execute(
            text("DELETE FROM conversations WHERE team_id = ANY(:team_ids)"),
            {"team_ids": team_ids}
        )
        db.flush()

        db.execute(
            text("DELETE FROM teams WHERE id = ANY(:team_ids)"),
            {"team_ids": team_ids}
        )
        db.flush()

    db.delete(agent)


@app.delete("/agents/{agent_id}")
async def delete_agent(
    agent_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Delete an agent and all related data"""
    try:
        agent = db.query(Agent).filter(
            Agent.id == agent_id,
            Agent.user_id == int(user_id)
        ).first()

        if not agent:
            raise HTTPException(status_code=404, detail="Agent not found")

        _delete_agent_and_related_data(agent, int(user_id), db)
        db.commit()

        return {"message": "Agent deleted successfully"}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error deleting agent: {e}")
        logger.error(f"Full traceback: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Internal server error")

@app.get("/agents/{agent_id}")
async def get_agent(
    agent_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get a specific agent (owner or shared user)"""
    try:
        uid = int(user_id)
        agent = _user_can_access_agent(uid, agent_id, db)
        is_owner = agent.user_id == uid

        # Check can_edit for shared agents
        can_edit = False
        if not is_owner:
            share = db.query(AgentShare).filter(
                AgentShare.agent_id == agent_id,
                AgentShare.user_id == uid
            ).first()
            can_edit = share.can_edit if share else False

        result = {
            "id": agent.id, "name": agent.name, "type": agent.type, "statut": agent.statut,
            "profile_photo": agent.profile_photo, "llm_provider": agent.llm_provider,
            "neo4j_enabled": agent.neo4j_enabled,
            "created_at": agent.created_at.isoformat() if agent.created_at else None,
            "shared": not is_owner,
            "can_edit": can_edit,
        }
        # Expose editable fields to owner OR shared user with can_edit
        if is_owner or can_edit:
            result.update({
                "contexte": agent.contexte, "biographie": agent.biographie,
                "neo4j_person_name": agent.neo4j_person_name,
                "neo4j_depth": agent.neo4j_depth, "email_tags": agent.email_tags,
                "weekly_recap_enabled": agent.weekly_recap_enabled,
            })
        if not is_owner:
            owner = db.query(User).filter(User.id == agent.user_id).first()
            result["owner_username"] = owner.username if owner else None

        return {"agent": result}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting agent: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")


@app.get("/teams")
async def list_teams(user_id: str = Depends(verify_token), db: Session = Depends(get_db)):
    """List teams for the current user."""
    try:
        # Return teams created by the current user
        teams = db.query(Team).filter(Team.user_id == int(user_id)).order_by(Team.created_at.desc()).all()
        out = []
        for t in teams:
            # Try to resolve leader and action agent names for convenience
            leader_name = None
            action_agent_names = []
            try:
                leader = db.query(Agent).filter(Agent.id == t.leader_agent_id).first()
                if leader:
                    leader_name = leader.name
            except Exception:
                leader_name = None
            try:
                ids = json.loads(t.action_agent_ids) if t.action_agent_ids else []
                for aid in ids:
                    a = db.query(Agent).filter(Agent.id == int(aid)).first()
                    if a:
                        action_agent_names.append(a.name)
            except Exception:
                action_agent_names = []

            out.append({
                "id": t.id,
                "name": t.name,
                "contexte": t.contexte,
                "leader_agent_id": t.leader_agent_id,
                "leader_name": leader_name,
                "action_agent_ids": json.loads(t.action_agent_ids) if t.action_agent_ids else [],
                "action_agent_names": action_agent_names,
                "created_at": t.created_at.isoformat() if t.created_at else None
            })
        return {"teams": out}
    except Exception as e:
        logger.exception(f"Error listing teams: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")


@app.post("/teams")
async def create_team(payload: TeamCreateValidated, user_id: str = Depends(verify_token), db: Session = Depends(get_db)):
    """Create a team. Expected payload: {name, contexte (opt), leader_agent_id, action_agent_ids: [id,id,id]}"""
    try:
        name = payload.name
        contexte = payload.contexte
        leader_agent_id = payload.leader_agent_id
        member_agent_ids = payload.action_agent_ids
        # On accepte n'importe quel nombre d'agents

        # Valider le chef (doit être conversationnel)
        leader = db.query(Agent).filter(Agent.id == int(leader_agent_id), Agent.user_id == int(user_id)).first()
        if not leader or getattr(leader, 'type', 'conversationnel') != 'conversationnel':
            raise HTTPException(status_code=400, detail="Leader agent must be a conversationnel agent belonging to you")

        # Valider les membres (uniquement conversationnels)
        member_agents = []
        for aid in member_agent_ids:
            a = db.query(Agent).filter(Agent.id == int(aid), Agent.user_id == int(user_id)).first()
            if not a or getattr(a, 'type', '') != 'conversationnel':
                raise HTTPException(status_code=400, detail=f"Agent {aid} doit être un agent conversationnel appartenant à vous")
            member_agents.append(a)

        team = Team(
            name=name,
            contexte=contexte,
            leader_agent_id=int(leader_agent_id),
            action_agent_ids=json.dumps([int(x) for x in member_agent_ids]),
            user_id=int(user_id)
        )
        db.add(team)
        db.commit()
        db.refresh(team)

        # Préparer la réponse avec les noms
        resp = {
            "team": {
                "id": team.id,
                "name": team.name,
                "contexte": team.contexte,
                "leader_agent_id": team.leader_agent_id,
                "leader_name": leader.name,
                "member_agent_ids": [int(x) for x in member_agent_ids],
                "member_agent_names": [a.name for a in member_agents],
                "created_at": team.created_at.isoformat() if team.created_at else None
            }
        }
        return resp
    except HTTPException:
        raise
    except Exception as e:
        logger.exception(f"Error creating team: {e}")
        # If the teams table does not exist, provide a helpful message
        if 'relation "teams"' in str(e) or 'does not exist' in str(e):
            raise HTTPException(status_code=500, detail="teams table not found in database. Please create the table before using this endpoint.")
        raise HTTPException(status_code=500, detail="Internal server error")


@app.get("/teams/{team_id}")
async def get_team(team_id: int, user_id: str = Depends(verify_token), db: Session = Depends(get_db)):
    try:
        t = db.query(Team).filter(Team.id == team_id, Team.user_id == int(user_id)).first()
        if not t:
            raise HTTPException(status_code=404, detail="Team not found")
        leader_name = None
        try:
            leader = db.query(Agent).filter(Agent.id == t.leader_agent_id).first()
            if leader:
                leader_name = leader.name
        except Exception:
            leader_name = None
        action_agent_names = []
        try:
            ids = json.loads(t.action_agent_ids) if t.action_agent_ids else []
            for aid in ids:
                a = db.query(Agent).filter(Agent.id == int(aid)).first()
                if a:
                    action_agent_names.append(a.name)
        except Exception:
            action_agent_names = []

        return {"team": {
            "id": t.id,
            "name": t.name,
            "contexte": t.contexte,
            "leader_agent_id": t.leader_agent_id,
            "leader_name": leader_name,
            "action_agent_ids": json.loads(t.action_agent_ids) if t.action_agent_ids else [],
            "action_agent_names": action_agent_names,
            "created_at": t.created_at.isoformat() if t.created_at else None
        }}
    except HTTPException:
        raise
    except Exception as e:
        logger.exception(f"Error fetching team {team_id}: {e}")
        raise HTTPException(status_code=500, detail="Internal server error")
# Endpoint pour modifier un agent existant
from fastapi import Form
@app.put("/agents/{agent_id}")
async def update_agent(
    agent_id: int,
    name: str = Form(...),
    contexte: str = Form(None),
    biographie: str = Form(None),
    type: str = Form("conversationnel"),
    email_tags: str = Form(None),  # JSON array ou liste séparée par virgules
    neo4j_enabled: str = Form("false"),
    neo4j_person_name: str = Form(None),
    neo4j_depth: str = Form("1"),
    weekly_recap_enabled: str = Form("false"),
    profile_photo: UploadFile = File(None),
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Met à jour un agent existant, y compris la photo de profil (GCS), le statut et les email_tags."""
    try:
        agent = _user_can_edit_agent(int(user_id), agent_id, db)

        agent.name = name
        agent.contexte = contexte
        agent.biographie = biographie
        agent.statut = "privé"
        agent.type = type
        agent.llm_provider = "perplexity" if type == "recherche_live" else "mistral"

        # Parser et mettre à jour les email_tags
        if email_tags is not None:
            if email_tags == "" or email_tags == "[]":
                agent.email_tags = None
            else:
                try:
                    parsed_tags = json.loads(email_tags)
                except json.JSONDecodeError:
                    tags_list = [t.strip() for t in email_tags.split(',') if t.strip()]
                    parsed_tags = [f"@{t.lstrip('@').lower()}" for t in tags_list]
                agent.email_tags = json.dumps(parsed_tags) if parsed_tags else None

        # Update Neo4j fields
        agent.neo4j_enabled = neo4j_enabled.lower() in ("true", "1", "yes")
        agent.neo4j_person_name = neo4j_person_name if neo4j_person_name and neo4j_person_name.strip() else None
        agent.neo4j_depth = int(neo4j_depth) if neo4j_depth else 1

        # Update Weekly Recap
        agent.weekly_recap_enabled = weekly_recap_enabled.lower() in ("true", "1", "yes")

        GCS_BUCKET_NAME = os.getenv("GCS_BUCKET_NAME", "applydi-agent-photos")
        def upload_profile_photo_to_gcs(file: UploadFile) -> str:
            client = storage.Client()
            bucket = client.bucket(GCS_BUCKET_NAME)
            safe_name = sanitize_filename(file.filename)
            filename = f"{int(time.time())}_{safe_name}"
            blob = bucket.blob(filename)
            blob.upload_from_file(file.file, content_type=file.content_type)
            try:
                blob.make_public()
            except Exception:
                logger.exception("Failed to make uploaded profile photo public; object may remain private")
            public_url = blob.public_url
            logger.info(f"Uploaded profile photo to GCS and set public URL: {public_url}")
            return public_url

        if profile_photo is not None:
            try:
                photo_url = upload_profile_photo_to_gcs(profile_photo)
                agent.profile_photo = photo_url
            except Exception as file_err:
                logger.error(f"[UPDATE_AGENT] Erreur lors de l'upload GCS: {file_err}")
                raise HTTPException(status_code=500, detail="Erreur lors de l'upload de la photo")

        db.commit()
        db.refresh(agent)
        logger.info(f"[UPDATE_AGENT] Agent modifié avec succès: id={agent.id}, statut={agent.statut}")
        return {"agent": agent}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[UPDATE_AGENT] Erreur inattendue: {e}")
        raise HTTPException(status_code=500, detail="Erreur lors de la mise à jour de l'agent")

## ---- Weekly Recap Endpoints ----

@app.post("/api/weekly-recap/trigger")
async def trigger_weekly_recap(
    request: Request,
    db: Session = Depends(get_db)
):
    """Trigger weekly recap for all enabled agents. Protected by X-API-Key."""
    api_key = request.headers.get("X-API-Key", "")
    expected_key = os.getenv("WEEKLY_RECAP_API_KEY", "")

    if not expected_key:
        raise HTTPException(status_code=500, detail="WEEKLY_RECAP_API_KEY not configured")
    if not api_key or not hmac.compare_digest(api_key, expected_key):
        raise HTTPException(status_code=401, detail="Invalid API Key")

    from weekly_recap import process_agent_recap

    agents = db.query(Agent).filter(Agent.weekly_recap_enabled == True).all()
    results = []
    for agent in agents:
        result = process_agent_recap(agent, db)
        results.append({"agent_id": agent.id, "agent_name": agent.name, **result})

    return {"processed": len(results), "results": results}


@app.post("/api/agents/{agent_id}/recap-preview")
async def recap_preview(
    agent_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Generate a recap preview without sending email. Auth required."""
    agent = db.query(Agent).filter(Agent.id == agent_id, Agent.user_id == int(user_id)).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")

    from weekly_recap import fetch_weekly_messages, fetch_traceability_documents, fetch_notion_content, build_recap_prompt, generate_recap_html, get_model_id_for_agent
    from openai_client import get_chat_response as _get_chat_response

    messages = fetch_weekly_messages(agent.id, db)
    docs = fetch_traceability_documents(agent.id, db)
    notion_pages = fetch_notion_content(agent.id, db)

    if not messages and not docs and not notion_pages:
        return {"status": "no_data", "message": "No messages or documents this week", "html": None}

    prompt_messages = build_recap_prompt(agent, messages, docs, notion_pages)
    model_id = get_model_id_for_agent(agent)
    recap_content = _get_chat_response(prompt_messages, model_id=model_id)
    html = generate_recap_html(agent.name, recap_content)

    return {
        "status": "success",
        "html": html,
        "message_count": len(messages),
        "doc_count": len(docs),
        "notion_count": len(notion_pages)
    }


@app.post("/api/agents/{agent_id}/recap-send")
async def recap_send(
    agent_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Send recap email now for a specific agent. Auth required."""
    agent = db.query(Agent).filter(Agent.id == agent_id, Agent.user_id == int(user_id)).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")
    if not agent.weekly_recap_enabled:
        raise HTTPException(status_code=400, detail="Weekly recap is not enabled for this agent")

    from weekly_recap import process_agent_recap
    result = process_agent_recap(agent, db)
    return result


## ---- Traceability Documents Endpoints ----

@app.post("/api/agents/{agent_id}/traceability-docs")
async def upload_traceability_doc(
    agent_id: int,
    file: UploadFile = File(...),
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Upload a traceability document (no chunking/embedding)."""
    agent = _user_can_edit_agent(int(user_id), agent_id, db)

    allowed_ext = {'.pdf', '.txt', '.docx', '.xlsx', '.xls', '.csv'}
    file_ext = '.' + file.filename.split('.')[-1].lower() if '.' in file.filename else ''
    if file_ext not in allowed_ext:
        raise HTTPException(status_code=400, detail=f"Unsupported format. Use: {', '.join(allowed_ext)}")

    content = await file.read()
    if len(content) > 20 * 1024 * 1024:  # 20 MB limit
        raise HTTPException(status_code=400, detail="File too large (max 20MB)")

    # Extract text content
    text_content = ""
    try:
        if file_ext == '.pdf':
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            try:
                from file_loader import load_text_from_pdf
                text_content = load_text_from_pdf(tmp_path) or ""
            finally:
                os.unlink(tmp_path)
        elif file_ext == '.docx':
            import docx
            doc_obj = docx.Document(io.BytesIO(content))
            text_content = "\n".join(p.text for p in doc_obj.paragraphs)
        elif file_ext in ('.xlsx', '.xls'):
            import pandas as pd
            df = pd.read_excel(io.BytesIO(content))
            text_content = df.to_string(index=False)
        elif file_ext == '.csv':
            import pandas as pd
            text_content = pd.read_csv(io.BytesIO(content)).to_string(index=False)
        else:
            text_content = content.decode('utf-8', errors='replace')
    except Exception as e:
        logger.warning(f"Could not extract text from traceability doc: {e}")
        text_content = ""

    # Upload to GCS
    gcs_url = None
    try:
        GCS_BUCKET_NAME = os.getenv("GCS_BUCKET_NAME", "applydi-documents")
        gcs_client = storage.Client()
        bucket = gcs_client.bucket(GCS_BUCKET_NAME)
        gcs_filename = f"traceability/{int(time.time())}_{sanitize_filename(file.filename)}"
        blob = bucket.blob(gcs_filename)
        blob.upload_from_string(content)
        gcs_url = blob.public_url
    except Exception as e:
        logger.warning(f"GCS upload failed for traceability doc: {e}")

    doc = Document(
        filename=file.filename,
        content=text_content,
        user_id=int(user_id),
        agent_id=agent_id,
        gcs_url=gcs_url,
        document_type="traceability"
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)

    return {"document": {"id": doc.id, "filename": doc.filename, "created_at": doc.created_at.isoformat()}}


@app.get("/api/agents/{agent_id}/traceability-docs")
async def list_traceability_docs(
    agent_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """List traceability documents for an agent."""
    agent = _user_can_edit_agent(int(user_id), agent_id, db)

    docs = db.query(Document).filter(
        Document.agent_id == agent_id,
        Document.document_type == "traceability"
    ).order_by(Document.created_at.desc()).all()

    return {"documents": [
        {"id": d.id, "filename": d.filename, "created_at": d.created_at.isoformat()}
        for d in docs
    ]}


@app.delete("/api/agents/{agent_id}/traceability-docs/{doc_id}")
async def delete_traceability_doc(
    agent_id: int,
    doc_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Delete a traceability document."""
    agent = _user_can_edit_agent(int(user_id), agent_id, db)

    doc = db.query(Document).filter(
        Document.id == doc_id,
        Document.agent_id == agent_id,
        Document.document_type == "traceability"
    ).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    # Delete from GCS if applicable
    if doc.gcs_url:
        try:
            GCS_BUCKET_NAME = os.getenv("GCS_BUCKET_NAME", "applydi-documents")
            gcs_client = storage.Client()
            bucket = gcs_client.bucket(GCS_BUCKET_NAME)
            blob_name = doc.gcs_url.split(f"{GCS_BUCKET_NAME}/")[-1] if GCS_BUCKET_NAME in (doc.gcs_url or "") else None
            if blob_name:
                bucket.blob(blob_name).delete()
        except Exception as e:
            logger.warning(f"Could not delete GCS blob for traceability doc {doc_id}: {e}")

    db.delete(doc)
    db.commit()

    return {"message": "Document deleted"}


## ---- Notion Links Endpoints ----

@app.post("/api/agents/{agent_id}/notion-links")
async def add_notion_link(
    agent_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db),
    body: dict = Body(...)
):
    """Link a Notion page or database to an agent."""
    agent = _user_can_edit_agent(int(user_id), agent_id, db)

    url = body.get("url", "").strip()
    resource_type = body.get("type", "page").strip()
    if resource_type not in ("page", "database"):
        raise HTTPException(status_code=400, detail="type must be 'page' or 'database'")
    if not url:
        raise HTTPException(status_code=400, detail="url is required")

    from notion_client import extract_notion_id, fetch_page_title, fetch_database_title, get_notion_token
    user = db.query(User).filter(User.id == int(user_id)).first()
    user_company_id = user.company_id if user else None
    if not get_notion_token(user_company_id):
        raise HTTPException(status_code=503, detail="Notion integration is not configured. Ask your organization owner to configure it.")

    try:
        notion_id = extract_notion_id(url)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))

    try:
        if resource_type == "page":
            label = fetch_page_title(notion_id, company_id=user_company_id)
        else:
            label = fetch_database_title(notion_id, company_id=user_company_id)
    except Exception as e:
        logger.warning(f"Could not fetch Notion title for {notion_id}: {e}")
        raise HTTPException(status_code=400, detail=f"Could not access Notion resource: {e}")

    link = NotionLink(
        agent_id=agent_id,
        notion_resource_id=notion_id,
        resource_type=resource_type,
        label=label,
    )
    db.add(link)
    db.commit()
    db.refresh(link)

    return {
        "link": {
            "id": link.id,
            "notion_resource_id": link.notion_resource_id,
            "resource_type": link.resource_type,
            "label": link.label,
            "created_at": link.created_at.isoformat(),
        }
    }


@app.get("/api/agents/{agent_id}/notion-links")
async def list_notion_links(
    agent_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """List Notion links for an agent."""
    agent = _user_can_edit_agent(int(user_id), agent_id, db)

    links = db.query(NotionLink).filter(NotionLink.agent_id == agent_id).order_by(NotionLink.created_at.desc()).all()
    # Check which links have been ingested into RAG
    ingested_ids = set()
    if links:
        link_ids = [l.id for l in links]
        ingested_docs = db.query(Document.notion_link_id).filter(Document.notion_link_id.in_(link_ids)).all()
        ingested_ids = {d[0] for d in ingested_docs}
    return {
        "links": [
            {
                "id": l.id,
                "notion_resource_id": l.notion_resource_id,
                "resource_type": l.resource_type,
                "label": l.label,
                "created_at": l.created_at.isoformat(),
                "ingested": l.id in ingested_ids,
            }
            for l in links
        ]
    }


@app.delete("/api/agents/{agent_id}/notion-links/{link_id}")
async def delete_notion_link(
    agent_id: int,
    link_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Delete a Notion link from an agent."""
    agent = _user_can_edit_agent(int(user_id), agent_id, db)

    link = db.query(NotionLink).filter(NotionLink.id == link_id, NotionLink.agent_id == agent_id).first()
    if not link:
        raise HTTPException(status_code=404, detail="Notion link not found")

    # Also delete any RAG document ingested from this Notion link
    ingested_doc = db.query(Document).filter(Document.notion_link_id == link_id).first()
    if ingested_doc:
        db.delete(ingested_doc)

    db.delete(link)
    db.commit()
    return {"message": "Notion link deleted"}


@app.post("/api/agents/{agent_id}/notion-links/{link_id}/preview")
async def preview_notion_link(
    agent_id: int,
    link_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Preview the live content of a Notion link."""
    agent = _user_can_edit_agent(int(user_id), agent_id, db)

    link = db.query(NotionLink).filter(NotionLink.id == link_id, NotionLink.agent_id == agent_id).first()
    if not link:
        raise HTTPException(status_code=404, detail="Notion link not found")

    from notion_client import fetch_page_content, fetch_database_entries, blocks_to_text, database_entries_to_text
    user = db.query(User).filter(User.id == int(user_id)).first()
    user_company_id = user.company_id if user else None

    try:
        if link.resource_type == "page":
            blocks = fetch_page_content(link.notion_resource_id, company_id=user_company_id)
            content = blocks_to_text(blocks)
        else:
            entries = fetch_database_entries(link.notion_resource_id, company_id=user_company_id)
            content = database_entries_to_text(entries)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Failed to fetch Notion content: {e}")

    return {"label": link.label, "resource_type": link.resource_type, "content": content}


@app.get("/api/agents/{agent_id}/sources")
async def get_agent_sources(
    agent_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Consolidated sources page: documents + notion links + permissions."""
    agent = _user_can_access_agent(int(user_id), agent_id, db)

    # Check edit permission
    can_edit = agent.user_id == int(user_id)
    if not can_edit:
        share = db.query(AgentShare).filter(
            AgentShare.agent_id == agent_id,
            AgentShare.user_id == int(user_id),
            AgentShare.can_edit == True
        ).first()
        can_edit = share is not None

    # RAG documents for this agent
    docs = db.query(Document).filter(
        Document.agent_id == agent_id,
        Document.document_type == "rag"
    ).order_by(Document.created_at.desc()).all()

    # Notion links
    notion_links = db.query(NotionLink).filter(NotionLink.agent_id == agent_id).order_by(NotionLink.created_at.desc()).all()

    # Set of notion_link_ids that already have a Document
    ingested_link_ids = set()
    for d in docs:
        if d.notion_link_id:
            ingested_link_ids.add(d.notion_link_id)

    return {
        "agent_name": agent.name,
        "documents": [
            {
                "id": d.id,
                "filename": d.filename,
                "created_at": d.created_at.isoformat() if d.created_at else None,
                "has_file": bool(d.gcs_url),
                "notion_link_id": d.notion_link_id,
            }
            for d in docs
        ],
        "notion_links": [
            {
                "id": nl.id,
                "label": nl.label,
                "resource_type": nl.resource_type,
                "created_at": nl.created_at.isoformat() if nl.created_at else None,
                "ingested": nl.id in ingested_link_ids,
            }
            for nl in notion_links
        ],
        "can_edit": can_edit,
    }


@app.post("/api/agents/{agent_id}/notion-links/{link_id}/ingest")
async def ingest_notion_link(
    agent_id: int,
    link_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Ingest a Notion link's content into the RAG pipeline."""
    agent = _user_can_edit_agent(int(user_id), agent_id, db)

    link = db.query(NotionLink).filter(NotionLink.id == link_id, NotionLink.agent_id == agent_id).first()
    if not link:
        raise HTTPException(status_code=404, detail="Notion link not found")

    # Check not already ingested
    existing = db.query(Document).filter(Document.notion_link_id == link_id).first()
    if existing:
        raise HTTPException(status_code=409, detail="This Notion link has already been ingested")

    from notion_client import fetch_page_content, fetch_database_entries, blocks_to_text, database_entries_to_text
    user = db.query(User).filter(User.id == int(user_id)).first()
    user_company_id = user.company_id if user else None

    try:
        if link.resource_type == "page":
            blocks = fetch_page_content(link.notion_resource_id, company_id=user_company_id)
            text_content = blocks_to_text(blocks)
        else:
            entries = fetch_database_entries(link.notion_resource_id, company_id=user_company_id)
            text_content = database_entries_to_text(entries)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Failed to fetch Notion content: {e}")

    if not text_content or not text_content.strip():
        raise HTTPException(status_code=400, detail="Notion resource returned empty content")

    from rag_engine import ingest_text_content
    safe_label = (link.label or "untitled").replace("/", "_")[:80]
    filename = f"notion_{link.resource_type}_{safe_label}.txt"

    # Upload text to GCS so the document is downloadable
    gcs_url = None
    try:
        from google.cloud import storage as gcs_storage
        import time as _time
        bucket_name = os.getenv("GCS_BUCKET_NAME", "applydi-documents")
        gcs_client = gcs_storage.Client()
        bucket = gcs_client.bucket(bucket_name)
        gcs_filename = f"{int(_time.time())}_{filename.replace(' ', '_')}"
        blob = bucket.blob(gcs_filename)
        blob.upload_from_string(text_content.encode("utf-8"), content_type="text/plain; charset=utf-8")
        gcs_url = blob.public_url
    except Exception as e:
        logger.warning(f"Could not upload Notion text to GCS: {e}")

    doc_id = ingest_text_content(
        text_content=text_content,
        filename=filename,
        user_id=int(user_id),
        agent_id=agent_id,
        db=db,
        gcs_url=gcs_url,
        notion_link_id=link.id,
    )

    doc = db.query(Document).filter(Document.id == doc_id).first()
    chunk_count = len(doc.chunks) if doc else 0

    return {
        "document_id": doc_id,
        "filename": filename,
        "chunk_count": chunk_count,
        "message": "Notion content ingested successfully",
    }


@app.post("/api/agents/{agent_id}/notion-links/{link_id}/resync")
async def resync_notion_link(
    agent_id: int,
    link_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Re-sync a Notion link: delete existing document then re-ingest from Notion."""
    agent = _user_can_edit_agent(int(user_id), agent_id, db)

    link = db.query(NotionLink).filter(NotionLink.id == link_id, NotionLink.agent_id == agent_id).first()
    if not link:
        raise HTTPException(status_code=404, detail="Notion link not found")

    # Delete existing ingested document if any
    existing = db.query(Document).filter(Document.notion_link_id == link_id).first()
    if existing:
        db.delete(existing)
        db.commit()

    # Re-fetch content from Notion
    from notion_client import fetch_page_content, fetch_database_entries, blocks_to_text, database_entries_to_text
    user = db.query(User).filter(User.id == int(user_id)).first()
    user_company_id = user.company_id if user else None

    try:
        if link.resource_type == "page":
            blocks = fetch_page_content(link.notion_resource_id, company_id=user_company_id)
            text_content = blocks_to_text(blocks)
        else:
            entries = fetch_database_entries(link.notion_resource_id, company_id=user_company_id)
            text_content = database_entries_to_text(entries)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Failed to fetch Notion content: {e}")

    if not text_content or not text_content.strip():
        raise HTTPException(status_code=400, detail="Notion resource returned empty content")

    from rag_engine import ingest_text_content
    safe_label = (link.label or "untitled").replace("/", "_")[:80]
    filename = f"notion_{link.resource_type}_{safe_label}.txt"

    gcs_url = None
    try:
        from google.cloud import storage as gcs_storage
        import time as _time
        bucket_name = os.getenv("GCS_BUCKET_NAME", "applydi-documents")
        gcs_client = gcs_storage.Client()
        bucket = gcs_client.bucket(bucket_name)
        gcs_filename = f"{int(_time.time())}_{filename.replace(' ', '_')}"
        blob = bucket.blob(gcs_filename)
        blob.upload_from_string(text_content.encode("utf-8"), content_type="text/plain; charset=utf-8")
        gcs_url = blob.public_url
    except Exception as e:
        logger.warning(f"Could not upload Notion text to GCS: {e}")

    doc_id = ingest_text_content(
        text_content=text_content,
        filename=filename,
        user_id=int(user_id),
        agent_id=agent_id,
        db=db,
        gcs_url=gcs_url,
        notion_link_id=link.id,
    )

    doc = db.query(Document).filter(Document.id == doc_id).first()
    chunk_count = len(doc.chunks) if doc else 0

    logger.info(f"Notion link {link_id} re-synced for agent {agent_id} by user {user_id}")
    return {
        "document_id": doc_id,
        "filename": filename,
        "chunk_count": chunk_count,
        "message": "Notion content re-synced successfully",
    }


## Suppression des endpoints de génération de fichiers CSV et PDF

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8080)

class ConversationCreate(BaseModel):
    agent_id: Optional[int] = None
    team_id: Optional[int] = None
    title: Optional[str] = None

    @classmethod
    def validate_one_id(cls, values):
        if not values.get('agent_id') and not values.get('team_id'):
            raise ValueError('agent_id ou team_id doit être fourni')
        return values

    from pydantic import root_validator
    @root_validator(pre=True)
    def check_ids(cls, values):
        return cls.validate_one_id(values)

class MessageCreate(BaseModel):
    conversation_id: int
    role: str
    content: str


def _verify_conversation_owner(conversation_id: int, user_id: str, db: Session) -> Conversation:
    """Load a conversation and verify the authenticated user owns it (via agent, share, or team)."""
    conv = db.query(Conversation).filter(Conversation.id == conversation_id).first()
    if not conv:
        raise HTTPException(status_code=404, detail="Conversation not found")
    uid = int(user_id)
    # If the conversation has a user_id, check it matches
    if conv.user_id and conv.user_id == uid:
        return conv
    if conv.agent_id:
        agent = db.query(Agent).filter(Agent.id == conv.agent_id).first()
        if not agent:
            raise HTTPException(status_code=404, detail="Conversation not found")
        if agent.user_id == uid:
            return conv
        # Check if user has a share on this agent
        share = db.query(AgentShare).filter(
            AgentShare.agent_id == conv.agent_id,
            AgentShare.user_id == uid
        ).first()
        if share:
            return conv
        raise HTTPException(status_code=404, detail="Conversation not found")
    elif conv.team_id:
        team = db.query(Team).filter(Team.id == conv.team_id).first()
        if not team or team.user_id != uid:
            raise HTTPException(status_code=404, detail="Conversation not found")
    else:
        raise HTTPException(status_code=404, detail="Conversation not found")
    return conv


@app.post("/conversations", response_model=dict)
async def create_conversation(conv: ConversationCreate, db: Session = Depends(get_db), user_id: str = Depends(verify_token)):
    uid = int(user_id)
    if conv.agent_id:
        _user_can_access_agent(uid, conv.agent_id, db)
    if conv.team_id:
        team = db.query(Team).filter(Team.id == conv.team_id).first()
        if not team or team.user_id != uid:
            raise HTTPException(status_code=404, detail="Team not found")
    conversation = Conversation(agent_id=conv.agent_id, team_id=conv.team_id, title=conv.title, user_id=uid)
    db.add(conversation)
    db.commit()
    db.refresh(conversation)
    return {"conversation_id": conversation.id}


from fastapi import Query

@app.get("/conversations", response_model=List[dict])
async def list_conversations(
    agent_id: int = Query(None),
    team_id: int = Query(None),
    db: Session = Depends(get_db),
    user_id: str = Depends(verify_token)
):
    uid = int(user_id)
    if agent_id is not None:
        agent = _user_can_access_agent(uid, agent_id, db)
        from sqlalchemy import or_
        # Show user's own conversations + legacy conversations (no user_id) if user is owner
        user_filter = [Conversation.user_id == uid]
        if agent.user_id == uid:
            user_filter.append(Conversation.user_id.is_(None))
        conversations = db.query(Conversation).filter(
            Conversation.agent_id == agent_id,
            or_(*user_filter)
        ).order_by(Conversation.created_at.desc()).all()
    elif team_id is not None:
        team = db.query(Team).filter(Team.id == team_id).first()
        if not team or team.user_id != uid:
            raise HTTPException(status_code=404, detail="Team not found")
        conversations = db.query(Conversation).filter(Conversation.team_id == team_id).order_by(Conversation.created_at.desc()).all()
    else:
        raise HTTPException(status_code=422, detail="agent_id ou team_id doit être fourni")
    return [{"id": c.id, "title": c.title, "created_at": c.created_at} for c in conversations]

@app.post("/conversations/{conversation_id}/messages", response_model=dict)
async def add_message(conversation_id: int, msg: MessageCreate, db: Session = Depends(get_db), user_id: str = Depends(verify_token)):
    _verify_conversation_owner(conversation_id, user_id, db)
    message = Message(conversation_id=conversation_id, role=msg.role, content=msg.content)
    db.add(message)
    db.commit()
    db.refresh(message)
    return {"message_id": message.id}

@app.get("/conversations/{conversation_id}/messages", response_model=List[dict])
async def get_messages(conversation_id: int, db: Session = Depends(get_db), user_id: str = Depends(verify_token)):
    _verify_conversation_owner(conversation_id, user_id, db)
    messages = db.query(Message).filter(Message.conversation_id == conversation_id).order_by(Message.timestamp.asc()).all()
    return [{"id": m.id, "role": m.role, "content": m.content, "timestamp": m.timestamp} for m in messages]
# Endpoint de connexion agent (email + password)
class FeedbackRequest(BaseModel):
    feedback: str  # 'like' ou 'dislike'

@app.patch("/messages/{message_id}/feedback")
async def set_message_feedback(message_id: int, req: FeedbackRequest, db: Session = Depends(get_db), user_id: str = Depends(verify_token)):
    msg = db.query(Message).filter(Message.id == message_id).first()
    if not msg:
        raise HTTPException(status_code=404, detail="Message not found")
    # Verify ownership via the message's conversation
    _verify_conversation_owner(msg.conversation_id, user_id, db)
    if req.feedback not in ("like", "dislike"):
        raise HTTPException(status_code=400, detail="Feedback must be 'like' or 'dislike'")
    msg.feedback = req.feedback
    # Si feedback = like, bufferise le message et le message user précédent
    if req.feedback == "like":
        msg.buffered = 1
        # Cherche le message user juste avant dans la même conversation
        prev_user_msg = db.query(Message).filter(
            Message.conversation_id == msg.conversation_id,
            Message.timestamp < msg.timestamp,
            Message.role == "user"
        ).order_by(Message.timestamp.desc()).first()
        if prev_user_msg:
            prev_user_msg.feedback = "like"
            prev_user_msg.buffered = 1
    db.commit()
    return {"message_id": msg.id, "feedback": msg.feedback, "buffered": msg.buffered}
# --- Endpoints pour conversations et messages ---
# Use validated model from validation.py
# class ConversationTitleUpdate is now ConversationTitleValidated

@app.put("/conversations/{conversation_id}/title")
async def update_conversation_title(conversation_id: int, data: ConversationTitleValidated, db: Session = Depends(get_db), user_id: str = Depends(verify_token)):
    conv = _verify_conversation_owner(conversation_id, user_id, db)
    conv.title = data.title
    db.commit()
    db.refresh(conv)
    return {"id": conv.id, "title": conv.title}

# Endpoint pour supprimer une conversation
@app.delete("/conversations/{conversation_id}")
async def delete_conversation(conversation_id: int, db: Session = Depends(get_db), user_id: str = Depends(verify_token)):
    conv = _verify_conversation_owner(conversation_id, user_id, db)
    db.delete(conv)
    db.commit()
    return {"message": "Conversation deleted"}

# --- SLACK CONFIG ENDPOINTS ---

@app.get("/api/agents/{agent_id}/slack-config")
async def get_slack_config(agent_id: int, request: Request, db: Session = Depends(get_db)):
    """Get Slack configuration status for an agent. Never returns tokens in clear."""
    from auth import verify_token_from_cookie
    user_id = verify_token_from_cookie(request)
    agent = db.query(Agent).filter(Agent.id == agent_id).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")
    if str(agent.user_id) != str(user_id):
        raise HTTPException(status_code=403, detail="Not authorized")

    raw_token = agent.slack_bot_token
    raw_secret = agent.slack_signing_secret
    is_configured = bool(raw_token and raw_secret)

    masked_token = ""
    masked_secret = ""
    if raw_token and len(raw_token) > 8:
        masked_token = raw_token[:4] + "****" + raw_token[-4:]
    elif raw_token:
        masked_token = "****"
    if raw_secret and len(raw_secret) > 8:
        masked_secret = raw_secret[:4] + "****" + raw_secret[-4:]
    elif raw_secret:
        masked_secret = "****"

    return {
        "is_configured": is_configured,
        "team_id": agent.slack_team_id or "",
        "bot_user_id": agent.slack_bot_user_id or "",
        "masked_token": masked_token,
        "masked_secret": masked_secret,
    }


class SlackConfigRequest(BaseModel):
    slack_bot_token: str
    slack_signing_secret: str


@app.put("/api/agents/{agent_id}/slack-config")
async def update_slack_config(agent_id: int, payload: SlackConfigRequest, request: Request, db: Session = Depends(get_db)):
    """Save Slack credentials after validating the bot token via auth.test."""
    from auth import verify_token_from_cookie
    user_id = verify_token_from_cookie(request)
    agent = db.query(Agent).filter(Agent.id == agent_id).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")
    if str(agent.user_id) != str(user_id):
        raise HTTPException(status_code=403, detail="Not authorized")

    # Validate the token via Slack auth.test
    resp = requests.post(
        "https://slack.com/api/auth.test",
        headers={"Authorization": f"Bearer {payload.slack_bot_token}"},
    )
    slack_data = resp.json()
    if not slack_data.get("ok"):
        raise HTTPException(status_code=400, detail=f"Invalid Slack token: {slack_data.get('error', 'unknown error')}")

    team_id = slack_data.get("team_id", "")
    bot_user_id = slack_data.get("user_id", "")
    team_name = slack_data.get("team", "")

    # Store via property setters (auto-encryption)
    agent.slack_bot_token = payload.slack_bot_token
    agent.slack_signing_secret = payload.slack_signing_secret
    agent.slack_team_id = team_id
    agent.slack_bot_user_id = bot_user_id
    db.commit()

    logger.info(f"Slack config saved for agent {agent_id}: team={team_name} ({team_id}), bot_user={bot_user_id}")

    return {
        "ok": True,
        "team_id": team_id,
        "bot_user_id": bot_user_id,
        "team_name": team_name,
    }


@app.delete("/api/agents/{agent_id}/slack-config")
async def delete_slack_config(agent_id: int, request: Request, db: Session = Depends(get_db)):
    """Remove all Slack configuration from an agent."""
    from auth import verify_token_from_cookie
    user_id = verify_token_from_cookie(request)
    agent = db.query(Agent).filter(Agent.id == agent_id).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")
    if str(agent.user_id) != str(user_id):
        raise HTTPException(status_code=403, detail="Not authorized")

    agent._slack_bot_token = None
    agent._slack_signing_secret = None
    agent.slack_team_id = None
    agent.slack_bot_user_id = None
    db.commit()

    logger.info(f"Slack config removed for agent {agent_id}")
    return {"ok": True}


@app.post("/api/agents/{agent_id}/slack-test")
async def test_slack_connection(agent_id: int, request: Request, db: Session = Depends(get_db)):
    """Test Slack connection using the stored bot token."""
    from auth import verify_token_from_cookie
    user_id = verify_token_from_cookie(request)
    agent = db.query(Agent).filter(Agent.id == agent_id).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")
    if str(agent.user_id) != str(user_id):
        raise HTTPException(status_code=403, detail="Not authorized")

    slack_token = agent.slack_bot_token
    if not slack_token:
        raise HTTPException(status_code=400, detail="No Slack token configured")

    resp = requests.post(
        "https://slack.com/api/auth.test",
        headers={"Authorization": f"Bearer {slack_token}"},
    )
    slack_data = resp.json()

    if not slack_data.get("ok"):
        return {"is_valid": False, "error": slack_data.get("error", "unknown error")}

    return {
        "is_valid": True,
        "team_name": slack_data.get("team", ""),
        "bot_name": slack_data.get("user", ""),
        "team_id": slack_data.get("team_id", ""),
    }


# --- SLACK WEBHOOK ENDPOINT ---


# On garde les 500 derniers event_id pour éviter les doublons
_recent_event_ids = deque(maxlen=500)
_event_ids_lock = threading.Lock()


def verify_slack_signature(request_body: bytes, timestamp: str, signature: str, signing_secret: str) -> bool:
    """
    Verify Slack request signature to prevent fake events.

    Security: Each Slack webhook request is signed with a unique signature based on:
    - Request body
    - Timestamp
    - Signing secret (unique per agent/workspace)

    This prevents:
    - Fake events from attackers
    - Replay attacks (timestamp validation)
    - Timing attacks (constant-time comparison)

    Multi-tenant: Each agent has their own signing_secret stored in DB.
    """
    if not signing_secret:
        logger.error("Slack signing_secret not configured for this agent - REJECTING request for security")
        return False  # Strict mode: reject if not configured

    # Check timestamp to prevent replay attacks (max 5 minutes old)
    try:
        request_timestamp = int(timestamp)
        current_timestamp = int(time.time())
        if abs(current_timestamp - request_timestamp) > 60 * 5:
            logger.warning(f"Slack request timestamp too old: {timestamp}")
            return False
    except (ValueError, TypeError):
        logger.error(f"Invalid Slack timestamp: {timestamp}")
        return False

    # Compute expected signature using HMAC-SHA256
    sig_basestring = f"v0:{timestamp}:".encode() + request_body
    expected_signature = "v0=" + hmac.new(
        signing_secret.encode(),
        sig_basestring,
        hashlib.sha256
    ).hexdigest()

    # Constant-time comparison to prevent timing attacks
    if not hmac.compare_digest(expected_signature, signature):
        logger.error("Slack signature mismatch — request rejected")
        return False

    return True


@app.post("/slack/events")
async def slack_events(request: Request, db: Session = Depends(get_db)):
    # Get raw body and headers for signature verification
    timestamp = request.headers.get("X-Slack-Request-Timestamp", "")
    signature = request.headers.get("X-Slack-Signature", "")
    body = await request.body()

    # Parse JSON data
    data = json.loads(body.decode('utf-8'))

    # Vérification du challenge lors de l'installation (no signature needed for this)
    if data.get("type") == "url_verification":
        return {"challenge": data["challenge"]}

    # Extract team_id and event data for agent lookup
    event = data.get("event", {})
    team_id = data.get("team_id") or event.get("team")
    event_text = event.get("text", "")

    # Parse bot mentions from text (format: <@U123ABC>)
    mentioned_bot_user_ids = re.findall(r"<@([A-Z0-9]+)>", event_text)

    # Find agent for signature verification
    # Priority 1: Match by bot_user_id (most specific, handles multi-bot workspaces)
    agent_for_verification = None
    matched_bot_user_id = None

    if mentioned_bot_user_ids:
        for bot_user_id in mentioned_bot_user_ids:
            agent = db.query(Agent).filter(
                Agent.slack_bot_user_id == bot_user_id,
                Agent.slack_team_id == team_id
            ).first()
            if agent:
                agent_for_verification = agent
                matched_bot_user_id = bot_user_id
                logger.info(f"Found agent by bot_user_id: {bot_user_id} (team: {team_id}) -> {agent.name}")
                break

    # Priority 2: Fallback to team_id only (less specific, for backwards compatibility)
    if not agent_for_verification:
        agent_for_verification = db.query(Agent).filter(
            Agent.slack_team_id == team_id
        ).first()
        if agent_for_verification:
            logger.info(f"Found agent by team_id only: {team_id} -> {agent_for_verification.name}")

    if not agent_for_verification:
        logger.warning(f"No agent found for Slack team_id: {team_id}")
        raise HTTPException(status_code=403, detail="No agent configured for this Slack workspace")

    # SECURITY: Verify Slack signature with this agent's signing_secret
    if not verify_slack_signature(body, timestamp, signature, agent_for_verification.slack_signing_secret):
        logger.error(f"Slack signature verification failed for team_id: {team_id}, agent: {agent_for_verification.name}")
        raise HTTPException(status_code=403, detail="Invalid Slack signature")

    logger.info(f"Slack signature verified for team: {team_id}, agent: {agent_for_verification.name}")

    # Check for duplicate events
    event_id = data.get("event_id")
    if event_id:
        with _event_ids_lock:
            if event_id in _recent_event_ids:
                logger.info(f"Event déjà traité, on ignore: {event_id}")
                return {"ok": True, "info": "Duplicate event ignored"}
            _recent_event_ids.append(event_id)

    # On ne traite que les mentions du bot (app_mention)
    if event.get("type") == "app_mention" and "text" in event:
        user_message = event["text"]
        channel = event["channel"]
        thread_ts = event.get("thread_ts")  # timestamp du thread si présent

        # Use the agent we already found for verification
        agent = agent_for_verification
        agent_id = agent.id
        slack_token = agent.slack_bot_token

        if not slack_token:
            logger.warning(f"No Slack token found for agent with team_id={team_id}")
            return {"ok": False, "error": "No Slack token for agent"}
        # 1. Récupère l'historique du channel ou du thread
        history = []
        try:
            headers = {"Authorization": f"Bearer {slack_token}"}
            messages = []
            if thread_ts:
                # Récupère tous les messages du thread
                resp = requests.get(
                    "https://slack.com/api/conversations.replies",
                    headers=headers,
                    params={"channel": channel, "ts": thread_ts}
                )
                messages = resp.json().get("messages", [])
                logger.info(f"Slack thread history (count={len(messages)}): {[m.get('text','') for m in messages]}")
            else:
                # Récupère les derniers messages du channel
                resp = requests.get(
                    "https://slack.com/api/conversations.history",
                    headers=headers,
                    params={"channel": channel, "limit": 10}
                )
                messages = resp.json().get("messages", [])
                logger.info(f"Slack channel history (count={len(messages)}): {[m.get('text','') for m in messages]}")
            # Formate l'historique pour le modèle dans l'ordre du plus ancien au plus récent
            for msg in sorted(messages, key=lambda m: float(m.get("ts", 0))):
                role = "user" if msg.get("user") else "assistant"
                content = msg.get("text", "")
                history.append({"role": role, "content": content})
        except Exception as e:
            logger.error(f"Erreur récupération historique Slack: {e}")
            history = []
        # Log le contenu de l'historique avant get_answer
        logger.info(f"Slack context sent to get_answer: {history}")
        # 2. Resolve model_id from agent's type
        slack_model_id = None
        if agent.finetuned_model_id:
            slack_model_id = agent.finetuned_model_id
        else:
            atype = getattr(agent, 'type', 'conversationnel')
            if atype == 'recherche_live':
                slack_model_id = os.getenv('PERPLEXITY_MODEL', 'perplexity:sonar')
            else:
                slack_model_id = os.getenv('MISTRAL_MODEL', 'mistral:mistral-small-latest')
        # Appel direct à la fonction get_answer avec l'historique Slack
        answer = get_answer(user_message, None, db, agent_id=agent_id, history=history, model_id=slack_model_id)
        # 3. Envoie la réponse sur Slack avec le bon token
        resp = requests.post(
            "https://slack.com/api/chat.postMessage",
            headers={"Authorization": f"Bearer {slack_token}"},
            json={"channel": channel, "text": answer, "thread_ts": thread_ts} if thread_ts else {"channel": channel, "text": answer}
        )
        logger.info(f"Slack API response: status={resp.status_code}")
    return {"ok": True}



class ForgotPasswordRequest(BaseModel):
    email: str

class ResetPasswordRequest(BaseModel):
    token: str
    new_password: str


# Endpoint pour demander la réinitialisation du mot de passe (DB version)
@app.post("/forgot-password")
async def forgot_password(req: ForgotPasswordRequest, request: Request, db: Session = Depends(get_db)):
    # Rate limiting: prevent abuse of password reset
    ip = request.client.host if hasattr(request, 'client') and request.client else 'unknown'
    if not _check_auth_rate_limit(ip):
        logger.warning(f"Rate limit exceeded for forgot-password from IP: {ip}")
        raise HTTPException(status_code=429, detail="Too many password reset attempts. Please try again in 15 minutes.")

    user = db.query(User).filter(User.email == req.email).first()
    if not user:
        raise HTTPException(status_code=404, detail="Utilisateur introuvable")

    # Generate token and hash it before storing (security: prevent token theft if DB compromised)
    token = str(uuid4())
    token_hash = hash_reset_token(token)
    expires_at = datetime.utcnow() + timedelta(minutes=15)

    # Store HASHED token in DB
    reset_token = PasswordResetToken(
        user_id=user.id,
        token=token_hash,  # Store hash, not plaintext
        expires_at=expires_at,
        used=False
    )
    db.add(reset_token)
    db.commit()

    # Send PLAINTEXT token in email (user needs it to reset password)
    frontend_url = os.getenv("FRONTEND_URL", "https://taic.ai")
    reset_link = f"{frontend_url}/reset-password?token={token}"
    try:
        send_password_reset_email(user.email, reset_link)
        logger.info(f"Password reset email sent to {user.email}")
        return {"message": "Un lien de réinitialisation a été envoyé par email"}
    except Exception as e:
        logger.error(f"Erreur lors de l'envoi de l'email: {e}")
        return {"message": "Erreur lors de l'envoi de l'email"}

# Endpoint pour réinitialiser le mot de passe (DB version)
@app.post("/reset-password")
async def reset_password(req: ResetPasswordRequest, request: Request, db: Session = Depends(get_db)):
    # Rate limiting: prevent brute force token guessing
    ip = request.client.host if hasattr(request, 'client') and request.client else 'unknown'
    if not _check_auth_rate_limit(ip):
        logger.warning(f"Rate limit exceeded for reset-password from IP: {ip}")
        raise HTTPException(status_code=429, detail="Too many reset attempts. Please try again in 15 minutes.")

    # Hash the received token to look it up in DB (tokens are stored hashed)
    token_hash = hash_reset_token(req.token)
    reset_token = db.query(PasswordResetToken).filter(PasswordResetToken.token == token_hash).first()

    if not reset_token:
        raise HTTPException(status_code=400, detail="Token invalide ou expiré")
    if reset_token.expires_at < datetime.utcnow():
        raise HTTPException(status_code=400, detail="Token expiré")
    if reset_token.used:
        raise HTTPException(status_code=400, detail="Token déjà utilisé")

    user = db.query(User).filter(User.id == reset_token.user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="Utilisateur introuvable")

    user.hashed_password = hash_password(req.new_password)
    reset_token.used = True
    db.commit()

    logger.info(f"Password reset successful for user {user.email}")
    return {"message": "Mot de passe réinitialisé avec succès"}
# Mémoire temporaire pour les event_id déjà traités (reset à chaque redémarrage du serveur)


@app.get("/health-nltk")
async def health_nltk():
    """Health check for NLTK and chunking logic"""
    from file_loader import chunk_text
    test_text = "Hello world. This is a test.\n\nNew paragraph."
    try:
        chunks = chunk_text(test_text)
        return {"status": "ok", "chunks": chunks, "n_chunks": len(chunks)}
    except Exception as e:
        logger.error(f"Health check NLTK failed: {e}")
        return {"status": "error", "detail": "NLTK health check failed"}

#test
##### Public agents endpoints (no auth) #####

# Public chat rate limiting configuration
_PUBLIC_CHAT_LIMIT = 60  # messages per hour per IP
_PUBLIC_CHAT_WINDOW = 3600  # 1 hour in seconds


def _check_rate_limit(ip: str):
    """
    Check if IP has exceeded rate limit for public chat using Redis.
    Returns True if allowed, False if rate limited.

    Security: Prevents abuse of public chat endpoints via distributed rate limiting.
    Redis ensures limits work across multiple Cloud Run instances.
    """
    key = f"rate_limit:public_chat:{ip}"

    # Try Redis first
    if redis_client:
        try:
            current = redis_client.get(key)
            current = int(current) if current else 0

            if current >= _PUBLIC_CHAT_LIMIT:
                logger.warning(f"Public chat rate limit exceeded for IP: {ip}")
                return False

            # Increment and set expiration
            pipe = redis_client.pipeline()
            pipe.incr(key)
            pipe.expire(key, _PUBLIC_CHAT_WINDOW)
            pipe.execute()
            return True

        except Exception as e:
            logger.error(f"Redis rate limit check failed: {e}. Using fallback.")

    # Fallback to in-memory (not distributed, but better than nothing)
    now = time.time()
    q = _public_chat_rate_fallback.get(ip, [])
    q = [t for t in q if now - t < _PUBLIC_CHAT_WINDOW]

    if len(q) >= _PUBLIC_CHAT_LIMIT:
        logger.warning(f"Public chat rate limit exceeded for IP (fallback): {ip}")
        return False

    q.append(now)
    _public_chat_rate_fallback[ip] = q
    return True


@app.get("/public/agents/{agent_id}")
async def public_get_agent(agent_id: int, request: Request, db: Session = Depends(get_db)):
    """Return public agent profile if statut == 'public'"""
    ip = request.client.host if hasattr(request, 'client') and request.client else 'unknown'
    if not _check_rate_limit(ip):
        raise HTTPException(status_code=429, detail="Too many requests. Please try again later.")
    agent = db.query(Agent).filter(Agent.id == agent_id, Agent.statut == 'public').first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found or not public")
    # Only expose non-sensitive fields
    return {
        "id": agent.id,
        "name": agent.name,
        "contexte": agent.contexte,
        "biographie": agent.biographie,
        "profile_photo": agent.profile_photo,
        "created_at": agent.created_at.isoformat() if hasattr(agent, 'created_at') else None,
        "slug": getattr(agent, 'slug', None),
    }


class PublicChatRequest(BaseModel):
    message: str
    history: Optional[List[dict]] = None


@app.post("/public/agents/{agent_id}/chat")
async def public_agent_chat(agent_id: int, req: PublicChatRequest, request: Request, db: Session = Depends(get_db)):
    """Public chat endpoint for a public agent. Rate-limited by IP."""
    agent = db.query(Agent).filter(Agent.id == agent_id, Agent.statut == 'public').first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found or not public")

    # Rate limiting (distributed via Redis)
    ip = request.client.host if hasattr(request, 'client') and request.client else 'unknown'
    if not _check_rate_limit(ip):
        raise HTTPException(status_code=429, detail="Rate limit exceeded. Please try again in 1 hour.")

    # Build history for the model if provided
    history = req.history or []
    # Append the current user message as last user message in history
    history.append({"role": "user", "content": req.message})

    # Resolve model_id from agent type (mistral for conversationnel, perplexity for recherche_live)
    public_model_id = agent.finetuned_model_id or resolve_model_id(agent)

    try:
        answer = get_answer(req.message, None, db, agent_id=agent_id, history=history, model_id=public_model_id)
    except Exception as e:
        logger.exception(f"Error generating public chat answer for agent {agent_id}: {e}")
        raise HTTPException(status_code=500, detail="Error generating answer")

    return {"answer": answer}

#test
@app.get("/documents/{document_id}/download-url")
async def get_signed_download_url(document_id: int, user_id: str = Depends(verify_token), db: Session = Depends(get_db)):
    """Retourne une URL signée pour télécharger le document depuis GCS"""
    from urllib.parse import urlparse
    logger = logging.getLogger("main.download_url")

    try:
        document = db.query(Document).filter(Document.id == document_id).first()
        if not document or not document.gcs_url:
            raise HTTPException(status_code=404, detail="Document non trouvé ou pas de fichier GCS")
        # Owner or has access to the agent
        if document.user_id != int(user_id):
            if document.agent_id:
                _user_can_access_agent(int(user_id), document.agent_id, db)
            else:
                raise HTTPException(status_code=403, detail="Access denied")

        gcs_url = document.gcs_url
        logger.info(f"Generating signed URL for document {document_id}, gcs_url={gcs_url}")

        # Parse bucket and blob name (supports storage.googleapis.com and gs:// formats)
        from urllib.parse import unquote
        parsed = urlparse(gcs_url)
        if gcs_url.startswith('gs://'):
            parts = gcs_url[5:].split('/', 1)
            bucket_name = parts[0]
            blob_name = parts[1] if len(parts) > 1 else ''
        else:
            path = parsed.path.lstrip('/')
            path_parts = path.split('/')
            bucket_name = path_parts[0]
            blob_name_encoded = '/'.join(path_parts[1:])
            # URL-decode the blob name (handles %C3%A9, %2B, etc.)
            blob_name = unquote(blob_name_encoded)

        logger.info(f"Blob name (encoded)={locals().get('blob_name_encoded', None)}, decoded={blob_name}")

        logger.info(f"Parsed bucket={bucket_name}, blob={blob_name}")

        client = storage.Client()
        bucket = client.bucket(bucket_name)
        blob = bucket.blob(blob_name)

        # Existence check
        try:
            exists = blob.exists()
        except Exception as e:
            logger.exception("Error checking blob existence (possible permission issue)")
            raise HTTPException(status_code=500, detail="Erreur lors de la vérification de l'existence du fichier GCS (vérifiez les permissions du service account)")

        if not exists:
            logger.error(f"Blob not found: {bucket_name}/{blob_name}")
            raise HTTPException(status_code=404, detail="Fichier introuvable dans le bucket GCS")

        try:
            url = blob.generate_signed_url(version="v4", expiration=600, method="GET")
        except Exception as e:
            logger.exception("Error generating signed URL (permission or signing issue)")
            # Provide a helpful hint without exposing sensitive info
            detail_msg = (
                "Impossible de générer le lien signé. Vérifiez que le service account a les droits GCS "
                "et la capacité de signer des URL (roles/storage.objectViewer et permissions de signature)."
            )
            # Fallback: offer a proxied download endpoint (secure, authenticated)
            proxy_url = f"/documents/{document_id}/download"
            logger.info(f"Falling back to proxy download for document {document_id}")
            return {"proxy_url": proxy_url, "note": "Signed URL generation failed; using authenticated proxy download."}

        logger.info(f"Signed URL generated for document {document_id}")
        return {"signed_url": url}

    except HTTPException:
        raise
    except Exception as e:
        tb = traceback.format_exc()
        logger = logging.getLogger("main.download_url")
        logger.error(f"Unexpected error generating signed URL for document {document_id}: {e}\n{tb}")
        raise HTTPException(status_code=500, detail="Erreur interne lors de la génération du lien de téléchargement. Vérifiez les logs du backend.")


@app.get("/documents/{document_id}/download")
async def proxy_download_document(document_id: int, user_id: str = Depends(verify_token), db: Session = Depends(get_db)):
    """Stream the object from GCS through the backend as an authenticated proxy.
    This is a secure fallback when signed URL generation is not possible from the environment.
    """
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document or not document.gcs_url:
        raise HTTPException(status_code=404, detail="Document non trouvé ou pas de fichier GCS")
    # Owner or has access to the agent
    if document.user_id != int(user_id):
        if document.agent_id:
            _user_can_access_agent(int(user_id), document.agent_id, db)
        else:
            raise HTTPException(status_code=403, detail="Access denied")

    from urllib.parse import urlparse, unquote
    gcs_url = document.gcs_url
    parsed = urlparse(gcs_url)
    path = parsed.path.lstrip('/')
    path_parts = path.split('/')
    bucket_name = path_parts[0]
    blob_name_encoded = '/'.join(path_parts[1:])
    blob_name = unquote(blob_name_encoded)
    logger = logging.getLogger("main.download_url")
    logger.info(f"Proxy download: bucket={bucket_name}, blob_encoded={blob_name_encoded}, blob_decoded={blob_name}")

    client = storage.Client()
    bucket = client.bucket(bucket_name)

    # Attempt to get blob and check existence
    def get_blob(name: str):
        return bucket.blob(name)

    blob = get_blob(blob_name)

    # Existence check (may raise if permission issues)
    exists = None
    try:
        exists = blob.exists()
    except Exception:
        logger.exception(f"Error checking existence for blob {bucket_name}/{blob_name} (possible permission issue)")
        # keep exists as None and try download below

    if exists is False:
        logger.error(f"Blob not found for proxy: {bucket_name}/{blob_name}")
        raise HTTPException(status_code=404, detail="Fichier introuvable dans le bucket GCS")

    data = None
    # Try direct download
    try:
        data = blob.download_as_bytes()
    except Exception as exc:
        logger.exception(f"Initial download attempt failed for {bucket_name}/{blob_name}: {exc}")

        # If download failed, try unicode normalization variants (NFC/NFD)
        try:
            import unicodedata
            tried = []
            for norm in ("NFC", "NFD"):
                alt_name = unicodedata.normalize(norm, blob_name)
                if alt_name in tried or alt_name == blob_name:
                    continue
                tried.append(alt_name)
                logger.info(f"Retrying download with normalized blob name ({norm}): {alt_name}")
                alt_blob = get_blob(alt_name)
                try:
                    data = alt_blob.download_as_bytes()
                    # successful: use this blob
                    blob = alt_blob
                    blob_name = alt_name
                    logger.info(f"Download succeeded with normalized name ({norm})")
                    break
                except Exception as exc2:
                    logger.exception(f"Download failed with normalized name {alt_name}: {exc2}")
        except Exception as norm_exc:
            logger.exception(f"Error during unicode normalization retries: {norm_exc}")

    if data is None:
        # Determine if likely permission issue vs not found
        # If exists is None, we couldn't determine existence due to permission; respond with 403 hint
        if exists is None:
            logger.error(f"Download failed and existence unknown for {bucket_name}/{blob_name}. Likely permission issue.")
            raise HTTPException(status_code=403, detail="Le service n'a pas les permissions nécessaires pour lire l'objet GCS. Vérifiez roles/storage.objectViewer.")
        else:
            logger.error(f"All attempts to download blob failed for {bucket_name}/{blob_name}")
            raise HTTPException(status_code=500, detail="Impossible de récupérer le fichier depuis GCS")

    # Guess mimetype
    mime, _ = mimetypes.guess_type(document.filename)
    if not mime:
        mime = 'application/octet-stream'

    from fastapi.responses import StreamingResponse
    from io import BytesIO
    # Ensure filename is safe; use the stored document filename
    safe_filename = document.filename or os.path.basename(blob_name)
    headers = {
        'Content-Disposition': f'attachment; filename="{safe_filename}"'
    }
    return StreamingResponse(BytesIO(data), media_type=mime, headers=headers)

# Endpoint pour extraire le texte des fichiers uploadés (PDF, TXT, DOCX, XLSX, PPTX, etc.)

# Nouvelle version : accepte un seul fichier UploadFile
@app.post("/api/agent/extractText")
async def extract_text_from_file(
    file: UploadFile = File(...),
    user_id: str = Depends(verify_token)
):
    """Extrait le texte d'un fichier uploadé et renvoie le texte extrait.

    Security: Requires authentication to prevent unauthorized file processing.
    """
    if not _check_api_rate_limit(user_id, "extract", _API_EXTRACT_LIMIT):
        raise HTTPException(status_code=429, detail="Too many extraction requests. Please try again later.")
    logger = logging.getLogger("extractText")
    ext = file.filename.lower().split('.')[-1]
    logger.info(f"Appel reçu sur /api/agent/extractText : filename={file.filename}, ext={ext}, content_type={getattr(file, 'content_type', 'unknown')}")
    text = ""
    try:
        if ext == "pdf":
            MAX_PDF_PAGES = 500
            logger.info(f"Tentative extraction PDF: {file.filename}")
            try:
                import pdfplumber
                with pdfplumber.open(file.file) as pdf:
                    if len(pdf.pages) > MAX_PDF_PAGES:
                        raise HTTPException(status_code=400, detail=f"PDF too large ({len(pdf.pages)} pages, max {MAX_PDF_PAGES})")
                    for i, page in enumerate(pdf.pages):
                        page_text = page.extract_text()
                        logger.info(f"Page {i+1} PDF: longueur={len(page_text) if page_text else 0}, aperçu='{page_text[:100] if page_text else ''}'")
                        if page_text:
                            text += page_text + "\n"
            except Exception as e:
                logger.error(f"Erreur extraction PDF {file.filename}: {e}")
        elif ext in ["txt", "md", "json", "xml", "csv"]:
            logger.info(f"Tentative extraction texte: {file.filename}")
            try:
                raw = await file.read()
                text = raw.decode(errors="ignore")
                logger.info(f"Texte extrait: longueur={len(text)}, aperçu='{text[:200]}'")
            except Exception as e:
                logger.error(f"Erreur extraction texte {file.filename}: {e}")
        elif ext in ["doc", "docx"]:
            logger.info(f"Tentative extraction DOCX: {file.filename}")
            try:
                doc = DocxDocument(file.file)
                text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                logger.info(f"Texte DOCX extrait: longueur={len(text)}, aperçu='{text[:200]}'")
            except Exception as e:
                logger.error(f"Erreur extraction DOCX {file.filename}: {e}")
        elif ext in ["xls", "xlsx"]:
            logger.info(f"Tentative extraction XLSX: {file.filename}")
            try:
                wb = openpyxl.load_workbook(file.file, read_only=True)
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                        text += row_text + "\n"
                logger.info(f"Texte XLSX extrait: longueur={len(text)}, aperçu='{text[:200]}'")
            except Exception as e:
                logger.error(f"Erreur extraction XLSX {file.filename}: {e}")
        elif ext in ["ppt", "pptx"]:
            logger.info(f"Tentative extraction PPTX: {file.filename}")
            try:
                pres = Presentation(file.file)
                for slide in pres.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                logger.info(f"Texte PPTX extrait: longueur={len(text)}, aperçu='{text[:200]}'")
            except Exception as e:
                logger.error(f"Erreur extraction PPTX {file.filename}: {e}")
        else:
            logger.warning(f"Type de fichier non supporté: {file.filename}")
            text = f"[Type de fichier non supporté: {file.filename}]"
    except Exception as e:
        logger.error(f"Erreur extraction {file.filename}: {e}")
        text = f"[Erreur extraction {file.filename}: {e}]"
    logger.info(f"Résultat extraction {file.filename}: longueur={len(text.strip())}, aperçu='{text.strip()[:200]}'")
    return {"text": text.strip()}


def update_agent_embedding(agent, db):
    if agent.contexte:
        agent.embedding = json.dumps(get_embedding(agent.contexte))
        db.commit()

# Security: /debug/test-openai-embeddings endpoint removed (exposed API connectivity)


# ============================================================================
# EMAIL INGESTION API - Cloud Function Integration
# ============================================================================

# Pydantic models for email ingestion
class EmailTag(BaseModel):
    name: str
    category: Optional[str] = None
    confidence: Optional[float] = None
    source: Optional[str] = None

class EmailMetadata(BaseModel):
    from_email: Optional[str] = None
    to: Optional[str] = None
    date: Optional[str] = None
    thread_id: Optional[str] = None
    has_attachments: Optional[bool] = False
    attachment_count: Optional[int] = 0
    priority: Optional[int] = None

    class Config:
        populate_by_name = True

    def __init__(self, **data):
        # Handle 'from' field mapping to 'from_email'
        if 'from' in data:
            data['from_email'] = data.pop('from')
        super().__init__(**data)

class EmailIngestRequest(BaseModel):
    source: str
    source_id: str
    title: str
    content: str
    metadata: Optional[EmailMetadata] = None
    tags: Optional[List[EmailTag]] = None
    agent_id: Optional[int] = None  # Optionnel maintenant, on route par tags


def extract_email_tags_from_title(title: str) -> List[str]:
    """Extrait les @tags du titre de l'email (case-insensitive)"""
    import re
    pattern = r'@([a-zA-Z0-9_-]+)'
    matches = re.findall(pattern, title)
    # Normaliser en lowercase avec @
    return [f"@{tag.lower()}" for tag in matches]


def find_agents_by_email_tags(db: Session, tags: List[str]) -> List[Agent]:
    """Trouve tous les agents dont email_tags contient au moins un des tags"""
    if not tags:
        return []

    all_agents = db.query(Agent).filter(Agent.email_tags.isnot(None)).all()
    matching_agents = []

    for agent in all_agents:
        try:
            agent_tags = json.loads(agent.email_tags) if agent.email_tags else []
            # Normaliser les tags de l'agent en lowercase
            agent_tags_lower = [t.lower() for t in agent_tags]
            # Vérifier si au moins un tag match
            if any(tag.lower() in agent_tags_lower for tag in tags):
                matching_agents.append(agent)
        except (json.JSONDecodeError, TypeError):
            continue

    return matching_agents


def verify_email_api_key(request: Request) -> bool:
    """Vérifie l'API Key pour l'ingestion d'emails

    Security: Uses constant-time comparison (hmac.compare_digest) to prevent
    timing attacks that could be used to guess the API key character by character.
    """
    api_key = request.headers.get("X-API-Key", "")
    expected_key = os.getenv("EMAIL_INGEST_API_KEY", "")

    if not expected_key:
        logger.error("EMAIL_INGEST_API_KEY not configured in environment")
        raise HTTPException(status_code=500, detail="API Key not configured on server")

    if not api_key:
        logger.warning(f"Missing API Key for email ingestion")
        raise HTTPException(status_code=401, detail="Invalid API Key")

    # Constant-time comparison to prevent timing attacks
    if not hmac.compare_digest(api_key, expected_key):
        logger.warning(f"Invalid API Key attempt for email ingestion")
        raise HTTPException(status_code=401, detail="Invalid API Key")

    return True


@app.post("/api/emails/ingest")
async def ingest_email(
    payload: EmailIngestRequest,
    request: Request,
    db: Session = Depends(get_db)
):
    """
    Ingère un email depuis la Cloud Function Gmail.

    Authentification via X-API-Key header.
    Route l'email vers les agents basé sur les @tags dans le titre.
    Si aucun tag trouvé ou aucun agent matché, ne rien faire.
    """
    # Vérifier l'API Key
    verify_email_api_key(request)

    try:
        logger.info(f"Ingesting email: {payload.title} (source_id: {payload.source_id})")

        # 1. Extraire les @tags du titre
        extracted_tags = extract_email_tags_from_title(payload.title)
        logger.info(f"Extracted tags from title: {extracted_tags}")

        # 2. Trouver les agents correspondants
        if extracted_tags:
            target_agents = find_agents_by_email_tags(db, extracted_tags)
            logger.info(f"Found {len(target_agents)} agents matching tags: {[a.name for a in target_agents]}")
        else:
            target_agents = []
            logger.info("No tags found in email title")

        # 3. Si aucun tag ou aucun agent matché, ne rien faire
        if not target_agents:
            logger.info(f"No matching agents for email: {payload.title} - skipping ingestion")
            return {
                "success": True,
                "document_ids": [],
                "agents_matched": 0,
                "message": "Aucun tag trouvé ou aucun companion correspondant - email ignoré"
            }

        # 4. Construire le contenu enrichi avec les métadonnées
        enriched_content = f"Sujet: {payload.title}\n"
        if payload.metadata:
            if payload.metadata.from_email:
                enriched_content += f"De: {payload.metadata.from_email}\n"
            if payload.metadata.date:
                enriched_content += f"Date: {payload.metadata.date}\n"
        enriched_content += f"\n{payload.content}"

        # Importer les fonctions nécessaires
        from file_loader import chunk_text
        from openai_client import get_embedding_fast
        from database import DocumentChunk
        import numpy as np

        # Fonction pour découper les gros chunks
        def split_for_embedding(chunk, max_tokens=8192):
            chunk = chunk.replace('\x00', '')
            max_chars = max_tokens * 4
            return [chunk[i:i+max_chars] for i in range(0, len(chunk), max_chars)]

        # Préparer les chunks une seule fois (réutilisés pour chaque agent)
        chunks = chunk_text(enriched_content)
        logger.info(f"Created {len(chunks)} chunks for email")

        # Calculer les embeddings une seule fois
        chunk_embeddings = []
        max_immediate_chunks = 20
        for i, chunk in enumerate(chunks):
            if i < max_immediate_chunks:
                try:
                    sub_chunks = split_for_embedding(chunk, 8192)
                    embeddings = []
                    for sub in sub_chunks:
                        embedding = get_embedding_fast(sub)
                        embeddings.append(embedding)
                    if embeddings:
                        avg_embedding = list(np.mean(np.array(embeddings), axis=0))
                    else:
                        avg_embedding = [0.0] * 1536
                except Exception as e:
                    logger.warning(f"Failed to get embedding for chunk {i}: {e}")
                    avg_embedding = [0.0] * 1536
            else:
                avg_embedding = None
            chunk_embeddings.append(avg_embedding)

        # 5. Créer un document pour CHAQUE agent trouvé
        document_ids = []
        for agent in target_agents:
            # Vérifier les doublons pour cet agent spécifique
            # L'identifiant unique est stocké dans gcs_url pour la déduplication
            unique_id = f"email_{payload.source_id}_agent_{agent.id}"
            existing_doc = db.query(Document).filter(
                Document.gcs_url == unique_id
            ).first()

            if existing_doc:
                logger.info(f"Email already ingested for agent {agent.name}: {payload.source_id}")
                document_ids.append(existing_doc.id)
                continue

            # Créer le document pour cet agent
            # Le filename affiche le titre de l'email (sujet) pour une meilleure lisibilité
            # On garde unique_id dans gcs_url pour la déduplication
            document = Document(
                filename=payload.title,  # Titre de l'email comme nom du document
                content=enriched_content,
                user_id=agent.user_id,
                agent_id=agent.id,
                gcs_url=unique_id  # Identifiant unique pour la déduplication
            )
            db.add(document)
            db.commit()
            db.refresh(document)

            logger.info(f"Document created for agent {agent.name} with ID: {document.id}")

            # Créer les chunks avec les embeddings pré-calculés
            for i, chunk in enumerate(chunks):
                doc_chunk = DocumentChunk(
                    document_id=document.id,
                    chunk_text=chunk,
                    embedding=json.dumps(chunk_embeddings[i]) if chunk_embeddings[i] else None,
                    chunk_index=i
                )
                db.add(doc_chunk)

            db.commit()
            document_ids.append(document.id)

        logger.info(f"Email ingested successfully to {len(document_ids)} agents: {payload.title}")

        return {
            "success": True,
            "document_ids": document_ids,
            "agents_matched": len(target_agents),
            "agents": [{"id": a.id, "name": a.name} for a in target_agents],
            "tags_extracted": extracted_tags,
            "message": f"Email ingéré avec succès vers {len(target_agents)} companion(s)"
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error ingesting email: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        db.rollback()
        raise HTTPException(status_code=500, detail="Erreur lors de l'ingestion de l'email")


@app.post("/api/emails/upload-attachment")
async def upload_email_attachment(
    request: Request,
    file: UploadFile = File(...),
    db: Session = Depends(get_db)
):
    """
    Upload une pièce jointe d'email comme document séparé.
    Le fichier est stocké sur GCS et devient téléchargeable.
    Route vers les agents basé sur les @tags passés dans le form data.

    Authentification via X-API-Key header.
    """
    # Vérifier l'API Key
    verify_email_api_key(request)

    try:
        # Récupérer les données du formulaire
        form = await request.form()
        email_subject = form.get("email_subject", "")
        source_id = form.get("source_id", "")

        logger.info(f"Uploading attachment: {file.filename} for email: {email_subject}")

        # Extraire les @tags du sujet de l'email
        extracted_tags = extract_email_tags_from_title(email_subject) if email_subject else []
        logger.info(f"Extracted tags from email subject: {extracted_tags}")

        # Trouver les agents correspondants
        if extracted_tags:
            target_agents = find_agents_by_email_tags(db, extracted_tags)
            logger.info(f"Found {len(target_agents)} agents matching tags")
        else:
            target_agents = []

        # Si aucun agent matché, ignorer
        if not target_agents:
            logger.info(f"No matching agents for attachment: {file.filename} - skipping")
            return {
                "success": True,
                "document_ids": [],
                "agents_matched": 0,
                "message": "Aucun companion correspondant - pièce jointe ignorée"
            }

        # Vérifier le type de fichier
        allowed_types = ['.pdf', '.txt', '.docx']
        if not any(file.filename.lower().endswith(ext) for ext in allowed_types):
            logger.warning(f"Unsupported file type: {file.filename}")
            return {
                "success": False,
                "message": f"Type de fichier non supporté: {file.filename}"
            }

        # Lire le contenu du fichier
        content = await file.read()

        # Vérifier la taille (10MB max)
        if len(content) > 10 * 1024 * 1024:
            raise HTTPException(status_code=413, detail="Fichier trop volumineux (max 10MB)")

        # Importer les fonctions nécessaires
        from rag_engine import process_document_for_user

        # Créer un document pour CHAQUE agent trouvé
        document_ids = []
        for agent in target_agents:
            # Vérifier les doublons
            unique_id = f"attachment_{source_id}_{file.filename}_agent_{agent.id}"
            existing_doc = db.query(Document).filter(
                Document.gcs_url.contains(file.filename),
                Document.agent_id == agent.id
            ).first()

            if existing_doc and source_id:
                # Vérification plus stricte avec source_id si disponible
                logger.info(f"Attachment may already exist for agent {agent.name}: {file.filename}")

            try:
                # Utiliser process_document_for_user qui gère l'upload GCS
                doc_id = process_document_for_user(
                    filename=file.filename,
                    content=content,
                    user_id=agent.user_id,
                    db=db,
                    agent_id=agent.id
                )
                document_ids.append(doc_id)
                logger.info(f"Attachment uploaded for agent {agent.name}: {file.filename} (doc_id: {doc_id})")

            except Exception as e:
                logger.error(f"Failed to upload attachment for agent {agent.name}: {e}")
                continue

        return {
            "success": True,
            "document_ids": document_ids,
            "agents_matched": len(target_agents),
            "agents": [{"id": a.id, "name": a.name} for a in target_agents],
            "filename": file.filename,
            "message": f"Pièce jointe uploadée vers {len(document_ids)} companion(s)"
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error uploading attachment: {e}")
        logger.error(f"Traceback: {traceback.format_exc()}")


# ============================================================================
# COMPANY & NEO4J ENDPOINTS
# ============================================================================

@app.post("/api/companies")
async def create_company(
    request: Request,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Create a company and affiliate the creator as owner."""
    import secrets as _secrets
    body = await request.json()
    name = body.get("name", "").strip()
    if not name:
        raise HTTPException(status_code=400, detail="Company name is required")

    existing = db.query(Company).filter(Company.name == name).first()
    if existing:
        raise HTTPException(status_code=409, detail="A company with this name already exists")

    # Check user isn't already in an org
    existing_membership = db.query(CompanyMembership).filter(CompanyMembership.user_id == int(user_id)).first()
    if existing_membership:
        raise HTTPException(status_code=409, detail="You are already a member of an organization")

    company = Company(name=name, neo4j_enabled=True, invite_code=_secrets.token_urlsafe(16))
    db.add(company)
    db.commit()
    db.refresh(company)

    # Affiliate the creator
    user = db.query(User).filter(User.id == int(user_id)).first()
    if user:
        user.company_id = company.id
        membership = CompanyMembership(user_id=user.id, company_id=company.id, role="owner")
        db.add(membership)
        db.commit()

    return {"company": {"id": company.id, "name": company.name, "neo4j_enabled": company.neo4j_enabled, "invite_code": company.invite_code}}


@app.get("/api/companies/mine")
async def get_my_company(
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get the current user's company with role and integration status."""
    uid = int(user_id)
    membership = db.query(CompanyMembership).filter(CompanyMembership.user_id == uid).first()
    if not membership:
        return {"company": None}

    company = db.query(Company).filter(Company.id == membership.company_id).first()
    if not company:
        return {"company": None}

    result = {
        "id": company.id,
        "name": company.name,
        "neo4j_enabled": company.neo4j_enabled,
        "role": membership.role,
        "has_neo4j": bool(company._neo4j_uri),
        "has_notion": bool(company._notion_api_key),
    }
    # Include invite_code for admin/owner
    if membership.role in ("admin", "owner"):
        result["invite_code"] = company.invite_code
        result["invite_code_enabled"] = company.invite_code_enabled

    return {"company": result}


@app.put("/api/user/company")
async def affiliate_user_to_company(
    request: Request,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Affiliate user to an existing company by name."""
    body = await request.json()
    company_name = body.get("company_name", "").strip()
    if not company_name:
        raise HTTPException(status_code=400, detail="company_name is required")

    company = db.query(Company).filter(Company.name == company_name).first()
    if not company:
        raise HTTPException(status_code=404, detail="Company not found")

    user = db.query(User).filter(User.id == int(user_id)).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")

    user.company_id = company.id
    db.commit()

    return {"company": {"id": company.id, "name": company.name, "neo4j_enabled": company.neo4j_enabled}}


# ============================================================================
# ORGANIZATION MANAGEMENT ENDPOINTS
# ============================================================================

@app.post("/api/companies/invite")
async def invite_to_company(
    request: Request,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Send an email invitation to join the organization. Requires admin+."""
    import secrets as _secrets
    from permissions import require_role

    membership = require_role(int(user_id), db, "admin")
    body = await request.json()
    email = body.get("email", "").strip().lower()
    role = body.get("role", "member")

    if not email:
        raise HTTPException(status_code=400, detail="Email is required")
    if role not in ("member", "admin"):
        raise HTTPException(status_code=400, detail="Role must be 'member' or 'admin'")

    # Check if user with this email is already a member
    existing_user = db.query(User).filter(User.email == email).first()
    if existing_user:
        existing_m = db.query(CompanyMembership).filter(
            CompanyMembership.user_id == existing_user.id,
            CompanyMembership.company_id == membership.company_id
        ).first()
        if existing_m:
            raise HTTPException(status_code=409, detail="This user is already a member of the organization")

    # Check for pending invitation
    pending = db.query(CompanyInvitation).filter(
        CompanyInvitation.company_id == membership.company_id,
        CompanyInvitation.email == email,
        CompanyInvitation.status == "pending"
    ).first()
    if pending:
        raise HTTPException(status_code=409, detail="An invitation is already pending for this email")

    token = _secrets.token_urlsafe(48)
    invitation = CompanyInvitation(
        company_id=membership.company_id,
        email=email,
        role=role,
        token=token,
        invited_by_user_id=int(user_id),
        expires_at=datetime.utcnow() + timedelta(days=7)
    )
    db.add(invitation)
    db.commit()

    # Send invitation email
    company = db.query(Company).filter(Company.id == membership.company_id).first()
    try:
        frontend_url = os.getenv("FRONTEND_URL", "http://localhost:3000")
        join_link = f"{frontend_url}/join?token={token}"
        send_invitation_email(email, company.name, join_link)
    except Exception as e:
        logger.warning(f"Failed to send invitation email to {email}: {e}")

    return {"message": "Invitation sent", "token": token}


@app.post("/api/companies/join")
async def join_company(
    request: Request,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Join an organization via invite token (email) or invite_code (shareable link)."""
    import secrets as _secrets
    uid = int(user_id)
    body = await request.json()
    token = body.get("token", "").strip()
    invite_code = body.get("invite_code", "").strip()

    # Check user isn't already in an org
    existing = db.query(CompanyMembership).filter(CompanyMembership.user_id == uid).first()
    if existing:
        raise HTTPException(status_code=409, detail="You are already a member of an organization")

    user = db.query(User).filter(User.id == uid).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")

    if token:
        # Join via email invitation token
        invitation = db.query(CompanyInvitation).filter(
            CompanyInvitation.token == token,
            CompanyInvitation.status == "pending"
        ).first()
        if not invitation:
            raise HTTPException(status_code=404, detail="Invalid or expired invitation")
        if invitation.expires_at < datetime.utcnow():
            invitation.status = "expired"
            db.commit()
            raise HTTPException(status_code=410, detail="This invitation has expired")

        company_id = invitation.company_id
        role = invitation.role
        invitation.status = "accepted"

    elif invite_code:
        # Join via shareable invite code
        company = db.query(Company).filter(
            Company.invite_code == invite_code,
            Company.invite_code_enabled == True
        ).first()
        if not company:
            raise HTTPException(status_code=404, detail="Invalid invite code")
        company_id = company.id
        role = "member"

    else:
        raise HTTPException(status_code=400, detail="Either 'token' or 'invite_code' is required")

    user.company_id = company_id
    membership = CompanyMembership(user_id=uid, company_id=company_id, role=role)
    db.add(membership)
    db.commit()

    company = db.query(Company).filter(Company.id == company_id).first()
    return {"message": f"You have joined {company.name}", "company": {"id": company.id, "name": company.name, "role": role}}


@app.post("/api/companies/invite-code/regenerate")
async def regenerate_invite_code(
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Regenerate the shareable invite code. Owner only."""
    import secrets as _secrets
    from permissions import require_role

    membership = require_role(int(user_id), db, "owner")
    company = db.query(Company).filter(Company.id == membership.company_id).first()
    company.invite_code = _secrets.token_urlsafe(16)
    db.commit()

    return {"invite_code": company.invite_code}


@app.put("/api/companies/invite-code/toggle")
async def toggle_invite_code(
    request: Request,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Enable or disable the shareable invite code. Owner only."""
    from permissions import require_role

    membership = require_role(int(user_id), db, "owner")
    body = await request.json()
    enabled = body.get("enabled", True)

    company = db.query(Company).filter(Company.id == membership.company_id).first()
    company.invite_code_enabled = bool(enabled)
    db.commit()

    return {"invite_code_enabled": company.invite_code_enabled}


@app.get("/api/companies/members")
async def list_company_members(
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """List all members of the organization. Requires admin+."""
    from permissions import require_role
    from sqlalchemy import func

    membership = require_role(int(user_id), db, "admin")

    members = (
        db.query(CompanyMembership, User)
        .join(User, CompanyMembership.user_id == User.id)
        .filter(CompanyMembership.company_id == membership.company_id)
        .order_by(CompanyMembership.joined_at.asc())
        .all()
    )

    result = []
    for m, u in members:
        agent_count = db.query(Agent).filter(Agent.user_id == u.id).count()
        result.append({
            "id": m.id,
            "user_id": u.id,
            "username": u.username,
            "email": u.email,
            "role": m.role,
            "joined_at": m.joined_at.isoformat() if m.joined_at else None,
            "agent_count": agent_count,
        })

    return {"members": result}


@app.put("/api/companies/members/{member_id}/role")
async def update_member_role(
    member_id: int,
    request: Request,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Change a member's role. Owner only."""
    from permissions import require_role

    membership = require_role(int(user_id), db, "owner")
    body = await request.json()
    new_role = body.get("role", "")

    if new_role not in ("member", "admin"):
        raise HTTPException(status_code=400, detail="Role must be 'member' or 'admin'")

    target = db.query(CompanyMembership).filter(
        CompanyMembership.id == member_id,
        CompanyMembership.company_id == membership.company_id
    ).first()
    if not target:
        raise HTTPException(status_code=404, detail="Member not found")

    if target.role == "owner":
        raise HTTPException(status_code=403, detail="Cannot change the owner's role")

    target.role = new_role
    db.commit()

    return {"message": f"Role updated to {new_role}"}


@app.delete("/api/companies/members/{member_id}")
async def remove_member(
    member_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Remove a member from the organization. Requires admin+."""
    from permissions import require_role

    membership = require_role(int(user_id), db, "admin")
    target = db.query(CompanyMembership).filter(
        CompanyMembership.id == member_id,
        CompanyMembership.company_id == membership.company_id
    ).first()
    if not target:
        raise HTTPException(status_code=404, detail="Member not found")

    if target.role == "owner":
        raise HTTPException(status_code=403, detail="Cannot remove the owner")

    # Clean up agent shares for the departing user
    db.query(AgentShare).filter(AgentShare.user_id == target.user_id).delete()

    # Remove the company_id from the user record too
    user = db.query(User).filter(User.id == target.user_id).first()
    if user:
        user.company_id = None

    db.delete(target)
    db.commit()

    return {"message": "Member removed"}


@app.post("/api/companies/leave")
async def leave_company(
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Leave the current organization. Owner cannot leave."""
    from permissions import get_user_membership

    uid = int(user_id)
    membership = get_user_membership(uid, db)
    if not membership:
        raise HTTPException(status_code=404, detail="You are not a member of any organization")

    if membership.role == "owner":
        raise HTTPException(status_code=403, detail="Owner cannot leave the organization. Transfer ownership first or delete the organization.")

    # Clean up agent shares for the departing user
    db.query(AgentShare).filter(AgentShare.user_id == uid).delete()

    user = db.query(User).filter(User.id == uid).first()
    if user:
        user.company_id = None

    db.delete(membership)
    db.commit()

    return {"message": "You have left the organization"}


@app.put("/api/companies/integrations")
async def update_company_integrations(
    request: Request,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Configure org-level integration credentials. Owner only."""
    from permissions import require_role

    membership = require_role(int(user_id), db, "owner")
    company = db.query(Company).filter(Company.id == membership.company_id).first()
    body = await request.json()

    # Neo4j
    if "neo4j_uri" in body:
        company.org_neo4j_uri = body["neo4j_uri"] or None
    if "neo4j_user" in body:
        company.org_neo4j_user = body["neo4j_user"] or None
    if "neo4j_password" in body:
        company.org_neo4j_password = body["neo4j_password"] or None

    # Notion
    if "notion_api_key" in body:
        company.org_notion_api_key = body["notion_api_key"] or None

    db.commit()
    return {"message": "Integrations updated"}


@app.get("/api/companies/integrations")
async def get_company_integrations(
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get the integration status (without secrets). Owner only."""
    from permissions import require_role

    membership = require_role(int(user_id), db, "owner")
    company = db.query(Company).filter(Company.id == membership.company_id).first()

    return {
        "neo4j": {
            "configured": bool(company._neo4j_uri),
            "uri": company.org_neo4j_uri[:30] + "..." if company._neo4j_uri else None,
            "user": company.org_neo4j_user if company._neo4j_user else None,
        },
        "notion": {
            "configured": bool(company._notion_api_key),
            "key_preview": company.org_notion_api_key[:8] + "..." if company._notion_api_key else None,
        },
    }


@app.get("/api/companies/agents")
async def list_company_agents(
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """List all agents in the organization. Requires admin+."""
    from permissions import require_role

    membership = require_role(int(user_id), db, "admin")

    # Get all user IDs in this company
    member_ids = [
        m.user_id for m in
        db.query(CompanyMembership).filter(CompanyMembership.company_id == membership.company_id).all()
    ]

    agents = db.query(Agent, User).join(User, Agent.user_id == User.id).filter(
        Agent.user_id.in_(member_ids)
    ).order_by(Agent.created_at.desc()).all()

    result = []
    for agent, owner in agents:
        doc_count = db.query(Document).filter(Document.agent_id == agent.id).count()
        share_count = db.query(AgentShare).filter(AgentShare.agent_id == agent.id).count()
        result.append({
            "id": agent.id,
            "name": agent.name,
            "type": agent.type,
            "statut": agent.statut,
            "llm_provider": agent.llm_provider,
            "owner_username": owner.username,
            "owner_id": owner.id,
            "created_at": agent.created_at.isoformat() if agent.created_at else None,
            "document_count": doc_count,
            "shared_with_count": share_count,
        })

    return {"agents": result}


@app.delete("/api/companies/agents/{agent_id}")
async def delete_company_agent(
    agent_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Delete an agent from the organization. Requires admin+."""
    from permissions import require_role

    membership = require_role(int(user_id), db, "admin")

    agent = db.query(Agent).filter(Agent.id == agent_id).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")

    # Verify the agent owner is a member of the same company
    owner_membership = db.query(CompanyMembership).filter(
        CompanyMembership.user_id == agent.user_id,
        CompanyMembership.company_id == membership.company_id
    ).first()
    if not owner_membership:
        raise HTTPException(status_code=404, detail="Agent not found in this organization")

    try:
        _delete_agent_and_related_data(agent, agent.user_id, db)
        db.commit()
        return {"message": "Agent deleted successfully"}
    except Exception as e:
        logger.error(f"Error deleting org agent: {e}")
        logger.error(f"Full traceback: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Internal server error")


@app.post("/api/companies/agents/{agent_id}/share")
async def share_agent(
    agent_id: int,
    request: Request,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Share an agent with another user in the same org. Owner of agent or admin+ required."""
    from permissions import require_role, get_user_membership

    uid = int(user_id)
    body = await request.json()
    target_user_id = body.get("user_id")
    can_edit = bool(body.get("can_edit", False))
    if not target_user_id:
        raise HTTPException(status_code=422, detail="user_id is required")

    agent = db.query(Agent).filter(Agent.id == agent_id).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")

    # Check permission: owner of the agent OR admin+
    is_owner = agent.user_id == uid
    if not is_owner:
        require_role(uid, db, "admin")

    # Both users must be in the same org
    caller_membership = get_user_membership(uid, db)
    if not caller_membership:
        raise HTTPException(status_code=403, detail="You are not in an organization")

    target_membership = db.query(CompanyMembership).filter(
        CompanyMembership.user_id == target_user_id,
        CompanyMembership.company_id == caller_membership.company_id
    ).first()
    if not target_membership:
        raise HTTPException(status_code=404, detail="Target user not found in this organization")

    # Cannot share with the owner of the agent
    if target_user_id == agent.user_id:
        raise HTTPException(status_code=400, detail="Cannot share an agent with its owner")

    # Check if already shared
    existing = db.query(AgentShare).filter(
        AgentShare.agent_id == agent_id,
        AgentShare.user_id == target_user_id
    ).first()
    if existing:
        raise HTTPException(status_code=409, detail="Agent already shared with this user")

    share = AgentShare(
        agent_id=agent_id,
        user_id=target_user_id,
        shared_by_user_id=uid,
        can_edit=can_edit
    )
    db.add(share)
    db.commit()

    return {"message": "Agent shared successfully"}


@app.delete("/api/companies/agents/{agent_id}/share/{target_user_id}")
async def unshare_agent(
    agent_id: int,
    target_user_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Remove a share. Owner of agent, admin+, or the target user themselves."""
    from permissions import require_role, get_user_membership

    uid = int(user_id)

    agent = db.query(Agent).filter(Agent.id == agent_id).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")

    is_owner = agent.user_id == uid
    is_self = target_user_id == uid

    if not is_owner and not is_self:
        require_role(uid, db, "admin")

    share = db.query(AgentShare).filter(
        AgentShare.agent_id == agent_id,
        AgentShare.user_id == target_user_id
    ).first()
    if not share:
        raise HTTPException(status_code=404, detail="Share not found")

    db.delete(share)
    db.commit()

    return {"message": "Share removed"}


@app.put("/api/companies/agents/{agent_id}/share/{target_user_id}")
async def update_agent_share(
    agent_id: int,
    target_user_id: int,
    request: Request,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Update can_edit permission on a share. Owner of agent or admin+ required."""
    from permissions import require_role

    uid = int(user_id)

    agent = db.query(Agent).filter(Agent.id == agent_id).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")

    is_owner = agent.user_id == uid
    if not is_owner:
        require_role(uid, db, "admin")

    share = db.query(AgentShare).filter(
        AgentShare.agent_id == agent_id,
        AgentShare.user_id == target_user_id
    ).first()
    if not share:
        raise HTTPException(status_code=404, detail="Share not found")

    body = await request.json()
    if "can_edit" in body:
        share.can_edit = bool(body["can_edit"])
    db.commit()

    return {"message": "Share updated", "can_edit": share.can_edit}


@app.get("/api/companies/agents/{agent_id}/shares")
async def list_agent_shares(
    agent_id: int,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """List users an agent is shared with. Owner of agent or admin+ required."""
    from permissions import require_role, get_user_membership

    uid = int(user_id)

    agent = db.query(Agent).filter(Agent.id == agent_id).first()
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")

    is_owner = agent.user_id == uid
    if not is_owner:
        require_role(uid, db, "admin")

    shares = db.query(AgentShare, User).join(User, AgentShare.user_id == User.id).filter(
        AgentShare.agent_id == agent_id
    ).all()

    return {"shares": [
        {
            "user_id": u.id,
            "username": u.username,
            "email": u.email,
            "can_edit": s.can_edit,
            "shared_at": s.created_at.isoformat() if s.created_at else None
        }
        for s, u in shares
    ]}


@app.get("/api/neo4j/persons")
async def list_neo4j_persons(
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """List Person nodes from Neo4j for the user's company."""
    user = db.query(User).filter(User.id == int(user_id)).first()
    if not user or not user.company_id:
        return {"persons": []}

    try:
        from neo4j_client import get_persons_for_company
        persons = get_persons_for_company(user.company_id)
        return {"persons": persons}
    except Exception as e:
        logger.warning(f"Failed to list Neo4j persons: {e}")
        return {"persons": []}


# ============================================================================
# PASSWORD CHANGE ENDPOINT
# ============================================================================

# Rate limiting for password change (separate from auth rate limit)
_password_change_rate_limit_fallback = {}
_PASSWORD_CHANGE_LIMIT = 5
_PASSWORD_CHANGE_WINDOW = 300  # 5 minutes

def _check_password_change_rate_limit(user_id: str) -> bool:
    """Check if user has exceeded rate limit for password change. 5 attempts per 5 minutes."""
    key = f"rate_limit:password_change:{user_id}"

    if redis_client:
        try:
            current = redis_client.get(key)
            current = int(current) if current else 0
            if current >= _PASSWORD_CHANGE_LIMIT:
                return False
            pipe = redis_client.pipeline()
            pipe.incr(key)
            pipe.expire(key, _PASSWORD_CHANGE_WINDOW)
            pipe.execute()
            return True
        except Exception as e:
            logger.error(f"Redis rate limit check failed for password change: {e}. Using fallback.")

    now = time.time()
    attempts = _password_change_rate_limit_fallback.get(user_id, [])
    attempts = [t for t in attempts if now - t < _PASSWORD_CHANGE_WINDOW]
    if len(attempts) >= _PASSWORD_CHANGE_LIMIT:
        return False
    attempts.append(now)
    _password_change_rate_limit_fallback[user_id] = attempts
    return True


@app.post("/api/user/change-password")
async def change_password(
    request: Request,
    body: ChangePasswordRequest,
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Change the authenticated user's password."""
    if not _check_password_change_rate_limit(user_id):
        raise HTTPException(status_code=429, detail="Too many attempts. Please try again later.")

    db_user = db.query(User).filter(User.id == int(user_id)).first()
    if not db_user:
        raise HTTPException(status_code=404, detail="User not found")

    if not verify_password(body.current_password, db_user.hashed_password):
        raise HTTPException(status_code=401, detail="Current password is incorrect")

    if body.current_password == body.new_password:
        raise HTTPException(status_code=400, detail="New password must be different from current password")

    db_user.hashed_password = hash_password(body.new_password)
    db.commit()

    return {"message": "Password changed successfully"}


# ============================================================================
# GDPR COMPLIANCE ENDPOINTS (Phase 5)
# ============================================================================

@app.get("/api/user/export-data")
async def export_user_data(
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """
    Export all user data in JSON format (GDPR Article 15 - Right of Access).

    Returns complete user data package including:
    - Profile information
    - Agents created
    - Documents uploaded
    - Conversations and messages
    - Team memberships
    """
    try:
        user = db.query(User).filter(User.id == int(user_id)).first()
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        # Get all user agents
        agents = db.query(Agent).filter(Agent.user_id == int(user_id)).all()
        agents_data = []
        for agent in agents:
            # Get documents for this agent
            documents = db.query(Document).filter(Document.agent_id == agent.id).all()
            agents_data.append({
                "id": agent.id,
                "name": agent.name,
                "type": agent.type,
                "statut": agent.statut,
                "contexte": agent.contexte,
                "biographie": agent.biographie,
                "created_at": agent.created_at.isoformat() if agent.created_at else None,
                "documents_count": len(documents),
                "documents": [{
                    "id": doc.id,
                    "filename": doc.filename,
                    "created_at": doc.created_at.isoformat() if doc.created_at else None
                } for doc in documents]
            })

        # Get all conversations
        conversations = db.query(Conversation).filter(
            Conversation.agent_id.in_([a.id for a in agents])
        ).all()
        conversations_data = []
        for conv in conversations:
            messages = db.query(Message).filter(Message.conversation_id == conv.id).all()
            conversations_data.append({
                "id": conv.id,
                "title": conv.title,
                "agent_id": conv.agent_id,
                "created_at": conv.created_at.isoformat() if conv.created_at else None,
                "messages": [{
                    "role": msg.role,
                    "content": msg.content,
                    "timestamp": msg.timestamp.isoformat() if msg.timestamp else None,
                    "feedback": msg.feedback
                } for msg in messages]
            })

        # Get team memberships
        teams = db.query(Team).filter(Team.user_id == int(user_id)).all()
        teams_data = [{
            "id": team.id,
            "name": team.name,
            "contexte": team.contexte,
            "created_at": team.created_at.isoformat() if team.created_at else None
        } for team in teams]

        export_data = {
            "export_date": datetime.utcnow().isoformat(),
            "user": {
                "id": user.id,
                "username": user.username,
                "email": user.email,
                "created_at": user.created_at.isoformat() if user.created_at else None
            },
            "agents": agents_data,
            "conversations": conversations_data,
            "teams": teams_data,
            "statistics": {
                "total_agents": len(agents),
                "total_documents": sum(len(a["documents"]) for a in agents_data),
                "total_conversations": len(conversations_data),
                "total_messages": sum(len(c["messages"]) for c in conversations_data),
                "total_teams": len(teams_data)
            }
        }

        logger.info(f"User {user_id} exported their data (GDPR Art. 15)")
        return export_data

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error exporting user data: {e}")
        raise HTTPException(status_code=500, detail="Failed to export data")


@app.get("/api/user/stats")
async def get_user_stats(
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get analytics data for the current user's activity."""
    try:
        from sqlalchemy import func, cast, Date

        uid = int(user_id)
        thirty_days_ago = datetime.utcnow() - timedelta(days=30)

        # 1. Messages per day (last 30 days)
        messages_per_day_q = (
            db.query(
                cast(Message.timestamp, Date).label("date"),
                func.count(Message.id).label("count")
            )
            .join(Conversation, Message.conversation_id == Conversation.id)
            .filter(Conversation.user_id == uid)
            .filter(Message.timestamp >= thirty_days_ago)
            .group_by(cast(Message.timestamp, Date))
            .order_by(cast(Message.timestamp, Date))
            .all()
        )
        messages_per_day = [{"date": str(r.date), "count": r.count} for r in messages_per_day_q]

        # 2. Messages per agent
        messages_per_agent_q = (
            db.query(
                Agent.name.label("name"),
                func.count(Message.id).label("messages")
            )
            .join(Conversation, Message.conversation_id == Conversation.id)
            .join(Agent, Conversation.agent_id == Agent.id)
            .filter(Conversation.user_id == uid)
            .group_by(Agent.name)
            .order_by(func.count(Message.id).desc())
            .all()
        )
        messages_per_agent = [{"name": r.name, "messages": r.messages} for r in messages_per_agent_q]

        # 3. Conversations per agent
        conversations_per_agent_q = (
            db.query(
                Agent.name.label("name"),
                func.count(Conversation.id).label("conversations")
            )
            .join(Agent, Conversation.agent_id == Agent.id)
            .filter(Conversation.user_id == uid)
            .group_by(Agent.name)
            .order_by(func.count(Conversation.id).desc())
            .all()
        )
        conversations_per_agent = [{"name": r.name, "conversations": r.conversations} for r in conversations_per_agent_q]

        # 4. Feedback distribution (agent messages only)
        feedback_q = (
            db.query(
                Message.feedback,
                func.count(Message.id).label("count")
            )
            .join(Conversation, Message.conversation_id == Conversation.id)
            .filter(Conversation.user_id == uid)
            .filter(Message.role == "agent")
            .group_by(Message.feedback)
            .all()
        )
        feedback = {"like": 0, "dislike": 0, "none": 0}
        for r in feedback_q:
            key = r.feedback if r.feedback in ("like", "dislike") else "none"
            feedback[key] += r.count

        # 5. Average messages per conversation
        subq = (
            db.query(
                Conversation.id,
                func.count(Message.id).label("msg_count")
            )
            .join(Message, Message.conversation_id == Conversation.id)
            .filter(Conversation.user_id == uid)
            .group_by(Conversation.id)
            .subquery()
        )
        avg_result = db.query(func.avg(subq.c.msg_count)).scalar()
        avg_messages = round(float(avg_result), 1) if avg_result else 0

        # 6. Role distribution (user vs agent messages)
        role_q = (
            db.query(
                Message.role,
                func.count(Message.id).label("count")
            )
            .join(Conversation, Message.conversation_id == Conversation.id)
            .filter(Conversation.user_id == uid)
            .group_by(Message.role)
            .all()
        )
        role_distribution = {"user": 0, "agent": 0}
        for r in role_q:
            if r.role in role_distribution:
                role_distribution[r.role] = r.count

        # Most active agent
        most_active_agent = messages_per_agent[0]["name"] if messages_per_agent else None

        return {
            "messages_per_day": messages_per_day,
            "messages_per_agent": messages_per_agent,
            "conversations_per_agent": conversations_per_agent,
            "feedback": feedback,
            "avg_messages_per_conversation": avg_messages,
            "role_distribution": role_distribution,
            "most_active_agent": most_active_agent
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error fetching user stats: {e}")
        raise HTTPException(status_code=500, detail="Failed to fetch stats")


@app.delete("/api/user/delete-account")
async def delete_user_account(
    user_id: str = Depends(verify_token),
    db: Session = Depends(get_db),
    anonymize: bool = False
):
    """
    Delete user account and all associated data (GDPR Article 17 - Right to Erasure).

    Query params:
    - anonymize: If True, anonymize data instead of deleting (keeps analytics)

    Deletes/Anonymizes:
    - User profile
    - All agents created by user
    - All documents uploaded
    - All conversations and messages
    - All teams created by user
    """
    try:
        user = db.query(User).filter(User.id == int(user_id)).first()
        if not user:
            raise HTTPException(status_code=404, detail="User not found")

        if anonymize:
            # Anonymize instead of delete (GDPR-compliant)
            user.username = f"deleted_user_{user.id}"
            user.email = f"deleted_{user.id}@anonymized.local"
            user.hashed_password = "ANONYMIZED"

            # Anonymize agents
            agents = db.query(Agent).filter(Agent.user_id == int(user_id)).all()
            for agent in agents:
                agent.name = f"Anonymized Agent {agent.id}"
                agent.contexte = "ANONYMIZED"
                agent.biographie = "ANONYMIZED"

            db.commit()
            logger.info(f"User {user_id} account anonymized (GDPR Art. 17)")
            return {
                "message": "Account anonymized successfully",
                "anonymized": True,
                "user_id": user.id
            }
        else:
            # Complete deletion
            # Get all agents
            agents = db.query(Agent).filter(Agent.user_id == int(user_id)).all()
            agent_ids = [agent.id for agent in agents]

            # Delete conversations and messages for user's agents
            if agent_ids:
                conversations = db.query(Conversation).filter(
                    Conversation.agent_id.in_(agent_ids)
                ).all()
                conv_ids = [conv.id for conv in conversations]

                if conv_ids:
                    # Delete messages
                    db.query(Message).filter(Message.conversation_id.in_(conv_ids)).delete(synchronize_session=False)
                    # Delete conversations
                    db.query(Conversation).filter(Conversation.id.in_(conv_ids)).delete(synchronize_session=False)

            # Delete documents (will cascade to chunks)
            db.query(Document).filter(Document.user_id == int(user_id)).delete(synchronize_session=False)

            # Delete agents
            db.query(Agent).filter(Agent.user_id == int(user_id)).delete(synchronize_session=False)

            # Delete teams
            db.query(Team).filter(Team.user_id == int(user_id)).delete(synchronize_session=False)

            # Delete password reset tokens
            from database import PasswordResetToken
            db.query(PasswordResetToken).filter(PasswordResetToken.user_id == int(user_id)).delete(synchronize_session=False)

            # Finally delete user
            db.delete(user)
            db.commit()

            logger.info(f"User {user_id} account completely deleted (GDPR Art. 17)")
            return {
                "message": "Account deleted successfully",
                "anonymized": False,
                "deleted": True
            }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error deleting user account: {e}")
        db.rollback()
        raise HTTPException(status_code=500, detail="Failed to delete account")